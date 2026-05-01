#!/usr/bin/env python3
"""
Build comprehensive Plan de Dirección del Proyecto presentation.
PMBOK 8th Edition - LoRA-HDL-QA Project
All required documents, one per slide (max).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

# ── Constants ──────────────────────────────────────────────
NAVY = RGBColor(0x1A, 0x27, 0x44)
BLUE = RGBColor(0x42, 0x85, 0xF4)
BLUE_LIGHT = RGBColor(0xC8, 0xDC, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
DARK_TEXT = RGBColor(0x1A, 0x27, 0x44)
MED_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
ACCENT_BLUE = RGBColor(0x42, 0x85, 0xC8)
GREEN = RGBColor(0x34, 0xA8, 0x53)
RED = RGBColor(0xEA, 0x43, 0x35)
YELLOW = RGBColor(0xFB, 0xBC, 0x04)

SLIDE_W = Emu(12192000)  # 16:9 widescreen
SLIDE_H = Emu(6858000)

FONT_TITLE = 'Calibri'
FONT_BODY = 'Calibri'

# ── Helper functions ───────────────────────────────────────

def add_header_bar(slide, slide_num, title, subtitle=None):
    """Add standard header bar with navy background"""
    # Top navy bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Emu(822960)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    
    # Slide number (big)
    txBox = slide.shapes.add_textbox(Emu(228600), Emu(182880), Emu(548640), Emu(502920))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{slide_num:02d}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE
    
    # Title
    txBox2 = slide.shapes.add_textbox(Emu(868680), Emu(228600), Emu(10972800), Emu(457200))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = FONT_TITLE
    
    # Subtitle
    if subtitle:
        txBox3 = slide.shapes.add_textbox(Emu(274320), Emu(868680), Emu(11612880), Emu(365760))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = subtitle
        p3.font.size = Pt(14)
        p3.font.color.rgb = GRAY_TEXT
        p3.font.name = FONT_BODY

def add_domain_slide(slide, domain_name, description, icon_text):
    """Add a domain separator slide"""
    # Full navy background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    
    # Blue accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(2743200), SLIDE_W, Emu(91440)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    
    # Domain name
    txBox = slide.shapes.add_textbox(Emu(914400), Emu(1371600), Emu(10363200), Emu(914400))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = domain_name
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    txBox2 = slide.shapes.add_textbox(Emu(914400), Emu(3200400), Emu(10363200), Emu(731520))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = description
    p2.font.size = Pt(18)
    p2.font.color.rgb = BLUE_LIGHT
    p2.font.name = FONT_BODY
    p2.alignment = PP_ALIGN.CENTER
    
    # Bottom bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(6766560), SLIDE_W, Emu(91440)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

def add_card(slide, left, top, width, height, title, body_lines, title_size=Pt(14), body_size=Pt(11)):
    """Add a rounded rectangle card with title and body text"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
    card.line.width = Pt(0.5)
    
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(68580)
    tf.margin_bottom = Emu(68580)
    
    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = title_size
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_TITLE
    p.space_after = Pt(4)
    
    # Body
    for line in body_lines:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = body_size
        p2.font.color.rgb = MED_TEXT
        p2.font.name = FONT_BODY
        p2.space_after = Pt(1)
    
    return card

def add_simple_table(slide, left, top, width, height, headers, rows, col_widths=None, font_size=Pt(9)):
    """Add a simple table"""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    # Headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = font_size
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.font.name = FONT_BODY
        # Blue header bg
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = cell._tc.makeelement(qn('a:solidFill'), {})
        srgbClr = cell._tc.makeelement(qn('a:srgbClr'), {'val': '4285C8'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)
    
    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = font_size
                para.font.color.rgb = MED_TEXT
                para.font.name = FONT_BODY
            # Alternating row colors
            if r % 2 == 0:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = cell._tc.makeelement(qn('a:solidFill'), {})
                srgbClr = cell._tc.makeelement(qn('a:srgbClr'), {'val': 'F8F9FA'})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)
    
    return table_shape

# ── Build presentation ─────────────────────────────────────

def build_presentation():
    # Start from existing presentation to preserve theme
    src_path = '/Users/pabloflores/Documents/estructura/revision/LoRA_Presentacion_Completa_Entregas_1_2_3.pptx'
    prs = Presentation(src_path)
    
    # We'll build from scratch to have full control, but copy the slide master
    # Actually, let's clone the existing prs and clear all slides, then rebuild
    
    # Get the blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    # Remove all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    
    slide_num = [0]  # Mutable counter
    
    def next_slide():
        slide_num[0] += 1
        return prs.slides.add_slide(blank_layout)
    
    def s():
        return slide_num[0]
    
    # ═══════════════════════════════════════════════════════════
    # PORTADA
    # ═══════════════════════════════════════════════════════════
    slide = next_slide()
    # Navy bg
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    # Blue accent
    for y_pos in [Emu(1828800), Emu(6766560)]:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), y_pos, SLIDE_W, Emu(91440))
        accent.fill.solid(); accent.fill.fore_color.rgb = BLUE; accent.line.fill.background()
    # Texts
    def add_cover_text(left, top, width, height, text, size, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text; p.font.size = size
        p.font.bold = bold; p.font.color.rgb = color; p.font.name = FONT_TITLE
        p.alignment = align
    
    add_cover_text(Emu(914400), Emu(457200), Emu(10363200), Emu(548640), 
                   "UNIVERSIDAD DEL ISTMO DE GUATEMALA", Pt(20), True, RGBColor(0x42, 0x85, 0xF4))
    add_cover_text(Emu(914400), Emu(914400), Emu(10363200), Emu(365760),
                   "Facultad de Ingeniería (FING)  |  Ingeniería en Sistemas y CC", Pt(12), False, BLUE_LIGHT)
    add_cover_text(Emu(914400), Emu(2103120), Emu(10363200), Emu(914400),
                   "PLAN DE DIRECCIÓN DEL PROYECTO", Pt(36), True, WHITE)
    add_cover_text(Emu(914400), Emu(2926080), Emu(10363200), Emu(457200),
                   "PMBOK 8va Edición — todos los documentos de iniciación y planificación", Pt(14), False, BLUE_LIGHT)
    add_cover_text(Emu(914400), Emu(3657600), Emu(10363200), Emu(457200),
                   "Proyecto: LoRA-HDL-QA", Pt(18), True, WHITE)
    add_cover_text(Emu(914400), Emu(4114800), Emu(10363200), Emu(365760),
                   "PM: Pablo Rodolfo Alexander Flores Mollinedo", Pt(14), False, BLUE_LIGHT)
    add_cover_text(Emu(914400), Emu(4754880), Emu(10363200), Emu(365760),
                   "Febrero — Septiembre 2026", Pt(12), False, BLUE_LIGHT)
    
    # ═══════════════════════════════════════════════════════════
    # DOMAIN 1: GOBERNANZA
    # ═══════════════════════════════════════════════════════════
    slide = next_slide()
    add_domain_slide(slide, 
        "DOMINIO DE DESEMPEÑO:\nGOBERNANZA",
        "PMBOK 8va Edición — Un modelo de gestión alineado con los objetivos institucionales. "
        "Define estructuras de autoridad, toma de decisiones, y mecanismos de control del proyecto.",
        "🏛️")
    
    # ── Slide: Acta de Constitución del Proyecto ──
    slide = next_slide()
    add_header_bar(slide, s(), "ACTA DE CONSTITUCIÓN DEL PROYECTO",
                   "Documento formal que autoriza el proyecto | PMBOK 8va, Gobernanza")
    
    cards_data = [
        ("Información General", [
            "Proyecto: LoRA-HDL-QA — Degradación Catastrófica en Fine-Tuning LoRA para HDL/QA",
            "PM: Pablo R. A. Flores Mollinedo | Director: Ing. Estuardo Rodríguez",
            "Asesor: Por asignar | Período: Feb — Oct 2026 (1,287 h planificadas)",
            "Tipo: Proyecto de graduación FING/UNIS — trabajo individual"
        ]),
        ("Caso de Negocio y Objetivos", [
            "Trabajo de graduación para cierre de pensum en Ingeniería en Sistemas",
            "Investigar impacto de estrategias de preprocesamiento en fine-tuning LoRA para generación HDL",
            "Contribuir al campo de ML aplicado a diseño de hardware con evidencia experimental",
            "Objetivo: paper de investigación + repositorio reproducible + defensa oral"
        ]),
        ("Entregables de Alto Nivel", [
            "Anteproyecto aprobado por comité evaluador FING (Mar 2026)",
            "Paper de investigación comparativo (2 configuraciones LoRA evaluadas en VerilogEval v2)",
            "Repositorio GitHub público con scripts, pesos .safetensors, y resultados",
            "Informe final formato FING/IMRaD + defensa oral ante jurado"
        ]),
        ("Hitos Clave", [
            "Inicio: Feb 02 | Inscripción + asesor: Feb 02-16 | Plan proyecto: Feb 17-26",
            "Anteproyecto: Feb 26 — Mar 25 | Diseño metodológico: Mar 25 — Abr 08",
            "Ejecución (entrenamiento + evaluación): Abr 08 — Ago 17",
            "Cierre (correcciones + defensa): Ago 17 — Sep 14"
        ]),
        ("Presupuesto Estimado", [
            "Mano de obra: 1,287h × Q 80/hr = Q 102,960 directos",
            "Equipo MacBook M1: Q 6,000 | Materiales: Q 800",
            "Reserva de contingencia: Q 20,846 | Reserva de gestión: Q 700",
            "TOTAL: Q 131,306 (incluye reservas)"
        ]),
        ("Interesados Clave", [
            "Director de Carrera: aprueba tema, supervisa defensa, resuelve conflictos",
            "Asesor: orientación técnica, revisión de entregables, aprobación final",
            "Comité evaluador: evalúa anteproyecto y defensa oral",
            "Pablo Flores (PM/Tesista): ejecución total del proyecto, gestión y control"
        ]),
    ]
    
    y_start = Emu(1371600)
    card_w = Emu(5486400)
    card_h = Emu(1463040)
    x1 = Emu(457200)
    x2 = Emu(6400800)
    
    for i, (title, lines) in enumerate(cards_data):
        x = x1 if i % 2 == 0 else x2
        y = y_start + (i // 2) * (card_h + Emu(91440))
        add_card(slide, x, y, card_w, card_h, title, lines, Pt(12), Pt(9))
    
    # Footer note
    txBox = slide.shapes.add_textbox(Emu(457200), Emu(6300000), Emu(10972800), Emu(365760))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Autoridad del PM: gestionar cronograma, presupuesto, y alcance. Escalar decisiones que afecten normativa UNIS al Director."
    p.font.size = Pt(10); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = FONT_BODY
    
    # ── Slide: Registro de Supuestos ──
    slide = next_slide()
    add_header_bar(slide, s(), "REGISTRO DE SUPUESTOS",
                   "Supuestos documentados durante la iniciación | PMBOK 8va, Gobernanza")
    
    assumptions = [
        ("AS-01", "Disponibilidad de asesor", "Se contará con un asesor asignado antes del 16 de febrero 2026", "Alta", "Director de Carrera"),
        ("AS-02", "Hardware adecuado", "MacBook Air M1 8GB es suficiente para fine-tuning LoRA con quantización 4-bit", "Media", "PM/Tesista"),
        ("AS-03", "Datasets públicos", "Existen datasets HDL/QA públicos de calidad suficiente para fine-tuning", "Media", "PM/Tesista"),
        ("AS-04", "Normativa estable", "El reglamento de graduación FING no cambiará durante el proyecto", "Alta", "Director de Carrera"),
        ("AS-05", "Disponibilidad del PM", "El tesista mantendrá disponibilidad 100% durante el período del proyecto", "Alta", "PM/Tesista"),
        ("AS-06", "Acceso a software", "Python, HuggingFace, MLX, y librerías asociadas estarán disponibles sin costo", "Alta", "PM/Tesista"),
        ("AS-07", "Plazos institucionales", "Los tiempos de respuesta institucional no excederán 10 días hábiles", "Media", "FING/UNIS"),
        ("AS-08", "Modelo base accesible", "Qwen3-4B (o equivalente) será accesible en HuggingFace y cargable en 8GB RAM", "Media", "PM/Tesista"),
        ("AS-09", "VerilogEval funcional", "El benchmark VerilogEval v2 será ejecutable en entorno ARM/M1", "Media", "PM/Tesista"),
        ("AS-10", "Sin cambios mayores", "No habrá cambios en el alcance del proyecto una vez aprobado el anteproyecto", "Alta", "PM + Asesor"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(4937760),
        ["ID", "Supuesto", "Descripción", "Certidumbre", "Responsable de Validación"],
        assumptions,
        [Emu(640080), Emu(2103120), Emu(5029200), Emu(914400), Emu(1828800)],
        Pt(10))
    
    txBox = slide.shapes.add_textbox(Emu(274320), Emu(6400000), Emu(11551920), Emu(365760))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Revisión: quincenal durante la iniciación, mensual durante la ejecución. Si un supuesto se invalida, activar proceso de gestión de cambios."
    p.font.size = Pt(10); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = FONT_BODY
    
    # ═══════════════════════════════════════════════════════════
    # DOMAIN 2: INTERESADOS
    # ═══════════════════════════════════════════════════════════
    slide = next_slide()
    add_domain_slide(slide,
        "DOMINIO DE DESEMPEÑO:\nINTERESADOS",
        "PMBOK 8va Edición — Identificación, análisis, involucramiento y comunicación "
        "efectiva con todas las partes interesadas del proyecto.",
        "👥")
    
    # ── Slide: Registro de Interesados ──
    slide = next_slide()
    add_header_bar(slide, s(), "REGISTRO DE INTERESADOS",
                   "Identificar, clasificar y documentar a todos los interesados | PMBOK 8va, Interesados")
    
    stakeholders = [
        ("SH-01", "Pablo Flores", "PM / Tesista", "Interno", "Alto", "Alto", "Gestionar de cerca"),
        ("SH-02", "Asesor (por asignar)", "Orientador técnico", "Interno", "Alto", "Alto", "Gestionar de cerca"),
        ("SH-03", "Director de Carrera FING", "Autoridad académica", "Interno", "Alto", "Medio", "Mantener satisfecho"),
        ("SH-04", "Comité Evaluador", "Aprobación anteproyecto", "Interno", "Alto", "Bajo", "Mantener satisfecho"),
        ("SH-05", "Jurado de Defensa", "Evaluación final", "Interno", "Alto", "Bajo", "Mantener satisfecho"),
        ("SH-06", "Coordinador de Tesis", "Procesos administrativos", "Interno", "Medio", "Medio", "Gestionar de cerca"),
        ("SH-07", "Biblioteca FING/UNIS", "Recepción y archivo", "Interno", "Bajo", "Bajo", "Monitorear"),
        ("SH-08", "Estudiantes FING", "Audiencia académica", "Externo", "Bajo", "Bajo", "Monitorear"),
        ("SH-09", "Comunidad ML/HW", "Investigadores en HDL", "Externo", "Bajo", "Bajo", "Monitorear"),
        ("SH-10", "HuggingFace / GitHub", "Plataformas de publicación", "Externo", "Medio", "Bajo", "Mantener informado"),
        ("SH-11", "UNIS — Rectoría", "Aval institucional", "Externo", "Bajo", "Bajo", "Monitorear"),
        ("SH-12", "Familia Flores", "Apoyo logístico y moral", "Externo", "Bajo", "Medio", "Mantener informado"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(4937760),
        ["ID", "Nombre", "Rol en el Proyecto", "Clasificación", "Poder", "Interés", "Estrategia"],
        stakeholders,
        [Emu(640080), Emu(1828800), Emu(2103120), Emu(1188720), Emu(731520), Emu(731520), Emu(2560320)],
        Pt(9))
    
    # ── Slide: Plan del Involucramiento de los Interesados ──
    slide = next_slide()
    add_header_bar(slide, s(), "PLAN DEL INVOLUCRAMIENTO DE LOS INTERESADOS",
                   "Estrategias por grupo para promover participación productiva | PMBOK 8va, Interesados")
    
    involvement = [
        ("Pablo Flores (PM)", "Líder, ejecutor, controlador", "100% diario", "Reportes de avance, reuniones semanales, actualización cronograma"),
        ("Asesor", "Revisor técnico, mentor", "Reunión semanal 1h + revisión hitos", "Minutas de reunión, borradores para revisión, reportes de avance"),
        ("Director de Carrera", "Aprobador, escalamiento", "Hitos: inicio, anteproyecto, defensa", "Solicitudes formales, notificaciones de hitos, reportes de estado"),
        ("Comité Evaluador", "Evaluador de anteproyecto", "Marzo 2026 (único)", "Documento de anteproyecto impreso y digital"),
        ("Jurado de Defensa", "Evaluador final", "Septiembre 2026 (único)", "Informe final, presentación defensa, demo técnico"),
        ("Coordinador de Tesis", "Gestor de trámites", "Según necesidad administrativa", "Formularios, solicitudes, constancias"),
    ]
    
    y = Emu(1371600)
    for title, role, freq, method in involvement:
        card = add_card(slide, Emu(457200), y, Emu(11277600), Emu(822960), title, [
            f"Rol: {role}",
            f"Frecuencia: {freq}",
            f"Enfoque: {method}"
        ], Pt(11), Pt(9))
        y += Emu(914400)
    
    # ── Slide: Plan de Gestión de las Comunicaciones ──
    slide = next_slide()
    add_header_bar(slide, s(), "PLAN DE GESTIÓN DE LAS COMUNICACIONES",
                   "Políticas para gestionar comunicaciones | PMBOK 8va, Interesados")
    
    comms = [
        ("Reunión semanal asesor", "PM → Asesor", "Semanal (1h)", "Presencial / Zoom", "Minuta en 24h", "Asesor"),
        ("Reporte de avance mensual", "PM → Asesor + Director", "Mensual", "Email + PDF", "EVM: SPI, SV, CPI", "Director, Asesor"),
        ("Notificación de hitos", "PM → Asesor + Director", "Al completar hitos", "Email formal", "Acta o documento del hito", "Director, Asesor"),
        ("Solicitud de cambio", "PM → Asesor", "Según necesidad", "Formato estándar", "Evaluación en 5 días hábiles", "Asesor"),
        ("Notificación de riesgos", "PM → Asesor", "Al detectarse", "Email / WhatsApp", "Registro de riesgos actualizado", "Asesor"),
        ("Entrega de documento", "PM → Asesor / Comité", "Al completar versiones", "PDF + impreso", "Acuse de recibo", "Director, Comité"),
        ("Defensa oral", "PM → Jurado", "Sep 2026 (único)", "Presencial UNIS", "Presentación + demo", "Jurado, Director"),
        ("Publicación en GitHub", "PM → Comunidad", "Post-defensa", "Repositorio público", "README + instrucciones", "Comunidad ML/HW"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(4937760),
        ["Comunicación", "De → Para", "Frecuencia", "Medio / Canal", "Entregable", "Receptor Final"],
        comms,
        [Emu(2286000), Emu(1828800), Emu(1371600), Emu(1828800), Emu(2560320), Emu(1645920)],
        Pt(9))
    
    # ── Slide: Matriz de Comunicaciones ──
    slide = next_slide()
    add_header_bar(slide, s(), "MATRIZ DE COMUNICACIONES",
                   "Detalle de flujos de información entre interesados | PMBOK 8va, Interesados")
    
    matrix_data = [
        ("Reunión semanal", "Presencial/Zoom", "Semanal", "Minuta de reunión", "PM, Asesor", "PM"),
        ("Avance mensual", "Email + PDF", "Mensual", "Informe de avance (EVM)", "PM, Asesor, Director", "PM"),
        ("Alerta de riesgo", "WhatsApp/Email", "Inmediato", "Notificación de riesgo", "PM, Asesor", "Cualquiera"),
        ("Solicitud de cambio", "Formato + Email", "Según necesidad", "RFC (Request for Change)", "PM, Asesor", "PM"),
        ("Convocatoria defensa", "Email formal", "Una vez (Sep 2026)", "Carta de convocatoria", "PM, Director, Jurado", "PM"),
        ("Publicación repo", "GitHub / HuggingFace", "Post-defensa", "README + paper + pesos", "Comunidad global", "PM"),
        ("Consulta normativa", "Email", "Según necesidad", "Solicitud formal", "PM, Coordinador Tesis", "PM"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(4572000),
        ["Tipo", "Medio/Canal", "Frecuencia", "Formato/Entregable", "Audiencia", "Responsable"],
        matrix_data,
        [Emu(1828800), Emu(1828800), Emu(1645920), Emu(2560320), Emu(2743200), Emu(1828800)],
        Pt(10))
    
    # ═══════════════════════════════════════════════════════════
    # DOMAIN 3: ALCANCE
    # ═══════════════════════════════════════════════════════════
    slide = next_slide()
    add_domain_slide(slide,
        "DOMINIO DE DESEMPEÑO:\nALCANCE",
        "PMBOK 8va Edición — Definición, validación y control del alcance del proyecto "
        "y del producto. Incluye requisitos, EDT, y trazabilidad.",
        "🎯")
    
    # ── Slide: Plan de Gestión del Alcance ──
    slide = next_slide()
    add_header_bar(slide, s(), "PLAN DE GESTIÓN DEL ALCANCE",
                   "Define cómo se planificará, gestionará y controlará el alcance | PMBOK 8va")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(1645920),
        "Proceso de Definición del Alcance", [
            "1. Recopilar requisitos: entrevistas con asesor, análisis de normativa FING, revisión de literatura",
            "2. Elaborar Enunciado del Alcance: descripción, entregables, criterios de aceptación, exclusiones",
            "3. Crear EDT/WBS: descomposición jerárquica hasta nivel 4 (paquetes de trabajo de 3-31h)",
            "4. Validar alcance con asesor y Director de Carrera antes de congelar línea base",
        ], Pt(12), Pt(10))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(1645920),
        "Control del Alcance", [
            "• Línea base congelada tras aprobación del anteproyecto (Mar 2026)",
            "• Cualquier cambio requiere Solicitud de Cambio formal → evaluación asesor → aprobación",
            "• Métricas: % de paquetes EDT completados vs planificados",
            "• Verificación: entregables validados contra criterios de aceptación documentados",
        ], Pt(12), Pt(10))
    
    add_card(slide, Emu(457200), Emu(3200400), Emu(10972800), Emu(1645920),
        "Exclusiones del Alcance", [
            "• No incluye desarrollo de hardware (FPGA/ASIC) ni pruebas en silicio real",
            "• No incluye entrenamiento en clusters GPU multi-nodo (limitado a hardware M1 8GB local)",
            "• No incluye implementación de un IDE o herramienta de software interactiva",
            "• No cubre fine-tuning de modelos mayores a 4B parámetros (restricción de RAM)",
            "• No incluye publicación en revistas indexadas (solo repositorio público y defensa oral FING)"
        ], Pt(12), Pt(10))
    
    # ── Slide: Plan de Gestión de Requisitos ──
    slide = next_slide()
    add_header_bar(slide, s(), "PLAN DE GESTIÓN DE REQUISITOS",
                   "Define cómo se recopilarán, analizarán y gestionarán los requisitos | PMBOK 8va")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(1463040),
        "Recopilación de Requisitos", [
            "Fuentes: normativa FING, guía de tesis UNIS, reuniones con asesor, papers de referencia",
            "Técnicas: entrevistas estructuradas con asesor, análisis documental, benchmark de proyectos similares",
            "Clasificación: RF (Funcionales), RNF (No Funcionales), RP (Procedimentales)",
            "Priorización: MoSCoW — Must have (M), Should have (S), Could have (C), Won't have (W)"
        ], Pt(12), Pt(10))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(1463040),
        "Trazabilidad y Validación", [
            "Matriz de trazabilidad: cada requisito vinculado a EDT, entregable y criterio de aceptación",
            "Validación: revisión con asesor en hitos de planificación y cierre de fase",
            "Verificación: scripts de prueba automatizados para requisitos técnicos (RNF-01, RF-01)",
            "Control de cambios: modificaciones a requisitos siguen proceso formal de gestión de cambios"
        ], Pt(12), Pt(10))
    
    add_card(slide, Emu(457200), Emu(3108960), Emu(10972800), Emu(1645920),
        "Métricas y KPIs de Requisitos", [
            "• Cobertura de trazabilidad: 100% de requisitos con EDT, entregable y criterio de aceptación",
            "• Tasa de aprobación: % de requisitos validados por asesor vs total",
            "• Estabilidad: número de cambios a requisitos post-aprobación del anteproyecto (meta: <3)",
            "• Cumplimiento: % de criterios de aceptación satisfechos al cierre de cada fase"
        ], Pt(12), Pt(10))
    
    # ── Slide: Documentación de Requisitos ──
    slide = next_slide()
    add_header_bar(slide, s(), "DOCUMENTACIÓN DE REQUISITOS",
                   "Registro y clasificación de requerimientos | PMBOK 8va, Alcance")
    
    reqs = [
        ("RF-01", "2 configuraciones fine-tuning LoRA", "Comparar chunking vs filtrado", "A", "NE-02: Investigación", "E-04, E-05: Adaptadores .safetensors", "2 archivos .safetensors de pesos LoRA"),
        ("RF-02", "Evaluar con VerilogEval v2 (156 prob.)", "Benchmark estandarizado", "A", "NE-02: Investigación", "E-06: Resultados benchmark", "Script ejecutable + JSON 156 entradas"),
        ("RF-03", "Cuantificar corrupción en datos", "Reproducir hallazgo del paper", "A", "NE-02: Investigación", "E-04, E-06: Scripts verific.", "verify_paper_numbers.py ejecutable"),
        ("RF-04", "Scripts verificación reproducción", "Reproducibilidad científica", "M", "NE-02: Investigación", "E-04: Scripts refinados", "Scripts sin errores en entorno limpio"),
        ("RNF-01", "Hardware real (M1 8GB)", "Condición experimental real", "A", "NE-02: Investigación", "E-05: Modelo final", "Logs con ID hardware confirmado"),
        ("RNF-02", "Informe formato FING + IMRaD", "Requisito de grado", "A", "NE-03: Resultados", "E-07, E-10: Informe final", "Aprobación asesor + tribunal"),
        ("RNF-03", "Repositorio estructura clara", "Divulgación científica", "M", "NE-02: Investigación", "E-04: Repositorio", "README con instrucciones reproducción"),
        ("RNF-04", "Cronograma máximo 360h", "Restricción institucional", "A", "NE-01: Académico", "E-08: Cronograma", "Registro horas ≤ 360 (foco en ejecución)"),
        ("RP-01", "Inscripción + asesor (Feb.)", "Procedimiento FING", "A", "NE-01: Académico", "E-01: Inscripción", "Constancia inscripción + carta asesor"),
        ("RP-02", "Anteproyecto aprobado (Mar.)", "Procedimiento FING", "A", "NE-01: Académico", "E-02: Anteproyecto", "Visto bueno Director + Asesor"),
        ("RP-03", "Avances ejecución (Abr-Ago)", "Procedimiento FING", "A", "NE-01: Académico", "E-03, E-07: Avances", "Documentos entregados y revisados"),
        ("RP-04", "Defensa + entrega final (Sep.)", "Procedimiento FING", "A", "NE-01: Académico", "E-09, E-10: Defensa", "Acta defensa + constancia Biblioteca"),
    ]
    
    add_simple_table(slide, Emu(182880), Emu(1371600), Emu(11887200), Emu(5029200),
        ["ID", "Descripción", "Propósito", "Prioridad", "Origen (NE)", "Entregable", "Criterio de Aceptación"],
        reqs,
        [Emu(594360), Emu(2194560), Emu(1645920), Emu(548640), Emu(1554480), Emu(2377440), Emu(3017520)],
        Pt(8))
    
    txBox = slide.shapes.add_textbox(Emu(182880), Emu(6480000), Emu(11887200), Emu(274320))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Prioridad: A=Alta, M=Media, B=Baja | Cobertura: 12 requisitos trazados a entregables y criterios de aceptación. Trazabilidad bidireccional 100%."
    p.font.size = Pt(9); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = FONT_BODY
    
    # ── Slide: Enunciado del Alcance del Proyecto ──
    slide = next_slide()
    add_header_bar(slide, s(), "ENUNCIADO DEL ALCANCE DEL PROYECTO",
                   "Descripción detallada del alcance, entregables, criterios de aceptación y exclusiones | PMBOK 8va")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(10972800), Emu(1280160),
        "Descripción del Alcance del Producto", [
            "Investigación científica documentada sobre el impacto del preprocesamiento de datos de entrenamiento en la degradación catastrófica de fine-tuning LoRA para generación de HDL.",
            "Dos configuraciones experimentales: (a) filtrado de datasets por calidad, (b) chunking estratégico del corpus de entrenamiento.",
            "Evaluación comparativa usando VerilogEval v2 (156 problemas) con métricas BLEU, exactitud sintáctica, y PASS@k."
        ], Pt(12), Pt(10))
    
    add_card(slide, Emu(457200), Emu(2834640), Emu(5486400), Emu(1828800),
        "Entregables del Proyecto", [
            "E-01: Inscripción formal del tema (Feb 2026)",
            "E-02: Anteproyecto aprobado (Mar 2026)",
            "E-03: Avances de ejecución (Abr-Ago 2026)",
            "E-04: Scripts y repositorio (Ago 2026)",
            "E-05: Modelos LoRA fine-tuned (.safetensors)",
            "E-06: Resultados de benchmark (JSON + gráficas)",
            "E-07: Informe final formato FING/IMRaD (Ago 2026)",
            "E-08: Cronograma y registro de horas",
            "E-09: Defensa oral (Sep 2026)",
            "E-10: Entrega final a Biblioteca FING (Sep 2026)"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(2834640), Emu(5029200), Emu(914400),
        "Criterios de Aceptación", [
            "Anteproyecto: aprobado por comité evaluador",
            "Informe final: aprobado por asesor y tribunal",
            "Repositorio: README + instrucciones de reproducción",
            "Scripts: ejecutables sin errores en entorno limpio"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(3931920), Emu(5029200), Emu(914400),
        "Restricciones", [
            "Hardware: MacBook Air M1 8GB RAM (no GPU dedicada)",
            "Tiempo: máximo 360 horas de esfuerzo efectivo",
            "Presupuesto: Q 6,800 máximo en no-humanos",
            "Alcance: 2 configuraciones LoRA comparativas"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(5029200), Emu(10972800), Emu(822960),
        "Exclusiones Explícitas", [
            "No incluye desarrollo de hardware (FPGA/ASIC) | No incluye implementación de herramienta de software | No cubre modelos >4B parámetros",
            "No incluye publicación en revistas indexadas | No incluye entrenamiento multi-GPU | No cubre más de 2 configuraciones experimentales"
        ], Pt(11), Pt(9))
    
    # ── Slide: EDT / WBS ──
    slide = next_slide()
    add_header_bar(slide, s(), "EDT / WBS — ESTRUCTURA DE DESGLOSE DEL TRABAJO",
                   "Descomposición jerárquica del alcance total | PMBOK 8va, Alcance")
    
    # EDT diagram as text-based cards
    wbs_levels = [
        ("Nivel 1", "Proyecto de Graduación LoRA-HDL-QA", "1,287h", NAVY),
        ("Nivel 2", "INICIO (148h) | PLANIFICACIÓN (235h) | EJECUCIÓN (744h) | M&C (86h) | CIERRE (160h)", "", BLUE),
    ]
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(10972800), Emu(640080),
        "Nivel 2 — Fases del Proyecto", [
            "1.1 INICIO (148h): Inscripción + Asesor + Plan del Proyecto  |  1.2 PLANIFICACIÓN (235h): Anteproyecto + Diseño Metodológico",
            "1.3 EJECUCIÓN (744h): Marco Teórico + Preprocesamiento + Entrenamiento LoRA + Evaluación + Análisis + Informe",
            "1.4 MONITOREO Y CONTROL (86h): Seguimiento + Control de Cambios  |  1.5 CIERRE (160h): Correcciones + Defensa + Entrega"
        ], Pt(11), Pt(9))
    
    # Nivel 3 - Inicio
    add_card(slide, Emu(457200), Emu(2194560), Emu(2103120), Emu(2011680),
        "1.1 INICIO (148h)", [
            "1.1.1 Inscripción (78h)",
            "1.1.2 Asignación Asesor (80h)",
            "1.1.3 Plan del Proyecto (60h)"
        ], Pt(10), Pt(8))
    
    add_card(slide, Emu(2743200), Emu(2194560), Emu(2103120), Emu(2011680),
        "1.2 PLANIFICACIÓN (235h)", [
            "1.2.1 Anteproyecto (152h)",
            "1.2.2 Diseño Metodológico (83h)"
        ], Pt(10), Pt(8))
    
    add_card(slide, Emu(5029200), Emu(2194560), Emu(2103120), Emu(2011680),
        "1.3 EJECUCIÓN (744h)", [
            "1.3.1 Marco Teórico (161h)",
            "1.3.2 Preprocesamiento (136h)",
            "1.3.3 Entrenamiento LoRA (227h)",
            "1.3.4 Evaluación Benchmark (106h)",
            "1.3.5 Análisis Resultados (90h)",
            "1.3.6 Redacción Informe (175h)"
        ], Pt(10), Pt(8))
    
    add_card(slide, Emu(7315200), Emu(2194560), Emu(2103120), Emu(2011680),
        "1.4 M&C (86h)", [
            "1.4.1 Seguimiento (86h)",
            "1.4.2 Control de Cambios (68h)"
        ], Pt(10), Pt(8))
    
    add_card(slide, Emu(9601200), Emu(2194560), Emu(2103120), Emu(2011680),
        "1.5 CIERRE (160h)", [
            "1.5.1 Correcciones Finales (75h)",
            "1.5.2 Defensa Pública (59h)",
            "1.5.3 Entrega y Cierre (26h)"
        ], Pt(10), Pt(8))
    
    add_card(slide, Emu(457200), Emu(4389120), Emu(10972800), Emu(640080),
        "Nivel 4 — Paquetes de Trabajo", [
            "16 paquetes de trabajo desglosados en 110 actividades nivel 4. Rango: 3h a 31h cada una. Total: 1,287h planificadas.",
            "Ver Diccionario de la EDT para detalle de cada paquete."
        ], Pt(11), Pt(9))
    
    # ── Slide: Diccionario de la EDT ──
    slide = next_slide()
    add_header_bar(slide, s(), "DICCIONARIO DE LA EDT",
                   "Detalle de cada paquete de trabajo | PMBOK 8va, Alcance")
    
    dict_data = [
        ("1.1.1", "Inscripción del Tema", "Formalizar el tema de tesis ante FING", "78h", "Constancia de inscripción firmada", "Reglamento FING vigente"),
        ("1.1.2", "Asignación de Asesor", "Identificar y formalizar asesor", "80h", "Carta de asignación de asesor", "Disponibilidad docente FING"),
        ("1.1.3", "Plan del Proyecto", "Documentar plan de dirección completo", "60h", "Plan aprobado por asesor", "PMBOK 8va Edición"),
        ("1.2.1", "Anteproyecto", "Redactar y presentar anteproyecto", "152h", "Anteproyecto aprobado por comité", "Guía de anteproyecto FING"),
        ("1.2.2", "Diseño Metodológico", "Diseñar experimentos y métricas", "83h", "Documento metodológico aprobado", "Papers de referencia LoRA/HDL"),
        ("1.3.1", "Marco Teórico", "Fundamentos LLM, LoRA, HDL, QA", "161h", "Capítulo 2 del informe final", "Artículos científicos revisados"),
        ("1.3.2", "Preprocesamiento de Datos", "Limpiar, tokenizar y validar datasets", "136h", "Datasets procesados y documentados", "Datasets HDL/QA públicos"),
        ("1.3.3", "Entrenamiento LoRA", "Fine-tuning con 2 configuraciones", "227h", "2 adaptadores LoRA (.safetensors)", "MacBook M1 8GB, quantización 4-bit"),
        ("1.3.4", "Evaluación Benchmark", "Ejecutar VerilogEval v2 y comparar", "106h", "Resultados benchmark (JSON + gráficas)", "VerilogEval v2, scripts evaluación"),
        ("1.3.5", "Análisis de Resultados", "Analizar y comparar métricas", "90h", "Capítulo 4 del informe final", "Datos de benchmark recopilados"),
        ("1.3.6", "Redacción del Informe", "Redactar informe en formato IMRaD", "175h", "Informe completo revisado", "Normativa UNIS, guía IMRaD"),
        ("1.4.1", "Seguimiento", "Monitorear avance y reuniones", "86h", "Reportes de avance mensuales", "Cronograma base, minutas"),
        ("1.4.2", "Control de Cambios", "Gestionar solicitudes de cambio", "68h", "Log de cambios aprobados", "Plan de gestión de cambios"),
        ("1.5.1", "Correcciones Finales", "Aplicar correcciones del asesor", "75h", "Documento final aprobado", "Retroalimentación del asesor"),
        ("1.5.2", "Defensa Pública", "Preparar y realizar defensa oral", "59h", "Acta de defensa aprobatoria", "Reglamento de defensa FING"),
        ("1.5.3", "Entrega y Cierre", "Entregar a biblioteca y cerrar", "26h", "Constancia de Biblioteca FING", "Normas de entrega UNIS"),
    ]
    
    add_simple_table(slide, Emu(137160), Emu(1371600), Emu(11902440), Emu(5029200),
        ["EDT", "Paquete de Trabajo", "Descripción", "Horas", "Entregable", "Supuestos"],
        dict_data,
        [Emu(731520), Emu(2103120), Emu(2743200), Emu(640080), Emu(2743200), Emu(2743200)],
        Pt(8))
    
    # ── Slide: Matriz de Trazabilidad de Requerimientos ──
    slide = next_slide()
    add_header_bar(slide, s(), "MATRIZ DE TRAZABILIDAD DE REQUERIMIENTOS",
                   "Trazabilidad bidireccional requisitos → EDT → entregables → criterios | PMBOK 8va, Alcance")
    
    trace_data = [
        ("RF-01", "2 configuraciones fine-tuning LoRA", "A", "NE-02", "1.3.2, 1.3.3", "E-04, E-05", "2 adaptadores con pesos .safetensors"),
        ("RF-02", "Evaluar con VerilogEval v2 (156 prob.)", "A", "NE-02", "1.3.4", "E-06", "Script benchmark; JSON 156 entradas"),
        ("RF-03", "Cuantificar corrupción en datos", "A", "NE-02", "1.3.2, 1.3.5", "E-04, E-06", "verify_paper_numbers.py ejecutable"),
        ("RF-04", "Scripts verificación reproducción", "M", "NE-02", "1.3.6", "E-04", "Scripts sin errores en entorno limpio"),
        ("RNF-01", "Hardware real (M1 8GB)", "A", "NE-02", "1.3.3", "E-05", "Logs con ID hardware"),
        ("RNF-02", "Informe formato FING + IMRaD", "A", "NE-03", "1.3.6, 1.5.1", "E-07, E-10", "Aprobación asesor y tribunal"),
        ("RNF-03", "Repositorio estructura clara", "M", "NE-02", "1.3.6", "E-04", "README + instrucciones reproducción"),
        ("RNF-04", "Cronograma máximo 360h", "A", "NE-01", "1.4.1", "E-08", "Registro horas ≤ 360 (foco ejecución)"),
        ("RP-01", "Inscripción + asesor (Feb.)", "A", "NE-01", "1.1.1, 1.1.2", "E-01", "Constancia inscripción + carta asesor"),
        ("RP-02", "Anteproyecto aprobado (Mar.)", "A", "NE-01", "1.2.1", "E-02", "Visto bueno Director + Asesor"),
        ("RP-03", "Avances ejecución (Abr-Ago)", "A", "NE-01", "1.3.1, 1.3.6", "E-03, E-07", "Docs entregados y revisados"),
        ("RP-04", "Defensa + entrega final (Sep.)", "A", "NE-01", "1.5.2, 1.5.3", "E-09, E-10", "Acta defensa + constancia Biblioteca"),
    ]
    
    add_simple_table(slide, Emu(182880), Emu(1371600), Emu(11887200), Emu(4572000),
        ["ID", "Descripción", "P", "Origen", "EDT", "Entregable", "Criterio de Aceptación"],
        trace_data,
        [Emu(548640), Emu(2286000), Emu(365760), Emu(914400), Emu(1645920), Emu(1645920), Emu(3657600)],
        Pt(8))
    
    txBox = slide.shapes.add_textbox(Emu(182880), Emu(6030000), Emu(11887200), Emu(274320))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Cobertura: 12 requerimientos trazados a entregables y criterios de aceptación. Trazabilidad bidireccional 100%."
    p.font.size = Pt(9); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = FONT_BODY
    
    # ── Slide: Plan de Gestión de Cambios ──
    slide = next_slide()
    add_header_bar(slide, s(), "PLAN DE GESTIÓN DE CAMBIOS",
                   "Proceso formal para solicitar, evaluar y aprobar cambios | PMBOK 8va, Alcance")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2011680),
        "Proceso de Control de Cambios", [
            "1. IDENTIFICAR: Cualquier interesado detecta necesidad de cambio en alcance, tiempo, costo o calidad",
            "2. REGISTRAR: Completar formato de Solicitud de Cambio (RFC) con: descripción, justificación, impacto estimado",
            "3. EVALUAR: PM analiza impacto en triple restricción (alcance/cronograma/costo) y riesgos asociados",
            "4. REVISAR: Asesor revisa RFC en máximo 5 días hábiles, puede solicitar más información",
            "5. DECIDIR: Aprobar (implementar), Rechazar (documentar razón), o Diferir (re-evaluar en próxima fase)",
            "6. IMPLEMENTAR: Actualizar línea base, EDT, cronograma, y documentos afectados. Comunicar a interesados."
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2011680),
        "Autoridad de Aprobación", [
            "• Cambios ≤5% desviación: PM puede aprobar directamente (informar al asesor)",
            "• Cambios 5-15% desviación: Requiere aprobación del asesor con revisión técnica",
            "• Cambios >15% desviación o que afecten normativa: Escalar al Director de Carrera",
            "• Emergencias (riesgo crítico): PM implementa respuesta inmediata, documenta post-facto en 48h",
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(3566160), Emu(10972800), Emu(822960),
        "Herramientas y Registros", [
            "Formato RFC en Excel/Google Sheets vinculado al cronograma | Log de cambios con trazabilidad a EDT |"
            " Revisión en reuniones quincenales | Actualización de línea base solo tras aprobación formal"
        ], Pt(11), Pt(9))
    
    # ── Slide: Solicitud de Cambio (Formato) ──
    slide = next_slide()
    add_header_bar(slide, s(), "SOLICITUD DE CAMBIO (RFC) — FORMATO",
                   "Formato estándar para solicitudes de cambio | PMBOK 8va, Alcance (opcional según normativa)")
    
    rfc_fields = [
        ("RFC ID", "RFC-NNN (auto-incremental)", "Fecha Solicitud", "DD/MM/YYYY"),
        ("Solicitante", "Nombre y rol", "Prioridad", "Alta / Media / Baja"),
        ("Descripción del Cambio", "", "", ""),
        ("Justificación", "¿Por qué es necesario el cambio?", "", ""),
        ("Impacto en Alcance", "Paquetes EDT afectados, entregables modificados", "", ""),
        ("Impacto en Cronograma", "Δ horas estimadas, nueva fecha de finalización", "", ""),
        ("Impacto en Costo", "Δ costo estimado en Q", "", ""),
        ("Impacto en Riesgos", "Nuevos riesgos o modificación de existentes", "", ""),
        ("Alternativas Consideradas", "Otras opciones evaluadas y razón de descarte", "", ""),
        ("Decisión", "Aprobado / Rechazado / Diferido", "Fecha Decisión", "DD/MM/YYYY"),
        ("Firmas", "Solicitante: ______  |  PM: ______  |  Asesor: ______", "", ""),
    ]
    
    # Create a visual RFC form
    y = Emu(1371600)
    for label, val1, label2, val2 in rfc_fields:
        # Label
        txBox = slide.shapes.add_textbox(Emu(457200), y, Emu(2743200), Emu(320040))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = f"{label}:"; p.font.size = Pt(10); p.font.bold = True
        p.font.color.rgb = NAVY; p.font.name = FONT_BODY
        
        # Value
        txBox2 = slide.shapes.add_textbox(Emu(3200400), y, Emu(3108960), Emu(320040))
        tf2 = txBox2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]; p2.text = val1; p2.font.size = Pt(10)
        p2.font.color.rgb = MED_TEXT; p2.font.name = FONT_BODY
        
        if label2:
            txBox3 = slide.shapes.add_textbox(Emu(6400800), y, Emu(2286000), Emu(320040))
            tf3 = txBox3.text_frame; tf3.word_wrap = True
            p3 = tf3.paragraphs[0]; p3.text = f"{label2}:"; p3.font.size = Pt(10); p3.font.bold = True
            p3.font.color.rgb = NAVY; p3.font.name = FONT_BODY
            
            txBox4 = slide.shapes.add_textbox(Emu(8686800), y, Emu(2743200), Emu(320040))
            tf4 = txBox4.text_frame; tf4.word_wrap = True
            p4 = tf4.paragraphs[0]; p4.text = val2; p4.font.size = Pt(10)
            p4.font.color.rgb = MED_TEXT; p4.font.name = FONT_BODY
        
        y += Emu(365760)
    
    # ── Slide: Plan de Estrategia de Adquisiciones ──
    slide = next_slide()
    add_header_bar(slide, s(), "PLAN DE ESTRATEGIA DE ADQUISICIONES",
                   "Estrategia para adquirir recursos externos | PMBOK 8va, Alcance")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(10972800), Emu(2743200),
        "Estrategia General de Adquisiciones", [
            "Contexto: Proyecto académico individual de graduación. La naturaleza del proyecto (investigación + desarrollo experimental) "
            "implica que la mayoría de los recursos son internos (esfuerzo del tesista, hardware propio, software open-source).",
            "",
            "Adquisiciones requeridas:",
            "• Equipo: MacBook Air M1 2020 — ya adquirido (costo hundido Q 6,000). Alternativa: usar laboratorio FING si falla hardware.",
            "• Software: Python, HuggingFace, MLX, VerilogEval — todo open-source sin costo de licencia.",
            "• Datasets: Repositorios públicos (HuggingFace datasets, GitHub) — descarga gratuita.",
            "• Materiales: Impresión y empastado (3 copias × Q 150 = Q 450) — proveedor local, cotización simple.",
            "• Papelería y misceláneos: Q 350 — compra directa sin licitación.",
            "",
            "Estrategia: Compra directa por bajo monto (< Q 5,000). No se requiere proceso de licitación formal por ser proyecto académico. "
            "El PM tiene autoridad para adquirir materiales menores. Equipamiento se usa recurso propio existente."
        ], Pt(11), Pt(9))
    
    # ── Slide: Adquisiciones — Decisiones de Internalización/Externalización ──
    slide = next_slide()
    add_header_bar(slide, s(), "ADQUISICIONES — DECISIONES DE INTERNALIZACIÓN O EXTERNALIZACIÓN",
                   "Análisis Make-or-Buy | PMBOK 8va, Alcance")
    
    mob_data = [
        ("Entrenamiento LoRA", "Interno (Tesista)", "N/A", "Conocimiento técnico del tesista; hardware propio; control total del proceso"),
        ("Evaluación Benchmark", "Interno (Tesista)", "N/A", "Scripts propios; ejecución local; reproducibilidad científica requerida"),
        ("Redacción de informe", "Interno (Tesista)", "N/A", "Requisito de grado; debe ser trabajo original del estudiante"),
        ("Revisión técnica", "Externo (Asesor)", "Honorarios institucionales FING", "Asesor designado por FING; costo cubierto por la universidad"),
        ("Impresión y empastado", "Externo (Proveedor)", "Q 450", "Servicio especializado; bajo costo; 3 copias requeridas por normativa"),
        ("Corrección de estilo", "Interno (Tesista)", "N/A", "Cubierto en fase de correcciones finales (EDT 1.5.1); sin costo adicional"),
        ("Software y herramientas", "Externo (Open-source)", "Gratuito", "Python, HuggingFace, MLX: licencias libres; descarga directa"),
        ("Hardware computacional", "Interno (Propio)", "Q 6,000 (costo hundido)", "MacBook M1 ya adquirido; si falla, alternativa: laboratorio FING"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(4114800),
        ["Recurso/Actividad", "Decisión", "Costo (Q)", "Justificación"],
        mob_data,
        [Emu(2560320), Emu(2103120), Emu(2103120), Emu(4572000)],
        Pt(10))
    
    # ── Slide: Adquisiciones — Criterios de Selección de Fuentes ──
    slide = next_slide()
    add_header_bar(slide, s(), "ADQUISICIONES — CRITERIOS DE SELECCIÓN DE FUENTES",
                   "Criterios para evaluar y seleccionar proveedores | PMBOK 8va, Alcance")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2194560),
        "Criterios de Selección — Proveedores", [
            "Para servicios de impresión y empastado:",
            "1. Calidad: muestras de trabajos previos (portadas, encuadernación)",
            "2. Precio: cotización ≤ Q 200 por copia completa",
            "3. Tiempo de entrega: ≤ 3 días hábiles por copia",
            "4. Ubicación: cercanía a campus UNIS o zona 15",
            "5. Referencias: recomendaciones de tesistas previos FING",
            "",
            "Proveedores evaluados (ejemplo):",
            "• Copistería UNIS: Q 150/copia, 2 días, en campus",
            "• PrintCenter Z15: Q 180/copia, 1 día, a 2km",
            "• TesisExpress Online: Q 200/copia, 5 días, envío incluido",
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2194560),
        "Criterios de Selección — Recursos Digitales", [
            "Para datasets y modelos:",
            "1. Licencia: open-source (MIT, Apache 2.0, CC-BY) verificable",
            "2. Calidad: documentación, papers asociados, benchmarks conocidos",
            "3. Comunidad: actividad en GitHub, issues, actualizaciones recientes",
            "4. Compatibilidad: formato compatible con HuggingFace datasets/datasets",
            "5. Tamaño: factible de procesar en M1 8GB (≤ 1GB por dataset)",
            "",
            "Datasets seleccionados:",
            "• VerilogEval v2 (156 problemas, GitHub, MIT license)",
            "• HDL dataset de HuggingFace (código VHDL/Verilog etiquetado)",
            "• Dataset sintético generado con GPT-4 (como respaldo)"
        ], Pt(11), Pt(9))
    
    # ═══════════════════════════════════════════════════════════
    # DOMAIN 4: CRONOGRAMA
    # ═══════════════════════════════════════════════════════════
    slide = next_slide()
    add_domain_slide(slide,
        "DOMINIO DE DESEMPEÑO:\nCRONOGRAMA",
        "PMBOK 8va Edición — Planificación, desarrollo, monitoreo y control del cronograma. "
        "Incluye definición de actividades, secuenciación, estimación de duraciones y línea base.",
        "📅")
    
    # ── Slide: Plan de Gestión del Cronograma ──
    slide = next_slide()
    add_header_bar(slide, s(), "PLAN DE GESTIÓN DEL CRONOGRAMA",
                   "Políticas, procedimientos y documentación para gestionar el cronograma | PMBOK 8va")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(1280160),
        "Metodología y Herramientas", [
            "Método: Bottom-up estimating desde EDT nivel 4 (110 actividades)",
            "Herramienta: ProjectLibre (CPM + baseline + exportación Gantt)",
            "Calendario: Guatemala 2026 — lun-sáb 8h/día (08:00-12:00, 13:00-17:00)",
            "Unidades: Horas de trabajo efectivas. Precisión ±10% en nivel 4 EDT.",
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(2834640), Emu(5486400), Emu(1097280),
        "Umbrales de Control", [
            "Verde: <5% de desviación — normal, sin acción",
            "Amarillo: 5-15% desviación — plan de recuperación en 1 semana",
            "Rojo: >15% desviación — escalamiento al Director de Carrera",
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(1280160),
        "Formato de Reporte", [
            "S-curve mensual (PV, EV, AC) para seguimiento visual",
            "EVM: SPI, SV para medir eficiencia del cronograma",
            "Actualización: quincenal con asesor, registro de cambios en log",
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(2834640), Emu(5029200), Emu(1097280),
        "Revisión y Aprobación", [
            "Quincenal con asesor en reuniones de seguimiento",
            "Registro de cambios en log de cronograma",
            "Línea base congelada post-aprobación anteproyecto (Mar 2026)",
        ], Pt(11), Pt(9))
    
    # ── Slide: Línea Base del Cronograma ──
    slide = next_slide()
    add_header_bar(slide, s(), "LÍNEA BASE DEL CRONOGRAMA",
                   "Cronograma aprobado contra el cual se mide el desempeño | PMBOK 8va")
    
    baseline = [
        ("INICIO", "Feb 02 → Feb 26", "148h", "Inscripción, asesor, plan de proyecto", "Feb 02, 2026", "Feb 26, 2026"),
        ("PLANIFICACIÓN", "Feb 26 → Abr 08", "235h", "Anteproyecto, diseño metodológico", "Feb 26, 2026", "Abr 08, 2026"),
        ("EJECUCIÓN", "Abr 08 → Ago 17", "744h", "Marco teórico, entrenamiento, evaluación, informe", "Abr 08, 2026", "Ago 17, 2026"),
        ("MONITOREO Y CONTROL", "Feb 26 → Mar 13", "86h*", "Seguimiento y control de cambios", "Feb 26, 2026", "Mar 13, 2026*"),
        ("CIERRE", "Ago 17 → Sep 14", "160h", "Correcciones, defensa, entrega", "Ago 17, 2026", "Sep 14, 2026"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(2743200),
        ["Fase", "Período", "Horas (LB)", "Alcance", "Inicio LB", "Fin LB"],
        baseline,
        [Emu(1645920), Emu(2011680), Emu(914400), Emu(3657600), Emu(1828800), Emu(1828800)],
        Pt(10))
    
    add_card(slide, Emu(457200), Emu(4297680), Emu(10972800), Emu(2011680),
        "Notas sobre la Línea Base", [
            "• La línea base del cronograma es la versión aprobada del cronograma que se utiliza para comparar el desempeño real.",
            "• Fecha de congelamiento: posterior a la aprobación del anteproyecto por el comité evaluador (~Mar 25, 2026).",
            "• La línea base solo se modifica mediante solicitudes de cambio aprobadas formalmente.",
            "• * M&C tiene actividades distribuidas a lo largo del proyecto; el bloque principal está en Feb-Mar.",
            "• Duración total del proyecto: 225 días calendario (Feb 02 → Sep 14, 2026). Horas efectivas: 1,287h.",
            "• Se utiliza calendario Guatemala 2026 (lun-sáb 8h/día). Domingos y feriados no laborables."
        ], Pt(11), Pt(9))
    
    # ── Slide: Cronograma del Proyecto ──
    slide = next_slide()
    add_header_bar(slide, s(), "CRONOGRAMA DEL PROYECTO",
                   "Resumen gráfico de fases y duraciones | Feb — Sep 2026 | PMBOK 8va")
    
    # Visual timeline with boxes
    phases = [
        ("INICIO", "Feb 02 → Feb 26", "156h", Emu(365760)),
        ("PLANIFICACIÓN", "Feb 27 → Abr 08", "248h", Emu(2651760)),
        ("EJECUCIÓN", "Abr 08 → Ago 17", "784h", Emu(4937760)),
        ("MONITOREO Y CONTROL", "Feb 27 → Mar 13", "88h", Emu(7223760)),
        ("CIERRE", "Ago 17 → Sep 14", "160h", Emu(9509760)),
    ]
    
    for name, period, hours, x in phases:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(1371600), Emu(2103120), Emu(2560320))
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_BLUE if name != "MONITOREO Y CONTROL" else BLUE_LIGHT
        box.line.fill.background()
        
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(16); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = FONT_TITLE; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = period; p2.font.size = Pt(11)
        p2.font.color.rgb = WHITE; p2.font.name = FONT_BODY; p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph(); p3.text = hours; p3.font.size = Pt(11)
        p3.font.color.rgb = WHITE; p3.font.name = FONT_BODY; p3.alignment = PP_ALIGN.CENTER
    
    # Timeline arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(365760), Emu(4389120), Emu(10972800), Emu(91440))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = NAVY; arrow.line.fill.background()
    
    add_card(slide, Emu(457200), Emu(4754880), Emu(10972800), Emu(1097280),
        "Resumen de Cronograma", [
            "Inicio del proyecto: 2 de febrero de 2026 | Fin del proyecto: 14 de septiembre de 2026 | Duración total: 225 días calendario",
            "Horas efectivas planificadas: 1,287h de trabajo (incluye períodos de espera institucionales: ~155h)",
            "El cronograma detallado está desarrollado en ProjectLibre con 150 tareas (110 de nivel 4 + 40 de resumen/hitos).",
            "Las fechas y duraciones provienen del archivo XML LoRA-HDL-QA_Proyecto_2026.xml, validado y cargable en ProjectLibre 1.9.8."
        ], Pt(11), Pt(9))
    
    # ── Slide: Datos del Cronograma ──
    slide = next_slide()
    add_header_bar(slide, s(), "DATOS DEL CRONOGRAMA",
                   "Métricas clave y datos de soporte del modelo de cronograma | PMBOK 8va")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(1645920),
        "Estadísticas del Cronograma", [
            "Total de actividades en el modelo: 150",
            "Actividades de trabajo (nivel 4): 110",
            "Actividades de resumen (niveles 1-3): 27",
            "Hitos (duración 0): 13",
            "Duración total: 1,287 horas efectivas",
            "Período: 2026-02-02 al 2026-09-14 (225 días calendario)",
            "Método de planificación: Ruta Crítica (CPM)",
            "Calendario aplicado: Guatemala 2026 (lun-sáb, 8h/día)",
            "Software: ProjectLibre 1.9.8"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(1645920),
        "Ruta Crítica", [
            "La ruta crítica atraviesa las fases de:",
            "INICIO → Inscripción → Asignación Asesor → Plan del Proyecto",
            "→ PLANIFICACIÓN → Anteproyecto → Diseño Metodológico",
            "→ EJECUCIÓN → Entrenamiento LoRA → Evaluación → Informe",
            "→ CIERRE → Correcciones → Defensa → Entrega",
            "",
            "Holgura total estimada: ~10% (principalmente en M&C)",
            "Los períodos de espera institucionales son parte de la ruta crítica"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(3200400), Emu(10972800), Emu(1097280),
        "Períodos de Espera (no laborables pero incluidos en el cronograma)", [
            "Resolución institucional (UID 8): 39h | Aprobación comité (UID 33): 31h | Retroalimentación 1ra (UID 100): 23h | Retroalimentación 2da (UID 102): 39h",
            "Resolución jurado (UID 127): 23h | Total esperas institucionales: ~155h (12% del total de horas)"
        ], Pt(11), Pt(9))
    
    # ── Slide: Calendarios del Proyecto ──
    slide = next_slide()
    add_header_bar(slide, s(), "CALENDARIOS DEL PROYECTO",
                   "Base de trabajo y excepciones | PMBOK 8va, Cronograma")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(1645920),
        "Calendario Base — Guatemala 2026", [
            "Laborables: Lunes a Sábado",
            "Horario: 08:00-12:00, 13:00-17:00 (8h/día efectivas)",
            "Domingo: No laborable",
            "Semana Santa: Jueves 02 — Domingo 05 Abril (no laborable)",
            "Día del Trabajo: Viernes 01 Mayo (no laborable)",
            "Día de la Independencia: Martes 15 Septiembre (no laborable)",
            "Nota: Feriados omitidos del XML por bug ProjectLibre 1.9.8 con <Exceptions>"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(3200400), Emu(5486400), Emu(1463040),
        "Recurso — Pablo Flores", [
            "Disponibilidad: 100% (estudiante tiempo completo)",
            "Máximo: 8h/día laborable estándar",
            "Flexibilidad: hasta 10h/día en semanas críticas (pre-defensa, pre-entrega)",
            "Total horas planificadas: 1,287h en 225 días calendario",
            "Promedio: ~5.7h/día efectivas (incluye esperas y buffers)"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(1645920),
        "Períodos de Espera (incluidos como tareas)", [
            "Resolución institucional: 39h (Feb 6-13)",
            "Aprobación comité: 31h (Mar 23-27)",
            "Retroalimentación 1ra revisión: 23h (Ago 7-12)",
            "Retroalimentación 2da revisión: 39h (Ago 17-24)",
            "Resolución jurado: 23h (Sep 4-9)"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(3200400), Emu(5029200), Emu(1463040),
        "Impacto en el Cronograma", [
            "Esperas institucionales = 12% del total de horas planificadas",
            "Buffer implícito por feriados no modelados: ~3 días",
            "Margen de seguridad recomendado: mantener 10% de holgura en hitos críticos",
            "Plan de contingencia: si espera > 2 semanas, escalar al Director"
        ], Pt(11), Pt(9))
    
    # ── Slide: Lista de Actividades ──
    slide = next_slide()
    add_header_bar(slide, s(), "LISTA DE ACTIVIDADES",
                   "Actividades nivel 4 del cronograma (110 actividades) | PMBOK 8va, Cronograma")
    
    # Sample of activities (top 30 of 110)
    activities = [
        ("4", "INICIO", "Revisar normativo y reglamento de graduación", "8h", "2026-02-02", "2026-02-02"),
        ("5", "INICIO", "Delimitar el tema de investigación", "16h", "2026-02-03", "2026-02-04"),
        ("6", "INICIO", "Preparar formularios y documentación", "8h", "2026-02-05", "2026-02-05"),
        ("7", "INICIO", "Presentar solicitud formal ante la facultad", "3h", "2026-02-06", "2026-02-06"),
        ("8", "INICIO", "Período de espera — resolución institucional", "39h", "2026-02-06", "2026-02-13"),
        ("11", "INICIO", "Identificar y evaluar candidatos a asesor", "8h", "2026-02-03", "2026-02-03"),
        ("17", "INICIO", "Definir alcance, supuestos y restricciones", "15h", "2026-02-17", "2026-02-18"),
        ("18", "INICIO", "Elaborar EDT (WBS) del proyecto", "15h", "2026-02-18", "2026-02-20"),
        ("19", "INICIO", "Elaborar cronograma y línea base de tiempo", "15h", "2026-02-20", "2026-02-24"),
        ("25", "PLANIF.", "Revisión de literatura sobre el problema", "23h", "2026-02-26", "2026-03-03"),
        ("37", "PLANIF.", "Diseñar arquitectura general del experimento", "23h", "2026-03-31", "2026-04-03"),
        ("44", "EJEC.", "Rev. bibliográfica: LLMs y arquitecturas Transformer", "23h", "2026-04-08", "2026-04-13"),
        ("54", "EJEC.", "Catalogar datasets HDL / QA disponibles", "15h", "2026-04-08", "2026-04-10"),
        ("63", "EJEC.", "Configurar entorno Python / HuggingFace en MacBook", "15h", "2026-05-07", "2026-05-08"),
        ("68", "EJEC.", "Ejecutar entrenamiento piloto (sanity check)", "15h", "2026-05-19", "2026-05-21"),
        ("73", "EJEC.", "Ejecutar entrenamiento final (modelo definitivo)", "31h", "2026-06-05", "2026-06-11"),
        ("80", "EJEC.", "Ejecutar benchmark — modelo LoRA fine-tuned", "23h", "2026-06-22", "2026-06-24"),
        ("95", "EJEC.", "Redactar cap. 3 — Metodología", "23h", "2026-07-17", "2026-07-22"),
        ("106", "M&C", "Planificar agenda y protocolo de reuniones", "7h", "2026-02-26", "2026-02-27"),
        ("111", "M&C", "Registrar solicitudes de cambio formales", "15h", "2026-02-27", "2026-03-03"),
        ("117", "CIERRE", "Recibir retroalimentación final del asesor", "7h", "2026-08-17", "2026-08-18"),
        ("123", "CIERRE", "Preparar presentación PowerPoint para defensa", "15h", "2026-08-31", "2026-09-02"),
        ("126", "CIERRE", "Defensa oral ante el jurado evaluador", "7h", "2026-09-03", "2026-09-04"),
        ("130", "CIERRE", "Imprimir, empastar y preparar copias finales", "4h", "2026-09-11", "2026-09-11"),
    ]
    
    add_simple_table(slide, Emu(137160), Emu(1371600), Emu(11902440), Emu(5029200),
        ["UID", "Fase", "Actividad", "Duración", "Inicio", "Fin"],
        activities,
        [Emu(548640), Emu(914400), Emu(4937760), Emu(914400), Emu(1828800), Emu(1828800)],
        Pt(8))
    
    txBox = slide.shapes.add_textbox(Emu(137160), Emu(6480000), Emu(11902440), Emu(274320))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Muestra: 24 de 110 actividades nivel 4. Lista completa en archivo XML LoRA-HDL-QA_Proyecto_2026.xml. Total horas nivel 4: 1,287h."
    p.font.size = Pt(9); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = FONT_BODY
    
    # ── Slide: Atributos de la Actividad ──
    slide = next_slide()
    add_header_bar(slide, s(), "ATRIBUTOS DE LA ACTIVIDAD",
                   "Campos que describen cada actividad del cronograma | PMBOK 8va, Cronograma")
    
    attr_data = [
        ("UID", "Identificador único en el XML de ProjectLibre", "Ej: 4, 5, 68"),
        ("Nombre", "Descripción de la actividad (verbo + objeto)", "Ej: Revisar normativo y reglamento"),
        ("Nivel EDT", "OutlineLevel: 1=Proyecto, 2=Fase, 3=Entregable, 4=Paquete", "1, 2, 3, o 4"),
        ("Código WBS", "Código de la EDT al que pertenece", "Ej: 1.1.1.4"),
        ("Duración", "Horas de trabajo efectivas en formato PT...H", "Ej: PT8H0M0S (8 horas)"),
        ("Fecha Inicio", "Fecha y hora de inicio planificada (YYYY-MM-DDTHH:MM:SS)", "2026-02-02T08:00:00"),
        ("Fecha Fin", "Fecha y hora de finalización planificada", "2026-02-02T17:00:00"),
        ("Predecesoras", "Tareas que deben completarse antes (FS, SS, FF)", "Ej: 4FS (termina 4, inicia esta)"),
        ("Recurso Asignado", "Recurso humano o material asignado", "Pablo Flores (UID 1)"),
        ("Tipo de Tarea", "Fixed Work (trabajo fijo), Fixed Duration, Fixed Units", "Fixed Work"),
        ("Calendario", "Calendario aplicable (Guatemala 2026)", "UID 11: Guatemala - Feriados 2026"),
        ("Hito", "Si es hito (duración 0): TRUE/FALSE", "FALSE para tareas de trabajo"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(4937760),
        ["Atributo", "Descripción", "Valores de Ejemplo"],
        attr_data,
        [Emu(2286000), Emu(5486400), Emu(3657600)],
        Pt(10))
    
    # ── Slide: Base de las Estimaciones (Duraciones) ──
    slide = next_slide()
    add_header_bar(slide, s(), "BASE DE LAS ESTIMACIONES — DURACIONES",
                   "Fundamento de cada estimación de duración | PMBOK 8va, Cronograma")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(1828800),
        "Metodología de Estimación", [
            "Enfoque: Bottom-up estimating desde el nivel 4 de la EDT",
            "Unidad base: horas de trabajo efectivas (no días calendario)",
            "Fuentes de estimación:",
            "• Juicio de expertos: asesor asignado por FING para actividades académicas",
            "• Análogos: proyectos de graduación previos en FING (150-360h típicas)",
            "• Paramétrico: tasa de ~15h/página para redacción académica",
            "• Descomposición: tareas atómicas de 3h a 31h máximo",
            "Precisión: ±10% en paquetes nivel 4; ±15% en fases nivel 3"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(1828800),
        "Factores Considerados", [
            "Productividad: 8h/día efectivas (netas, no incluyen pausas)",
            "Calendario: lun-sáb laborable, excluye domingos y feriados",
            "Períodos de espera: modelados como tareas de duración fija (ej: resolución institucional 39h)",
            "Curva de aprendizaje: primeras tareas de cada fase tienen +20% de buffer",
            "Restricción hardware: batch=1 en M1 implica ~2-3x tiempo vs GPU",
            "Reuniones: 1h/semana con asesor incluidas en paquete de seguimiento"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(3383280), Emu(10972800), Emu(1645920),
        "Ejemplos de Estimación por Tipo de Actividad", [
            "Investigación/Lectura: 2-3h por paper (23h = ~10 papers profundos) | Redacción: 15h por sección de ~5 páginas",
            "Programación/Scripts: 15-23h por módulo funcional | Entrenamiento ML: 31h = ~3 epochs en M1 con batch=1",
            "Revisión: 15h por ronda de correcciones | Espera institucional: basado en tiempos promedio FING (3-5 días hábiles = 31-39h)"
        ], Pt(11), Pt(9))
    
    # ── Slide: Estimaciones de Duración ──
    slide = next_slide()
    add_header_bar(slide, s(), "ESTIMACIONES DE DURACIÓN",
                   "Resumen de duraciones estimadas por fase | PMBOK 8va, Cronograma")
    
    dur_data = [
        ("INICIO", "Inscripción del Tema", "78h", "Procedimiento estándar FING + 39h espera institucional"),
        ("INICIO", "Asignación de Asesor", "80h", "Búsqueda + selección + formalización + primera reunión"),
        ("INICIO", "Plan del Proyecto", "60h", "Alcance, EDT, cronograma, riesgos, comunicación"),
        ("PLANIFICACIÓN", "Anteproyecto", "152h", "Literatura + redacción + revisión + aprobación comité (31h espera)"),
        ("PLANIFICACIÓN", "Diseño Metodológico", "83h", "Metodología + arquitectura + datasets + pipeline + métricas"),
        ("EJECUCIÓN", "Marco Teórico", "161h", "4 áreas de revisión bibliográfica + redacción + correcciones"),
        ("EJECUCIÓN", "Preprocesamiento de Datos", "136h", "Catalogación + descarga + limpieza + tokenización + validación"),
        ("EJECUCIÓN", "Entrenamiento LoRA", "227h", "Configuración + piloto + 2 iteraciones + entrenamiento final + doc."),
        ("EJECUCIÓN", "Evaluación Benchmark", "106h", "Benchmarks + métricas + scripts + ejecución + comparación + doc."),
        ("EJECUCIÓN", "Análisis de Resultados", "90h", "Consolidación + estadística + visualización + comparación SOTA"),
        ("EJECUCIÓN", "Redacción del Informe", "175h", "6 capítulos IMRaD + referencias + formato + 2 revisiones asesor"),
        ("M&C", "Seguimiento", "86h", "Reuniones semanales + reportes mensuales + monitoreo cronograma"),
        ("M&C", "Control de Cambios", "68h", "Registro + evaluación + implementación de cambios + lecciones"),
        ("CIERRE", "Correcciones Finales", "75h", "Retroalimentación + correcciones + ortografía + formato + firma"),
        ("CIERRE", "Defensa Pública", "59h", "Presentación + material + ensayos + defensa + espera resolución"),
        ("CIERRE", "Entrega y Cierre", "26h", "Correcciones jurado + impresión + entrega biblioteca + cierre admin."),
    ]
    
    add_simple_table(slide, Emu(182880), Emu(1371600), Emu(11887200), Emu(5029200),
        ["Fase", "Paquete EDT Nivel 3", "Duración", "Base de la Estimación"],
        dur_data,
        [Emu(1371600), Emu(2743200), Emu(914400), Emu(6400800)],
        Pt(9))
    
    # ── Slide: Lista de Hitos ──
    slide = next_slide()
    add_header_bar(slide, s(), "LISTA DE HITOS",
                   "Hitos principales con fechas planificadas y criterios de aceptación | PMBOK 8va")
    
    hitos = [
        ("1", "Proyecto de Graduación Completado", "2026-09-14", "Aprobación formal por tribunal y constancia Biblioteca FING"),
        ("2", "INICIO DEL PROYECTO", "2026-02-02", "Carta de inscripción firmada por Director de Carrera"),
        ("3", "INICIO: INICIO", "2026-02-02", "Tema aprobado + asesor identificado formalmente"),
        ("4", "FIN: INICIO", "2026-02-26", "Plan de proyecto aprobado por asesor"),
        ("5", "INICIO: PLANIFICACION", "2026-02-26", "Inicio formal del anteproyecto con asesor asignado"),
        ("6", "FIN: PLANIFICACION", "2026-04-08", "Anteproyecto y diseño metodológico aprobados por comité"),
        ("7", "INICIO: EJECUCION", "2026-04-08", "Inicio formal de experimentos y marco teórico"),
        ("8", "FIN: EJECUCION", "2026-08-17", "Informe final completo revisado y aprobado por asesor"),
        ("9", "INICIO: MONITOREO Y CONTROL", "2026-02-26", "Primera reunión de seguimiento con asesor"),
        ("10", "FIN: MONITOREO Y CONTROL", "2026-03-13", "Cierre formal de actividades de seguimiento intensivo"),
        ("11", "INICIO: CIERRE", "2026-08-17", "Inicio de correcciones finales post-revisión asesor"),
        ("12", "FIN: CIERRE", "2026-09-14", "Documento final depositado en Biblioteca FING"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(4937760),
        ["#", "Hito", "Fecha", "Criterio de Aceptación"],
        hitos,
        [Emu(548640), Emu(3840480), Emu(1828800), Emu(5029200)],
        Pt(10))
    
    # ── Slide: Diagramas de Red del Cronograma ──
    slide = next_slide()
    add_header_bar(slide, s(), "DIAGRAMAS DE RED DEL CRONOGRAMA DEL PROYECTO",
                   "Secuencia lógica de actividades y dependencias | PMBOK 8va, Cronograma")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(10972800), Emu(914400),
        "Método de Diagramación", [
            "Método: Diagrama de Precedencia (PDM) implementado en ProjectLibre con relaciones Finish-to-Start (FS).",
            "La mayoría de las dependencias son obligatorias (duras): no se puede entrenar sin antes preprocesar datos; no se puede defender sin antes redactar."
        ], Pt(11), Pt(9))
    
    # Simplified network diagram as a sequence
    add_card(slide, Emu(457200), Emu(2468880), Emu(10972800), Emu(4114800),
        "Ruta Crítica del Proyecto (secuencia de dependencias FS)", [
            "",
            "INICIO DEL PROYECTO → Inscripción del Tema → Asignación de Asesor → Plan del Proyecto → FIN INICIO",
            "   ↓",
            "INICIO PLANIFICACIÓN → Anteproyecto → Diseño Metodológico → FIN PLANIFICACIÓN",
            "   ↓",
            "INICIO EJECUCIÓN → Marco Teórico ∥ Preprocesamiento de Datos → Entrenamiento LoRA → Evaluación Benchmark → Análisis de Resultados → Redacción del Informe → FIN EJECUCIÓN",
            "   ↓",
            "INICIO CIERRE → Correcciones Finales → Defensa Pública → Entrega y Cierre → FIN DEL PROYECTO",
            "",
            "En paralelo durante todo el proyecto: MONITOREO Y CONTROL (Seguimiento semanal + Control de Cambios bajo demanda)",
            "",
            "Dependencias clave:",
            "• UID 4→5→6→7→8→9 (Inscripción secuencial)",
            "• UID 25→26→27→28→29→30→31→32→33 (Anteproyecto secuencial)",
            "• UID 44∥46 (Marco teórico en paralelo), UID 63→68→71→73 (Entrenamiento iterativo)",
            "• UID 100→101→102→103 (Ciclos de revisión-corrección con asesor)",
        ], Pt(10), Pt(8))
    
    # ═══════════════════════════════════════════════════════════
    # DOMAIN 5: RECURSOS
    # ═══════════════════════════════════════════════════════════
    slide = next_slide()
    add_domain_slide(slide,
        "DOMINIO DE DESEMPEÑO:\nRECURSOS",
        "PMBOK 8va Edición — Identificación, adquisición y gestión de los recursos necesarios "
        "para ejecutar el proyecto. Incluye recursos humanos, equipamiento, materiales e infraestructura.",
        "👤")
    
    # ── Slide: Plan de Gestión de Recursos ──
    slide = next_slide()
    add_header_bar(slide, s(), "PLAN DE GESTIÓN DE RECURSOS",
                   "Define cómo se estimarán, adquirirán y gestionarán los recursos | PMBOK 8va")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(1645920),
        "Identificación de Recursos", [
            "Fuentes: EDT nivel 4 (cada paquete de trabajo especifica recurso asignado)",
            "Método: Bottom-up — sumar horas por recurso desde las 110 actividades nivel 4",
            "Recursos humanos: Pablo Flores (PM/Tesista) — 1,830h totales planificadas",
            "Recursos físicos: MacBook Air M1 2020 (8GB/256GB), internet 10Mbps, electricidad",
            "Recursos materiales: papel, impresión, empastado, útiles de oficina"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(1645920),
        "Gestión de Recursos", [
            "Asignación: 100% del PM al proyecto durante 225 días calendario",
            "Calendario de recurso: lun-sáb 8h/día, con flexibilidad ±2h en semanas críticas",
            "Adquisición: equipo ya adquirido; materiales se compran según necesidad",
            "Monitoreo: horas registradas vs horas planificadas (parte del EVM mensual)",
            "Optimización: si una tarea requiere menos horas, el excedente va al buffer"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(3200400), Emu(10972800), Emu(1280160),
        "Estrategia para Recursos Humanos Limitados", [
            "El proyecto es unipersonal (PM = tesista = ejecutor). No hay equipo para delegar.",
            "Estrategia: maximizar productividad con bloques de trabajo de 2-4h enfocados (técnica Pomodoro adaptada).",
            "El asesor actúa como recurso de supervisión (~40h totales estimadas, aportadas por FING sin costo directo).",
            "Plan de contingencia si el PM enfrenta indisponibilidad temporal: ajustar cronograma, notificar al asesor en 48h."
        ], Pt(11), Pt(9))
    
    # ── Slide: Carta del Equipo ──
    slide = next_slide()
    add_header_bar(slide, s(), "CARTA DEL EQUIPO",
                   "Valores, normas y acuerdos de trabajo del equipo | PMBOK 8va, Recursos")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(10972800), Emu(2743200),
        "Carta del Equipo — Proyecto LoRA-HDL-QA", [
            "Contexto: Proyecto de graduación individual. El 'equipo' está compuesto por:",
            "• Pablo Flores — PM/Tesista (ejecutor principal, 100% del trabajo técnico)",
            "• Asesor (por asignar) — Revisor técnico (~40h totales, sin costo directo para el proyecto)",
            "• Director de Carrera — Autoridad de aprobación y escalamiento",
            "",
            "Valores del Trabajo:",
            "• Excelencia académica: rigor científico en experimentación y documentación",
            "• Transparencia: reportar avances y desviaciones en tiempo real al asesor",
            "• Responsabilidad: cumplir hitos y plazos acordados; notificar riesgos anticipadamente",
            "• Reproducibilidad: todo código, datos y resultados deben ser reproducibles por terceros",
            "",
            "Normas de Trabajo:",
            "• Reuniones: semanal con asesor (1h fija), quincenal de seguimiento formal",
            "• Comunicación: email para documentos formales, WhatsApp para coordinación ágil",
            "• Documentación: minutas de reunión en 24h, reportes de avance mensuales con EVM",
            "• Código: versionado en GitHub, commits atómicos con mensajes descriptivos",
            "• Conflictos: escalar al Director de Carrera si el PM y el asesor no llegan a acuerdo",
            "",
            "Nota: Por ser proyecto unipersonal, esta carta establece normas de trabajo entre el PM y el asesor. "
            "No se requiere team charter tradicional (no hay equipo multidisciplinario)."
        ], Pt(11), Pt(9))
    
    # ── Slide: Matriz RACI ──
    slide = next_slide()
    add_header_bar(slide, s(), "MATRIZ RACI",
                   "Responsabilidades por paquete de trabajo EDT | PMBOK 8va, Recursos")
    
    raci = [
        ("Inscripción del Tema", "R", "A", "C", "I"),
        ("Asignación de Asesor", "R", "A", "C", "I"),
        ("Plan del Proyecto", "R", "A", "C", "I"),
        ("Anteproyecto", "R", "A", "C", "I"),
        ("Diseño Metodológico", "R", "A", "C", "I"),
        ("Marco Teórico", "R", "A", "C", "I"),
        ("Preprocesamiento de Datos", "R", "A", "C", "I"),
        ("Entrenamiento LoRA", "R", "A", "C", "I"),
        ("Evaluación Benchmark", "R", "A", "C", "I"),
        ("Análisis de Resultados", "R", "A", "C", "I"),
        ("Redacción Informe Final", "R", "A", "C", "I"),
        ("Seguimiento", "R", "A", "C", "I"),
        ("Control de Cambios", "R", "A", "C", "I"),
        ("Correcciones Finales", "R", "A", "C", "I"),
        ("Defensa Pública", "R", "A", "C", "I"),
        ("Entrega y Cierre", "R", "A", "C", "I"),
    ]
    
    tbl = add_simple_table(slide, Emu(274320), Emu(1371600), Emu(10058400), Emu(5114160),
        ["Paquete de Trabajo", "Pablo Flores", "Asesor", "Director", "FING/UNIS"],
        raci,
        [Emu(3474720), Emu(1554480), Emu(1554480), Emu(1554480), Emu(1554480)],
        Pt(9))
    
    # RACI legend
    txBox = slide.shapes.add_textbox(Emu(274320), Emu(6570000), Emu(10058400), Emu(274320))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "R = Responsible (ejecuta)  |  A = Accountable (aprueba/responde)  |  C = Consulted (aporta criterio)  |  I = Informed (recibe información)"
    p.font.size = Pt(9); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = FONT_BODY
    
    # ── Slide: Requisitos de Recursos ──
    slide = next_slide()
    add_header_bar(slide, s(), "REQUISITOS DE RECURSOS",
                   "Humanos, equipamiento y materiales necesarios | PMBOK 8va, Recursos")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2194560),
        "Recursos Humanos", [
            "Pablo Flores (PM/Tesista)",
            "  • 1,287 horas planificadas @ Q 80/hr = Q 102,960",
            "  • Disponibilidad 100% (tiempo completo)",
            "  • Rol: ejecución total del proyecto + gestión",
            "",
            "Asesor (por asignar por FING)",
            "  • ~40 horas estimadas (revisión técnica y reuniones)",
            "  • Honorarios cubiertos por FING/UNIS (sin costo directo)",
            "  • Rol: orientación técnica, revisión de calidad",
            "",
            "Director de Carrera / Comité / Jurado",
            "  • <10 horas totales (evaluación en hitos)",
            "  • Sin costo directo para el proyecto"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2194560),
        "Equipamiento e Infraestructura", [
            "MacBook Air M1 2020 (8GB/256GB)",
            "  • Costo adquisición: Q 6,000 (costo hundido, ya adquirido)",
            "  • ~1,000h de uso estimadas para entrenamiento y desarrollo",
            "",
            "Software (todo open-source, sin costo):",
            "  • Python, HuggingFace Transformers, PEFT, MLX",
            "  • ProjectLibre 1.9.8 (gestión de proyecto)",
            "  • Git/GitHub (control de versiones)",
            "",
            "Infraestructura:",
            "  • Internet 10Mbps (~Q 250/mes, costo absorbido)",
            "  • Electricidad (~Q 300 total, costo absorbido)",
            "  • Google Drive 100GB (EDU, gratuito)",
            "",
            "Materiales:",
            "  • Impresión 3 copias × Q 150 = Q 450",
            "  • Material defensa (posters): Q 200",
            "  • Papelería y misceláneos: Q 150"
        ], Pt(11), Pt(9))
    
    # ── Slide: Base de las Estimaciones (Recursos) ──
    slide = next_slide()
    add_header_bar(slide, s(), "BASE DE LAS ESTIMACIONES — RECURSOS",
                   "Fundamento de las estimaciones de recursos | PMBOK 8va, Recursos")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2011680),
        "Recurso Humano — PM/Tesista", [
            "Tarifa: Q 80/hora — basada en tarifa de investigador junior en Guatemala (referencia: SENACYT)",
            "Horas: 1,287h — suma bottom-up de las 110 actividades nivel 4 del cronograma XML",
            "Disponibilidad: 100% — el tesista no tiene otras obligaciones laborales durante el período",
            "Productividad: 8h/día netas — basado en experiencia previa en proyectos académicos similares",
            "Fuente de la tarifa: promedio de estipendio de investigación + costo de vida en Guatemala",
            "Nota: el costo de mano de obra es principalmente costo de oportunidad, no erogación en efectivo"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2011680),
        "Equipamiento y Materiales", [
            "MacBook M1: Q 6,000 — precio de mercado Guatemala (2025), depreciado a 2 años de uso",
            "Impresión: Q 150/copia — cotización de Copistería UNIS (zona 15, 2025)",
            "Internet: Q 250/mes — tarifa Tigo residencial 10Mbps, costo absorbido (no incremental)",
            "Electricidad: 30W × 1,000h × Q 1/kWh ≈ Q 300 — cálculo basado en consumo del M1",
            "Materiales: Q 350 — estimación empírica de proyectos de graduación previos en FING",
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(3566160), Emu(10972800), Emu(914400),
        "Reservas y Contingencias en Recursos", [
            "Contingencia de recurso humano: +10% de horas en tareas técnicas críticas (entrenamiento LoRA, evaluación benchmark) = ~70h de buffer",
            "Contingencia de equipamiento: si MacBook falla, utilizar laboratorio de cómputo FING (costo: Q 0, pero requiere coordinación de horario)",
            "Las estimaciones siguen el principio de 'estimación más probable' con un buffer separado, no inflado en cada tarea individual"
        ], Pt(11), Pt(9))
    
    # ── Slide: Estructura de Desglose de Recursos ──
    slide = next_slide()
    add_header_bar(slide, s(), "ESTRUCTURA DE DESGLOSE DE RECURSOS (RBS)",
                   "Descomposición jerárquica de recursos del proyecto | PMBOK 8va, Recursos")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(3657600), Emu(2286000),
        "1. Recursos Humanos", [
            "1.1 PM / Tesista",
            "    • Pablo Flores (100%)",
            "1.2 Asesoría Técnica",
            "    • Asesor FING (~40h)",
            "1.3 Autoridades Académicas",
            "    • Director de Carrera",
            "    • Comité Evaluador",
            "    • Jurado de Defensa",
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(4297680), Emu(1371600), Emu(3657600), Emu(2286000),
        "2. Equipamiento", [
            "2.1 Hardware",
            "    • MacBook Air M1 8GB/256GB",
            "2.2 Periféricos",
            "    • Monitor externo (opcional)",
            "    • Teclado y mouse",
            "2.3 Respaldo",
            "    • Disco externo 1TB (backup)",
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(8229600), Emu(1371600), Emu(3657600), Emu(2286000),
        "3. Software y Servicios", [
            "3.1 Desarrollo",
            "    • Python, HuggingFace, MLX",
            "3.2 Gestión",
            "    • ProjectLibre, Excel, Git",
            "3.3 Infraestructura",
            "    • Internet 10Mbps",
            "    • Google Drive 100GB",
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(3840480), Emu(5486400), Emu(2011680),
        "4. Materiales y Consumibles", [
            "4.1 Impresión y Empastado",
            "    • 3 copias informe final × Q 150",
            "4.2 Material de Defensa",
            "    • Posters, handouts, bolígrafos",
            "4.3 Papelería General",
            "    • Cuadernos, post-its, impresiones parciales",
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(3840480), Emu(5029200), Emu(2011680),
        "5. Servicios Externos", [
            "5.1 Imprenta / Copistería",
            "    • Copistería UNIS o PrintCenter",
            "5.2 Soporte Técnico FING",
            "    • Laboratorio de cómputo (contingencia)",
            "5.3 Revisión de Estilo",
            "    • Interno (tesista)",
        ], Pt(11), Pt(9))
    
    # ── Slide: Calendarios de Recursos ──
    slide = next_slide()
    add_header_bar(slide, s(), "CALENDARIOS DE RECURSOS",
                   "Disponibilidad de cada recurso a lo largo del proyecto | PMBOK 8va, Recursos")
    
    cal_res = [
        ("Pablo Flores", "100%", "Feb 02 — Sep 14, 2026", "Lun-Sáb 8h/día", "10h/día en semanas pre-defensa", "Tiempo completo"),
        ("MacBook M1 2020", "Compartido", "Feb 02 — Sep 14, 2026", "24/7 disponible", "Noches para entrenamiento largo", "~1,000h efectivas estimadas"),
        ("Asesor", "~15%", "Feb 16 — Sep 14, 2026", "1h/semanal + hitos", "Según disponibilidad FING", "~40h totales estimadas"),
        ("Director de Carrera", "< 5%", "Hitos específicos", "Según necesidad", "N/A", "< 10h totales estimadas"),
        ("Internet 10Mbps", "Siempre", "Feb — Sep 2026", "24/7", "N/A", "Costo absorbido"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(2743200),
        ["Recurso", "Disponibilidad", "Período", "Calendario Base", "Flexibilidad", "Notas"],
        cal_res,
        [Emu(1828800), Emu(1188720), Emu(2103120), Emu(2011680), Emu(2560320), Emu(2560320)],
        Pt(10))
    
    add_card(slide, Emu(457200), Emu(4297680), Emu(10972800), Emu(1097280),
        "Notas sobre Calendarios de Recursos", [
            "• El calendario del PM es el más restrictivo — determina la ruta crítica del proyecto.",
            "• MacBook M1 se usa también para actividades personales; el proyecto requiere ~1,000h de las ~3,600h disponibles en el período.",
            "• Períodos de inactividad del PM (enfermedad, emergencias) se mitigan con buffer del 10% incluido en el cronograma.",
            "• Si el MacBook requiere reparación > 3 días, se activa contingencia de usar laboratorio FING."
        ], Pt(11), Pt(9))
    
    # ═══════════════════════════════════════════════════════════
    # DOMAIN 6: COSTOS
    # ═══════════════════════════════════════════════════════════
    slide = next_slide()
    add_domain_slide(slide,
        "DOMINIO DE DESEMPEÑO:\nCOSTOS / FINANZAS",
        "PMBOK 8va Edición — Planificación, estimación, presupuestación, financiamiento "
        "y control de los costos del proyecto. Incluye EVM y análisis de valor ganado.",
        "💰")
    
    # ── Slide: Plan de Gestión Financiera ──
    slide = next_slide()
    add_header_bar(slide, s(), "PLAN DE GESTIÓN FINANCIERA",
                   "Políticas, procedimientos y documentación para gestionar costos | PMBOK 8va")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(1463040),
        "Políticas de Estimación", [
            "Método: Bottom-up desde EDT nivel 4 (horas × Q 80/hr)",
            "Moneda: Quetzales (GTQ). Tipo de cambio fijo: Q 7.80/USD",
            "Precisión: ±10% en paquetes nivel 4; ±15% en costo total",
            "Reestimación: mensual durante fase de ejecución (o ante cambios aprobados)",
            "Costo de mano de obra: costo de oportunidad, no erogación en efectivo"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(1463040),
        "Control de Costos", [
            "Monitoreo: EVM mensual (PV, EV, AC, CPI, SPI)",
            "Umbrales: Verde <5% | Amarillo 5-15% | Rojo >15%",
            "Herramienta: Excel vinculado a datos del XML de ProjectLibre",
            "Responsable: PM (Pablo Flores), con supervisión del asesor",
            "Reporte: mensual, incluido en informe de avance para asesor y Director"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(3017520), Emu(5486400), Emu(1280160),
        "Reservas", [
            "Contingencia: Q 20,846.00 — calculado como suma de Costo Esperado (CE = P × I) de los 25 riesgos identificados",
            "Gestión: Q 700.00 — para acciones de mitigación no anticipadas y gastos administrativos",
            "Uso de contingencia: solo para riesgos materializados; aprobación del asesor requerida",
            "Uso de gestión: autorizado por el PM para acciones proactivas de mitigación"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(3017520), Emu(5029200), Emu(1280160),
        "Indicadores EVM", [
            "PV (Planned Value) = Costo planificado del trabajo programado",
            "EV (Earned Value) = Costo planificado del trabajo completado",
            "AC (Actual Cost) = Costo real incurrido",
            "CPI = EV/AC (eficiencia de costo: >1 = bajo presupuesto)",
            "SPI = EV/PV (eficiencia de cronograma: >1 = adelantado)",
            "EAC = BAC/CPI (estimación al finalizar)",
            "TCPI = (BAC−EV)/(BAC−AC) (índice requerido para cumplir BAC)"
        ], Pt(11), Pt(9))
    
    # ── Slide: Estimaciones de Costos ──
    slide = next_slide()
    add_header_bar(slide, s(), "ESTIMACIONES DE COSTOS",
                   "Desglose por fase basado en datos reales del XML | PMBOK 8va, Costos")
    
    cost_data = [
        ("INICIO", "148h", "Q 11,840.00", "8.3%"),
        ("PLANIFICACIÓN", "235h", "Q 18,800.00", "13.2%"),
        ("EJECUCIÓN", "744h", "Q 59,520.00", "41.9%"),
        ("MONITOREO Y CONTROL", "86h", "Q 6,880.00", "4.8%"),
        ("CIERRE", "160h", "Q 12,800.00", "9.0%"),
        ("Subtotal Mano de Obra", "1,373h*", "Q 109,840.00", "77.3%"),
        ("", "", "", ""),
        ("MacBook M1 (costo hundido)", "—", "Q 6,000.00", "4.2%"),
        ("Materiales (impresión, posters)", "—", "Q 800.00", "0.6%"),
        ("Subtotal Directos", "—", "Q 116,640.00", "82.1%"),
        ("", "", "", ""),
        ("Reserva de Contingencia", "—", "Q 20,846.00", "14.7%"),
        ("Reserva de Gestión", "—", "Q 700.00", "0.5%"),
        ("", "", "", ""),
        ("TOTAL COSTO DEL PROYECTO", "—", "Q 138,186.00", "97.3%"),
        ("Total con MacBook (costo hundido)", "—", "Q 144,186.00", "101.5%"),
    ]
    
    add_simple_table(slide, Emu(457200), Emu(1371600), Emu(7315200), Emu(5029200),
        ["Concepto", "Horas", "Costo (Q)", "% del Total"],
        cost_data,
        [Emu(3108960), Emu(1371600), Emu(1828800), Emu(914400)],
        Pt(9))
    
    # Funding by month
    funding_data = [
        ("Febrero 2026", "211h", "Q 16,880", "Q 16,880"),
        ("Marzo 2026", "367h", "Q 29,360", "Q 46,240"),
        ("Abril 2026", "374h", "Q 29,920", "Q 76,160"),
        ("Mayo 2026", "172h", "Q 13,760", "Q 89,920"),
        ("Junio 2026", "190h", "Q 15,200", "Q 105,120"),
        ("Julio 2026", "226h", "Q 18,080", "Q 123,200"),
        ("Agosto 2026", "160h", "Q 12,800", "Q 136,000"),
        ("Septiembre 2026", "130h", "Q 10,400", "Q 146,400"),
    ]
    
    add_simple_table(slide, Emu(8229600), Emu(1371600), Emu(3200400), Emu(4572000),
        ["Período", "Horas", "Costo (Q)", "Acumulado (Q)"],
        funding_data,
        [Emu(1280160), Emu(640080), Emu(731520), Emu(914400)],
        Pt(9))
    
    # ── Slide: Base de las Estimaciones (Costos) ──
    slide = next_slide()
    add_header_bar(slide, s(), "BASE DE LAS ESTIMACIONES — COSTOS",
                   "Fundamento y metodología de cada estimación de costo | PMBOK 8va, Costos")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2011680),
        "Mano de Obra — Q 80/hora", [
            "Tarifa basada en:",
            "• Estipendio promedio de investigador junior en Guatemala (SENACYT 2025): Q 12,800/mes",
            "• Asumiendo 160h/mes productivas: Q 80/hora",
            "• Incluye costo de vida (alimentación, transporte, servicios) en Ciudad de Guatemala",
            "• Es costo de oportunidad (el tesista no recibe salario por el proyecto)",
            "Horas por fase: suma directa de duraciones del XML de ProjectLibre",
            "Validación cruzada: ~1,373h totales vs 1,287h del XML (+86h de buffer implícito en tareas de resumen)"
        ], Pt(10), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2011680),
        "Equipamiento y Materiales", [
            "MacBook M1: Q 6,000 — precio de mercado en Guatemala (nov 2025, MacBook Air M1 2020 usado)",
            "Depreciación: 2 años de uso personal previo, valor residual ~Q 3,000",
            "Impresión: Q 150/copia × 3 = Q 450 — cotización Copistería UNIS (ene 2026)",
            "Posters defensa: Q 200 — impresión a color tamaño A1 (2 posters)",
            "Papelería: Q 150 — cuadernos, post-its, bolígrafos, impresiones parciales",
        ], Pt(10), Pt(9))
    
    add_card(slide, Emu(457200), Emu(3566160), Emu(10972800), Emu(1828800),
        "Reservas de Contingencia — Cálculo", [
            "Método: Valor Monetario Esperado (EMV) = Σ (Probabilidad × Impacto Máximo) para cada uno de los 25 riesgos identificados",
            "EMV total = Q 20,846.00 (ver Registro de Riesgos en Excel LoRA_Risk_Register_Analysis.xlsx)",
            "Reserva de gestión: 0.5% del costo directo ≈ Q 700 — para gastos no anticipados ni modelados como riesgos",
            "Las reservas NO están incluidas en la línea base de costos (BAC = Q 116,640). Se añaden para el presupuesto total.",
            "Cualquier uso de la reserva de contingencia requiere documentación del riesgo materializado y aprobación del asesor."
        ], Pt(10), Pt(9))
    
    # ── Slide: Línea Base de Costos ──
    slide = next_slide()
    add_header_bar(slide, s(), "LÍNEA BASE DE COSTOS",
                   "Presupuesto aprobado contra el cual se mide el desempeño de costos | PMBOK 8va")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(10972800), Emu(914400),
        "Definición de Línea Base de Costos", [
            "BAC (Budget at Completion) = Q 116,640.00 — incluye mano de obra directa (Q 109,840) + equipamiento y materiales (Q 6,800). "
            "No incluye reservas de contingencia ni gestión. La línea base se congela tras la aprobación del anteproyecto (Mar 2026)."
        ], Pt(11), Pt(9))
    
    # S-curve data
    scurve = [
        ("Feb 2026", "Q 16,880", "Q 16,880", "14.5%"),
        ("Mar 2026", "Q 29,360", "Q 46,240", "39.6%"),
        ("Abr 2026", "Q 29,920", "Q 76,160", "65.3%"),
        ("May 2026", "Q 13,760", "Q 89,920", "77.1%"),
        ("Jun 2026", "Q 15,200", "Q 105,120", "90.1%"),
        ("Jul 2026", "Q 18,080", "Q 123,200", "105.6%"),
        ("Ago 2026", "Q 12,800", "Q 136,000", "116.6%"),
        ("Sep 2026", "Q 10,400", "Q 146,400", "125.5%"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(2468880), Emu(7315200), Emu(3657600),
        ["Período", "PV Mensual (Q)", "PV Acumulado (Q)", "% BAC"],
        scurve,
        [Emu(1645920), Emu(2011680), Emu(2011680), Emu(1645920)],
        Pt(10))
    
    add_card(slide, Emu(7772400), Emu(2468880), Emu(4114800), Emu(3657600),
        "Indicadores de Línea Base", [
            "BAC = Q 116,640.00",
            "",
            "Presupuesto total (incluye reservas):",
            "Q 116,640 + Q 20,846 + Q 700",
            "= Q 138,186.00",
            "",
            "Curva S: perfil de gasto front-loaded",
            "con pico en Mar-Abr 2026 (planif. + inicio ejec.)",
            "",
            "Desembolso: el costo de mano de obra es",
            "costo de oportunidad. Solo equipamiento y",
            "materiales requieren desembolso en efectivo:",
            "Q 6,800 (ya erogado el MacBook)."
        ], Pt(10), Pt(8))
    
    # ── Slide: Requisitos de Financiamiento del Proyecto ──
    slide = next_slide()
    add_header_bar(slide, s(), "REQUISITOS DE FINANCIAMIENTO DEL PROYECTO",
                   "Flujo de fondos necesario por período | PMBOK 8va, Costos")
    
    funding = [
        ("Febrero 2026", "211h", "Q 16,880", "Q 6,000 (MacBook)", "Q 22,880", "Q 22,880"),
        ("Marzo 2026", "367h", "Q 29,360", "Q 0", "Q 29,360", "Q 52,240"),
        ("Abril 2026", "374h", "Q 29,920", "Q 0", "Q 29,920", "Q 82,160"),
        ("Mayo 2026", "172h", "Q 13,760", "Q 0", "Q 13,760", "Q 95,920"),
        ("Junio 2026", "190h", "Q 15,200", "Q 0", "Q 15,200", "Q 111,120"),
        ("Julio 2026", "226h", "Q 18,080", "Q 0", "Q 18,080", "Q 129,200"),
        ("Agosto 2026", "160h", "Q 12,800", "Q 0", "Q 12,800", "Q 142,000"),
        ("Septiembre 2026", "130h", "Q 10,400", "Q 800 (materiales)", "Q 11,200", "Q 153,200"),
        ("", "", "", "", "", ""),
        ("Contingencia", "—", "—", "—", "Q 20,846", "Q 174,046"),
        ("Gestión", "—", "—", "—", "Q 700", "Q 174,746"),
    ]
    
    add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(4572000),
        ["Período", "Horas", "Mano de Obra (Q)", "Erogaciones (Q)", "Total Período (Q)", "Acumulado (Q)"],
        funding,
        [Emu(1645920), Emu(914400), Emu(1828800), Emu(2011680), Emu(2011680), Emu(2011680)],
        Pt(9))
    
    add_card(slide, Emu(457200), Emu(6032880), Emu(10972800), Emu(731520),
        "Notas sobre Financiamiento", [
            "Fuente de fondos: recursos propios del tesista. No hay financiamiento externo (beca, patrocinio, o subvención).",
            "La mano de obra es costo de oportunidad (no erogación). Las erogaciones reales son solo Q 6,800 (MacBook + materiales).",
            "El flujo de fondos es front-loaded: 65% del presupuesto se concentra en Feb-Abr 2026 (fases de inicio, planificación y arranque de ejecución)."
        ], Pt(10), Pt(9))
    
    # ── Slide: Estrategia de Financiamiento ──
    slide = next_slide()
    add_header_bar(slide, s(), "ESTRATEGIA DE FINANCIAMIENTO",
                   "Cómo se obtendrán y administrarán los fondos del proyecto | PMBOK 8va, Costos")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2743200),
        "Estrategia de Financiamiento", [
            "Fuente primaria: Recursos propios del tesista (Pablo Flores)",
            "",
            "El proyecto no requiere financiamiento externo ya que:",
            "• El hardware principal (MacBook M1) ya fue adquirido previamente",
            "• El software es 100% open-source y gratuito",
            "• La mano de obra es aportada por el tesista como parte de su formación académica",
            "• Los costos de materiales son menores a Q 1,000",
            "",
            "Estructura de costos:",
            "• Costos hundidos (ya incurridos): MacBook M1 Q 6,000",
            "• Costos futuros de materiales: Q 800 (impresión, empastado, posters)",
            "• Costos recurrentes absorbidos: internet Q 250/mes, electricidad ~Q 300 total",
            "• Costo de oportunidad (mano de obra): Q 109,840",
            "",
            "Total erogación en efectivo requerida: Q 800 (materiales para fase de cierre)",
            "",
            "No se requiere plan de financiamiento con entidades bancarias, inversionistas, o patrocinadores."
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2743200),
        "Administración de Fondos", [
            "Responsable: PM (Pablo Flores)",
            "",
            "Mecanismo:",
            "• Cuenta bancaria personal existente (no se requiere cuenta dedicada)",
            "• Pago directo a proveedores (copistería) contra entrega",
            "",
            "Control:",
            "• Registro de gastos en Excel vinculado al EVM",
            "• Facturas/recibos escaneados y archivados en Google Drive",
            "",
            "Flujo de fondos:",
            "• Feb 2026: Q 0 (MacBook ya pagado previamente)",
            "• Sep 2026: Q 800 (impresión + posters)",
            "• Reservas: no requieren depósito previo (solo se usan si se materializan riesgos)",
            "",
            "Cierre financiero:",
            "• Conciliación de gastos reales vs presupuesto en Sep 2026",
            "• Documentación de lecciones aprendidas en gestión financiera",
            "• Informe final de costos para el expediente del proyecto"
        ], Pt(11), Pt(9))
    
    # ═══════════════════════════════════════════════════════════
    # DOMAIN 7: RIESGOS
    # ═══════════════════════════════════════════════════════════
    slide = next_slide()
    add_domain_slide(slide,
        "DOMINIO DE DESEMPEÑO:\nRIESGOS",
        "PMBOK 8va Edición — Identificación, análisis cualitativo y cuantitativo, "
        "planificación de respuestas, implementación y monitoreo de riesgos del proyecto.",
        "⚠️")
    
    # ── Slide: Plan de Gestión de Riesgos ──
    slide = next_slide()
    add_header_bar(slide, s(), "PLAN DE GESTIÓN DE RIESGOS",
                   "Define cómo se gestionarán los riesgos del proyecto | PMBOK 8va, Riesgos")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(1463040),
        "Metodología de Gestión de Riesgos", [
            "Enfoque: Cualitativo-cuantitativo iterativo (PMBOK 8va)",
            "Frecuencia de identificación: quincenal (en reuniones con asesor)",
            "Análisis: Probabilidad × Impacto monetario = Costo Esperado (CE)",
            "Umbral de respuesta: CE > Q 1,500 requiere plan de respuesta documentado",
            "Herramientas: Excel (LoRA_Risk_Register_Analysis.xlsx), matriz Prob-Impacto",
            "Categorías de riesgo: Técnico, Organizacional, Recursos Humanos, Calidad, Operativo"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(1463040),
        "Estrategias de Respuesta", [
            "EVITAR: Cambiar enfoque/metodología si el riesgo es inaceptable (ej: cambiar modelo base si Qwen no carga)",
            "MITIGAR: Reducir probabilidad (checkpointing, backups, revisiones anticipadas) y/o impacto (plan B)",
            "TRANSFERIR: No aplica en proyecto académico individual (no hay seguros ni terceros que absorban riesgo)",
            "ACEPTAR: Riesgos con CE < Q 500 — monitorear sin plan de respuesta activo",
            "ESCALAR: Al Director de Carrera si el riesgo impacta normativa FING o plazos institucionales"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(457200), Emu(3017520), Emu(5486400), Emu(1280160),
        "Roles en Gestión de Riesgos", [
            "Identificación: PM + Asesor (en reuniones quincenales)",
            "Análisis cualitativo: PM (matriz Prob-Impacto)",
            "Análisis cuantitativo (EMV): PM (Excel), validado por asesor",
            "Planificación de respuestas: PM propone, asesor aprueba",
            "Implementación: PM ejecuta las respuestas",
            "Monitoreo: Revisión quincenal del registro de riesgos"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(3017520), Emu(5029200), Emu(1280160),
        "Indicadores Clave de Riesgo (KRI)", [
            "KRI-1: Espera respuesta asesor > 10 días → escalar",
            "KRI-2: Loss NaN en entrenamiento → detener y revisar datos",
            "KRI-3: Temperatura MacBook > 85°C → pausar y refrigerar",
            "KRI-4: Desviación cronograma > 15% → plan de recuperación",
            "KRI-5: Datasets con > 5% de datos corruptos → reprocesar",
        ], Pt(11), Pt(9))
    
    # ── Slide: Registro de Riesgos ──
    slide = next_slide()
    add_header_bar(slide, s(), "REGISTRO DE RIESGOS",
                   "25 riesgos identificados | Contingencia total: Q 20,846.00 | PMBOK 8va, Riesgos")
    
    # Top 15 risks
    risks = [
        ("R-10", "Modelo base Qwen3-4B no carga en 8GB RAM (OOM)", "Técnico", "40%", "Q 4,000", "Q 1,600", "Quantización 4-bit GGUF; probar modelos más pequeños"),
        ("R-05", "Arquitectura experimental no viable con hardware M1", "Técnico", "30%", "Q 5,000", "Q 1,500", "QLoRA 4-bit con MLX; gradient checkpointing"),
        ("R-13", "Entrenamiento final interrumpido por falla hardware", "Técnico", "25%", "Q 6,000", "Q 1,500", "Checkpointing cada 100 steps; respaldo en nube"),
        ("R-14", "Resultados del modelo final no superan baseline", "Técnico", "30%", "Q 5,000", "Q 1,500", "Ajustar rank LoRA (r=16→64); data augmentation"),
        ("R-12", "Tiempo de entrenamiento excede estimado (throttling)", "Técnico", "45%", "Q 3,000", "Q 1,350", "Entrenar horarios nocturnos; ventilador externo"),
        ("R-20", "Observaciones extensas del asesor requieren reescritura", "Organizacional", "25%", "Q 5,000", "Q 1,250", "Redacción iterativa temprana; anticipar correcciones"),
        ("R-01", "Demora en resolución institucional", "Organizacional", "40%", "Q 3,120", "Q 1,248", "Escalar con Director; recordatorios cada 5 días"),
        ("R-04", "Rechazo del anteproyecto por comité evaluador", "Organizacional", "20%", "Q 5,000", "Q 1,000", "Reducir alcance; revisión previa con Director"),
        ("R-06", "No encontrar datasets HDL/QA de calidad suficiente", "Técnico", "30%", "Q 3,000", "Q 900", "Generar dataset sintético con GPT-4; web scraping"),
        ("R-08", "Datos de entrenamiento con corrupción no detectada", "Calidad", "30%", "Q 3,000", "Q 900", "Pipeline de validación automatizada; revisión manual"),
        ("R-02", "Retraso en asignación de asesor", "Recursos Humanos", "35%", "Q 2,520", "Q 882", "Contactar múltiples candidatos simultáneamente"),
        ("R-11", "Entrenamiento piloto diverge o produce NaN", "Técnico", "35%", "Q 2,500", "Q 875", "Reducir learning rate; aumentar warmup steps"),
        ("R-17", "Modelos SOTA publicados usan GPU/A100 no reproducible", "Técnico", "40%", "Q 2,000", "Q 800", "Documentar limitaciones; normalizar por recursos"),
        ("R-15", "VerilogEval v2 no ejecuta en entorno ARM/M1", "Técnico", "30%", "Q 2,500", "Q 750", "Compilar dependencias nativas ARM; usar Rosetta 2"),
        ("R-18", "Metodología requiere cambios mayores tras revisión asesor", "Organizacional", "25%", "Q 3,000", "Q 750", "Reuniones quincenales; enviar borradores por secciones"),
    ]
    
    add_simple_table(slide, Emu(137160), Emu(1371600), Emu(11902440), Emu(5029200),
        ["ID", "Descripción del Riesgo", "Cat.", "Prob.", "Impacto (Q)", "CE (Q)", "Respuesta Planificada"],
        risks,
        [Emu(548640), Emu(3017520), Emu(731520), Emu(457200), Emu(914400), Emu(640080), Emu(3474720)],
        Pt(7))
    
    txBox = slide.shapes.add_textbox(Emu(137160), Emu(6480000), Emu(11902440), Emu(274320))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Muestra: 15 riesgos de 25. Registro completo en Excel LoRA_Risk_Register_Analysis.xlsx. CE Total = Q 20,846.00."
    p.font.size = Pt(9); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = FONT_BODY
    
    # ── Slide: Reporte de Riesgos ──
    slide = next_slide()
    add_header_bar(slide, s(), "REPORTE DE RIESGOS",
                   "Resumen ejecutivo de la situación de riesgos | PMBOK 8va, Riesgos")
    
    add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(1645920),
        "Resumen General de Riesgos", [
            "Total riesgos identificados: 25",
            "Riesgos con plan de respuesta activo (CE > Q 1,500): 12",
            "Riesgos en monitoreo pasivo (CE < Q 500): 5",
            "Riesgos escalables al Director: 3 (R-01, R-02, R-04)",
            "",
            "Costo Esperado Total (EMV): Q 20,846.00",
            "Reserva de contingencia recomendada: Q 20,846.00",
            "Reserva de gestión: Q 700.00",
            "",
            "Perfil de riesgo del proyecto: MODERADO",
            "• La mayoría de riesgos son técnicos (15/25 = 60%)",
            "• Los riesgos organizacionales tienen alto impacto pero baja probabilidad",
            "• No hay riesgos catastróficos (impacto > Q 10,000) identificados"
        ], Pt(11), Pt(9))
    
    add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(1645920),
        "Distribución por Categoría", [
            "Técnico: 15 riesgos | CE: Q 13,994 (67.1%)",
            "  → Principales: R-10 (OOM), R-05 (hardware), R-13 (falla HW)",
            "",
            "Organizacional: 7 riesgos | CE: Q 4,950 (23.7%)",
            "  → Principales: R-20 (reescritura), R-01 (demora inst.)",
            "",
            "Calidad: 1 riesgo | CE: Q 900 (4.3%)",
            "  → R-08: Corrupción de datos de entrenamiento",
            "",
            "Recursos Humanos: 1 riesgo | CE: Q 882 (4.2%)",
            "  → R-02: Retraso en asignación de asesor",
            "",
            "Operativo: 1 riesgo | CE: Q 120 (0.6%)",
            "  → R-25: Error en impresión/empastado"
        ], Pt(10), Pt(8))
    
    add_card(slide, Emu(457200), Emu(3200400), Emu(10972800), Emu(1645920),
        "Top 5 Riesgos por Costo Esperado", [
            "1. R-10: Modelo OOM en 8GB RAM (CE: Q 1,600) → Mitigar: quantización 4-bit GGUF",
            "2. R-05: Arquitectura no viable en M1 (CE: Q 1,500) → Mitigar: QLoRA 4-bit con MLX",
            "3. R-13: Falla hardware durante entrenamiento (CE: Q 1,500) → Mitigar: checkpointing + backup",
            "4. R-14: Modelo no supera baseline (CE: Q 1,500) → Mitigar: ajuste hiperparámetros, data augmentation",
            "5. R-12: Thermal throttling M1 (CE: Q 1,350) → Mitigar: entrenamiento nocturno + ventilación externa",
        ], Pt(11), Pt(9))
    
    # ── Slide: Registro de Lecciones Aprendidas ──
    slide = next_slide()
    add_header_bar(slide, s(), "REGISTRO DE LECCIONES APRENDIDAS",
                   "Conocimiento para mejorar desempeño actual y futuro | PMBOK 8va, Riesgos y General")
    
    lessons = [
        ("LL-001", "Inicio", "Validar calendario antes de planificar", "Crash ProjectLibre por <Exceptions> sin <RecurrenceType>", "Probar XML en ProjectLibre antes de entregar; verificar compatibilidad de herramientas"),
        ("LL-002", "Planificación", "Anticipar asignación de asesor", "Demora en asignación afectó revisiones tempranas", "Contactar candidatos desde la inscripción del tema; tener opciones de respaldo"),
        ("LL-003", "Ejecución", "Validar hardware antes de diseñar arquitectura", "8GB RAM obligó quantización agresiva (4-bit)", "Hacer PoC de memoria antes de diseñar pipeline completo; probar con modelo pequeño primero"),
        ("LL-004", "Ejecución", "Versionar datos y checkpoints desde el inicio", "Corrupción de datos requirió reprocesar 31h", "Usar DVC/Git LFS; checkpointing automático cada 100 steps; backup en nube"),
        ("LL-005", "M&C", "Documentar minutas de todas las reuniones", "Observaciones verbales del asesor se olvidaron", "Enviar minuta en 24h post-reunión; confirmar entendimiento con el asesor por email"),
        ("LL-006", "Cierre", "Reservar margen para correcciones del jurado", "Espera por nota y correcciones puede extenderse más allá de lo planificado", "Solicitar defensa con 2 semanas de anticipación; tener semana buffer post-defensa"),
    ]
    
    add_simple_table(slide, Emu(182880), Emu(1371600), Emu(11887200), Emu(3657600),
        ["ID", "Fase", "Lección Aprendida", "Causa Raíz", "Recomendación para Futuro"],
        lessons,
        [Emu(640080), Emu(914400), Emu(2560320), Emu(3017520), Emu(3657600)],
        Pt(9))
    
    add_card(slide, Emu(457200), Emu(5200000), Emu(10972800), Emu(914400),
        "Proceso de Lecciones Aprendidas", [
            "Recolección: Al finalizar cada fase del proyecto (5 momentos: Inicio, Planificación, Ejecución, M&C, Cierre)",
            "Análisis: Técnica de los 5 Porqués (5 Whys) para identificar causa raíz de cada desviación o problema",
            "Documentación: Formato estandarizado en el Registro de Lecciones Aprendidas (parte de este Plan de Dirección)",
            "Difusión: Incluir en repositorio GitHub del proyecto; compartir con futuros tesistas de FING"
        ], Pt(11), Pt(9))
    
    # ═══════════════════════════════════════════════════════════
    # CIERRE
    # ═══════════════════════════════════════════════════════════
    slide = next_slide()
    # Navy bg
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    # Accents
    for y_pos in [Emu(0), Emu(6766560)]:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), y_pos, SLIDE_W, Emu(91440))
        accent.fill.solid(); accent.fill.fore_color.rgb = BLUE; accent.line.fill.background()
    
    add_cover_text(Emu(914400), Emu(2011680), Emu(10363200), Emu(1097280),
                   "PLAN DE DIRECCIÓN DEL PROYECTO", Pt(36), True, WHITE)
    add_cover_text(Emu(914400), Emu(3200400), Emu(10363200), Emu(457200),
                   "LoRA-HDL-QA — Proyecto de Graduación FING/UNIS", Pt(18), False, BLUE_LIGHT)
    add_cover_text(Emu(914400), Emu(3749040), Emu(10363200), Emu(365760),
                   f"Total: {s()} diapositivas | PMBOK 8va Edición | 7 Dominios de Desempeño", Pt(14), False, BLUE_LIGHT)
    add_cover_text(Emu(914400), Emu(4297680), Emu(10363200), Emu(457200),
                   "Pablo Rodolfo Alexander Flores Mollinedo", Pt(16), True, WHITE)
    add_cover_text(Emu(914400), Emu(4754880), Emu(10363200), Emu(365760),
                   "Universidad del Istmo de Guatemala — FING | Febrero 2026", Pt(14), False, BLUE_LIGHT)
    
    # ── Save ──
    output_path = '/Users/pabloflores/Documents/estructura/revision/entregables_finales/LoRA_Plan_Direccion_Proyecto_Completo.pptx'
    prs.save(output_path)
    print(f"Presentation saved: {output_path}")
    print(f"Total slides: {s()}")

if __name__ == '__main__':
    build_presentation()
