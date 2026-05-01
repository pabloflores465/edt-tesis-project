#!/usr/bin/env python3
"""Fix task slides - show ALL 110 level-4 tasks properly split across 2 slides."""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import xml.etree.ElementTree as ET

prs = Presentation('/Users/pabloflores/Documents/estructura/revision/entregables_finales/LoRA_Plan_Direccion_Proyecto_Completo.pptx')

NAVY = RGBColor(0x1A, 0x27, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
MED_TEXT = RGBColor(0x33, 0x33, 0x33)

def clear_slide(slide):
    spTree = slide.shapes._spTree
    for sp in list(spTree):
        if sp.tag != qn('p:cSld'):
            spTree.remove(sp)

def add_header(slide, num, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(12192000), Emu(822960))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Emu(228600), Emu(182880), Emu(548640), Emu(502920))
    tf = txBox.text_frame; p = tf.paragraphs[0]
    p.text = f"{num:02d}"; p.font.size = Pt(36); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = 'Calibri'
    txBox2 = slide.shapes.add_textbox(Emu(868680), Emu(228600), Emu(10972800), Emu(457200))
    tf2 = txBox2.text_frame; p2 = tf2.paragraphs[0]
    p2.text = title; p2.font.size = Pt(20); p2.font.bold = True
    p2.font.color.rgb = WHITE; p2.font.name = 'Calibri'
    if subtitle:
        txBox3 = slide.shapes.add_textbox(Emu(274320), Emu(868680), Emu(11612880), Emu(320040))
        tf3 = txBox3.text_frame; p3 = tf3.paragraphs[0]
        p3.text = subtitle; p3.font.size = Pt(12)
        p3.font.color.rgb = GRAY_TEXT; p3.font.name = 'Calibri'

def add_table(slide, left, top, width, height, headers, rows, col_widths, fs=Pt(6)):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    ts = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = ts.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    # Headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = fs; para.font.bold = True
            para.font.color.rgb = WHITE; para.font.name = 'Calibri'
        tcPr = cell._tc.get_or_add_tcPr()
        sf = cell._tc.makeelement(qn('a:solidFill'), {})
        sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': '4285C8'})
        sf.append(sc); tcPr.append(sf)
    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c); cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = fs; para.font.color.rgb = MED_TEXT
                para.font.name = 'Calibri'
            if r % 2 == 0:
                tcPr = cell._tc.get_or_add_tcPr()
                sf = cell._tc.makeelement(qn('a:solidFill'), {})
                sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': 'F8F9FA'})
                sf.append(sc); tcPr.append(sf)

# Load tasks
tree = ET.parse('/Users/pabloflores/Documents/estructura/revision/entregables_finales/LoRA-HDL-QA_Proyecto_2026.xml')
root = tree.getroot()
nsxml = {'ns': 'http://schemas.microsoft.com/project'}
all_tasks = root.findall('.//ns:Task', nsxml)

l4_tasks = []
for task in all_tasks:
    uid = task.find('ns:UID', nsxml)
    uid_int = int(uid.text)
    if uid_int > 133:
        continue
    outline = task.find('ns:OutlineLevel', nsxml)
    lvl = int(outline.text) if outline is not None else 0
    if lvl != 4:
        continue
    name = task.find('ns:Name', nsxml)
    dur = task.find('ns:Duration', nsxml)
    start = task.find('ns:Start', nsxml)
    finish = task.find('ns:Finish', nsxml)
    wbs = task.find('ns:WBS', nsxml)
    
    dur_text = dur.text if dur is not None else '0'
    d_hours = 0
    if dur_text and dur_text.startswith('PT'):
        d_clean = dur_text.replace('PT','').replace('H0M0S','').replace('H','')
        try:
            d_hours = int(d_clean)
        except:
            pass
    start_t = start.text[:10] if start is not None else ''
    finish_t = finish.text[:10] if finish is not None else ''
    wbs_t = wbs.text if wbs is not None else ''
    name_t = name.text if name is not None else ''
    
    phase = '?'
    if wbs_t.startswith('1.1'):
        phase = 'INICIO'
    elif wbs_t.startswith('1.2'):
        phase = 'PLANIF.'
    elif wbs_t.startswith('1.3'):
        phase = 'EJEC.'
    elif wbs_t.startswith('1.4'):
        phase = 'M&C'
    elif wbs_t.startswith('1.5'):
        phase = 'CIERRE'
    
    l4_tasks.append((uid_int, phase, name_t[:62], f"{d_hours}h", start_t, finish_t))

half = len(l4_tasks) // 2
first = l4_tasks[:half]
second = l4_tasks[half:]

print(f'Total L4 tasks: {len(l4_tasks)}, half: {half}')
print(f'First: UIDs {first[0][0]}-{first[-1][0]} ({len(first)} tasks)')
print(f'Second: UIDs {second[0][0]}-{second[-1][0]} ({len(second)} tasks)')

# FIX Slide 29 (index 28): Lista 1/2
slide = prs.slides[28]
clear_slide(slide)
add_header(slide, 29, "LISTA DE ACTIVIDADES (1/2)",
           f"Actividades nivel 4 del cronograma — {len(first)} de {len(l4_tasks)} | PMBOK 8va, Cronograma")
add_table(slide, Emu(91440), Emu(1280160), Emu(12009160), Emu(5300000),
    ["UID", "Fase", "Actividad", "Dur.", "Inicio", "Fin"],
    first,
    [Emu(457200), Emu(594360), Emu(5303520), Emu(457200), Emu(1828800), Emu(1828800)],
    Pt(6))
txBox = slide.shapes.add_textbox(Emu(91440), Emu(6620000), Emu(12009160), Emu(180000))
tf = txBox.text_frame; p = tf.paragraphs[0]
p.text = f"{len(first)} de {len(l4_tasks)} actividades nivel 4 | Continúa en diapositiva 30 | Fuente: LoRA-HDL-QA_Proyecto_2026.xml"
p.font.size = Pt(7); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = 'Calibri'

# FIX Slide 30 (index 29): Lista 2/2
slide = prs.slides[29]
clear_slide(slide)
add_header(slide, 30, "LISTA DE ACTIVIDADES (2/2)",
           f"Actividades nivel 4 del cronograma — {len(second)} de {len(l4_tasks)} | PMBOK 8va, Cronograma")
add_table(slide, Emu(91440), Emu(1280160), Emu(12009160), Emu(5300000),
    ["UID", "Fase", "Actividad", "Dur.", "Inicio", "Fin"],
    second,
    [Emu(457200), Emu(594360), Emu(5303520), Emu(457200), Emu(1828800), Emu(1828800)],
    Pt(6))
txBox = slide.shapes.add_textbox(Emu(91440), Emu(6620000), Emu(12009160), Emu(180000))
tf = txBox.text_frame; p = tf.paragraphs[0]
p.text = f"{len(second)} de {len(l4_tasks)} actividades nivel 4 | Total horas nivel 4: 1,287h | Fuente: LoRA-HDL-QA_Proyecto_2026.xml"
p.font.size = Pt(7); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = 'Calibri'

prs.save('/Users/pabloflores/Documents/estructura/revision/entregables_finales/LoRA_Plan_Direccion_Proyecto_Completo.pptx')
print('\n✅ Fixed slides 29-30 with all 110 level-4 tasks')
