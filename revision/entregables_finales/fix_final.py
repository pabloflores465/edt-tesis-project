#!/usr/bin/env python3
"""Final fixes: EDT overlap, Gracias slide, per-task risks, per-task financial, overflow fixes."""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import openpyxl, xml.etree.ElementTree as ET

BASE = '/Users/pabloflores/Documents/estructura/revision/entregables_finales'
prs = Presentation(f'{BASE}/LoRA_Plan_Direccion_Proyecto_Completo.pptx')

NAVY = RGBColor(0x1A, 0x27, 0x44)
BLUE = RGBColor(0x42, 0x85, 0xF4)
BLUE_LIGHT = RGBColor(0xC8, 0xDC, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
MED_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
ACCENT_BLUE = RGBColor(0x42, 0x85, 0xC8)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN = Emu(457200)
CONTENT_W = SLIDE_W - 2 * MARGIN

def clear_slide(slide):
    spTree = slide.shapes._spTree
    for sp in list(spTree):
        if sp.tag != qn('p:cSld'):
            spTree.remove(sp)

def add_header(slide, num, title, sub=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Emu(822960))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Emu(228600), Emu(182880), Emu(548640), Emu(502920))
    tf = txBox.text_frame; p = tf.paragraphs[0]
    p.text = f"{num:02d}"; p.font.size = Pt(36); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = 'Calibri'
    txBox2 = slide.shapes.add_textbox(Emu(868680), Emu(228600), Emu(10972800), Emu(457200))
    tf2 = txBox2.text_frame; p2 = tf2.paragraphs[0]
    p2.text = title; p2.font.size = Pt(20); p2.font.bold = True
    p2.font.color.rgb = WHITE; p2.font.name = 'Calibri'
    if sub:
        txBox3 = slide.shapes.add_textbox(MARGIN, Emu(868680), CONTENT_W, Emu(320040))
        tf3 = txBox3.text_frame; tf3.word_wrap = True
        p3 = tf3.paragraphs[0]; p3.text = sub
        p3.font.size = Pt(12); p3.font.color.rgb = GRAY_TEXT; p3.font.name = 'Calibri'

def add_card(slide, left, top, width, height, title, lines, ts=Pt(11), bs=Pt(9)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6); card.line.width = Pt(0.5)
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(68580); tf.margin_right = Emu(68580)
    tf.margin_top = Emu(45720); tf.margin_bottom = Emu(45720)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = ts; p.font.bold = True; p.font.color.rgb = NAVY; p.font.name = 'Calibri'; p.space_after = Pt(3)
    for line in lines:
        p2 = tf.add_paragraph(); p2.text = line
        p2.font.size = bs; p2.font.color.rgb = MED_TEXT; p2.font.name = 'Calibri'; p2.space_after = Pt(1)
    return card

def add_table(slide, left, top, width, height, headers, rows, col_widths, fs=Pt(8)):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    ts = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = ts.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = fs; para.font.bold = True
            para.font.color.rgb = WHITE; para.font.name = 'Calibri'
        tcPr = cell._tc.get_or_add_tcPr()
        sf = cell._tc.makeelement(qn('a:solidFill'), {})
        sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': '4285C8'})
        sf.append(sc); tcPr.append(sf)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c); cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = fs; para.font.color.rgb = MED_TEXT; para.font.name = 'Calibri'
            if r % 2 == 0:
                tcPr = cell._tc.get_or_add_tcPr()
                sf = cell._tc.makeelement(qn('a:solidFill'), {})
                sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': 'F8F9FA'})
                sf.append(sc); tcPr.append(sf)
    return ts

def add_edt_node(slide, left, top, width, height, text, fill_color=ACCENT_BLUE, text_color=WHITE, fs=Pt(8)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = fill_color
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(22860); tf.margin_right = Emu(22860)
    tf.margin_top = Emu(11430); tf.margin_bottom = Emu(11430)
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = fs; p.font.bold = True
        p.font.color.rgb = text_color; p.font.name = 'Calibri'
        p.alignment = 1
    return box

def add_connector(slide, x1, y1, x2, y2, color=RGBColor(0x99, 0x99, 0x99)):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR_TYPE.STRAIGHT = 1
    conn.line.color.rgb = color; conn.line.width = Pt(1)

# ═══════════════════════════════════════════════════
# 1. REBUILD EDT (slide 15) — fix overlaps, add numbering
# ═══════════════════════════════════════════════════
slide = prs.slides[14]
clear_slide(slide)
add_header(slide, 15, "EDT / WBS — ESTRUCTURA DE DESGLOSE DEL TRABAJO",
           "Descomposición jerárquica del alcance total en 16 paquetes de trabajo | PMBOK 8va, Alcance")

# Label
txBox = slide.shapes.add_textbox(Emu(274320), Emu(1188720), Emu(1828800), Emu(274320))
tf = txBox.text_frame; p = tf.paragraphs[0]
p.text = "Diagrama EDT"; p.font.size = Pt(12); p.font.bold = True
p.font.color.rgb = NAVY; p.font.name = 'Calibri'

# Root node
root_x = Emu(4572000)
add_edt_node(slide, root_x, Emu(1554480), Emu(2743200), Emu(411480),
             "Proyecto de Graduación\nLoRA-HDL-QA", NAVY, WHITE, Pt(9))

# Phase level nodes (level 2) — wider to fit text
phases = [
    ("1.1 INICIO\n148h", Emu(183000)),
    ("1.2 PLANIFICACIÓN\n235h", Emu(2378000)),
    ("1.3 EJECUCIÓN\n744h", Emu(4634000)),
    ("1.4 MONITOREO\nY CONTROL\n86h", Emu(6858000)),
    ("1.5 CIERRE\n160h", Emu(9052560)),
]
phase_y = Emu(2286000)
phase_w = Emu(2100000)
phase_h = Emu(530000)

for name, x in phases:
    add_edt_node(slide, x, phase_y, phase_w, phase_h, name, ACCENT_BLUE, WHITE, Pt(8))

# Connectors root → phases
root_btm = Emu(1965960)
root_cx = Emu(5943600)
for _, px in phases:
    pcx = px + phase_w // 2
    add_connector(slide, root_cx, root_btm, pcx, phase_y)

# Work package nodes — repositioned to avoid overlaps
wp_y = Emu(3018000)
wp_y2 = Emu(3650000)  # Second row for EJECUCIÓN row 2 and M&C

# INICIO children (1.1.x)
inic_w = Emu(550000)
inic_h = Emu(457200)
inic_xs = [Emu(280000), Emu(900000), Emu(1520000)]
inic_names = ["1.1.1 Inscripción\ndel Tema", "1.1.2 Asignación\nde Asesor", "1.1.3 Plan del\nProyecto"]
for name, ix in zip(inic_names, inic_xs):
    add_edt_node(slide, ix, wp_y, inic_w, inic_h, name, BLUE_LIGHT, NAVY, Pt(6))

# PLANIF children (1.2.x)
plan_xs = [Emu(2470000), Emu(3300000)]
plan_names = ["1.2.1 Anteproyecto", "1.2.2 Diseño\nMetodológico"]
for name, px in zip(plan_names, plan_xs):
    add_edt_node(slide, px, wp_y, Emu(730000), inic_h, name, BLUE_LIGHT, NAVY, Pt(6))

# EJECUCIÓN row 1 (1.3.1-1.3.3)
ej_w = Emu(650000)
ej_xs1 = [Emu(4700000), Emu(5430000), Emu(6160000)]
ej_names1 = ["1.3.1 Marco\nTeórico", "1.3.2 Preproc.\nde Datos", "1.3.3 Entrenam.\nLoRA"]
for name, ex in zip(ej_names1, ej_xs1):
    add_edt_node(slide, ex, wp_y, ej_w, inic_h, name, BLUE_LIGHT, NAVY, Pt(6))

# EJECUCIÓN row 2 (1.3.4-1.3.6)
ej_names2 = ["1.3.4 Evaluac.\nBenchmark", "1.3.5 Análisis\nResultados", "1.3.6 Redacción\nInforme Final"]
for name, ex in zip(ej_names2, ej_xs1):
    add_edt_node(slide, ex, wp_y2, ej_w, inic_h, name, BLUE_LIGHT, NAVY, Pt(6))

# M&C children (moved to second row to avoid overlap with EJECUCIÓN)
mc_xs = [Emu(6920000), Emu(7720000)]
mc_names = ["1.4.1 Seguimiento", "1.4.2 Control\nde Cambios"]
for name, mx in zip(mc_names, mc_xs):
    add_edt_node(slide, mx, wp_y2, Emu(700000), inic_h, name, BLUE_LIGHT, NAVY, Pt(6))

# CIERRE children
cie_xs = [Emu(9120000), Emu(9920000), Emu(10720000)]
cie_names = ["1.5.1 Correcc.\nFinales", "1.5.2 Defensa\nPública", "1.5.3 Entrega\ny Cierre"]
for name, cx in zip(cie_names, cie_xs):
    add_edt_node(slide, cx, wp_y, Emu(700000), inic_h, name, BLUE_LIGHT, NAVY, Pt(6))

# Phase → child connectors
for name, px in phases:
    pcx = px + phase_w // 2
    py_btm = phase_y + phase_h
    if "INICIO" in name:
        for ix in inic_xs: add_connector(slide, pcx, py_btm, ix + inic_w//2, wp_y)
    elif "PLANIF" in name:
        for ppx in plan_xs: add_connector(slide, pcx, py_btm, ppx + Emu(730000)//2, wp_y)
    elif "EJECUCIÓN" in name:
        for ex in ej_xs1:
            add_connector(slide, pcx, py_btm, ex + ej_w//2, wp_y)
            add_connector(slide, pcx, py_btm, ex + ej_w//2, wp_y2)
    elif "MONITOREO" in name:
        for mx in mc_xs: add_connector(slide, pcx, py_btm, mx + Emu(700000)//2, wp_y2)
    elif "CIERRE" in name:
        for cx in cie_xs: add_connector(slide, pcx, py_btm, cx + Emu(700000)//2, wp_y)

print('✓ EDT rebuilt — no overlaps, numbered phases 1.1-1.5, children numbered')

# ═══════════════════════════════════════════════════
# 2. ADD "GRACIAS" SLIDE (new slide 56)
# ═══════════════════════════════════════════════════
# Add new slide at end
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
for y in [Emu(0), Emu(6766560)]:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), y, SLIDE_W, Emu(91440))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

def cover_text(slide, l, t, w, h, text, size, bold=False, color=WHITE, align=1):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = size
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = 'Calibri'; p.alignment = align

cover_text(slide, Emu(914400), Emu(2011680), Emu(10363200), Emu(1097280), "GRACIAS", Pt(44), True)
cover_text(slide, Emu(914400), Emu(3200400), Emu(10363200), Emu(457200),
           "Plan de Dirección del Proyecto — LoRA-HDL-QA", Pt(18), False, BLUE_LIGHT)
cover_text(slide, Emu(914400), Emu(3749040), Emu(10363200), Emu(365760),
           "56 diapositivas | PMBOK 8va Edición | 7 Dominios de Desempeño", Pt(14), False, BLUE_LIGHT)
cover_text(slide, Emu(914400), Emu(4297680), Emu(10363200), Emu(457200),
           "Pablo Rodolfo Alexander Flores Mollinedo", Pt(16), True)
cover_text(slide, Emu(914400), Emu(4754880), Emu(10363200), Emu(365760),
           "Universidad del Istmo de Guatemala — FING | Abril 2026", Pt(14), False, BLUE_LIGHT)
print('✓ Gracias slide added (slide 56)')

# ═══════════════════════════════════════════════════
# 3. REBUILD REGISTRO DE RIESGOS — per-task analysis
# ═══════════════════════════════════════════════════
wb = openpyxl.load_workbook(f'{BASE}/LoRA_Risk_Register_Analysis.xlsx', data_only=True)
ws = wb['Registro de Riesgos']

risk_rows = []
for row in ws.iter_rows(min_row=5, max_row=30, values_only=True):
    if row[0] is not None:
        task_uid = int(row[0])
        risk_id = row[2]
        desc = row[3][:65]
        prob = row[5]
        impact = row[6]
        ce = row[9]
        resp = row[7][:55] if row[7] else ''
        risk_rows.append((task_uid, risk_id, desc, f"{prob*100:.0f}%", f"Q {impact:,.0f}", f"Q {ce:,.0f}", resp))

# Slide 52: rebuild as per-task risk analysis
slide = prs.slides[51]
clear_slide(slide)
add_header(slide, 52, "REGISTRO DE RIESGOS — ANÁLISIS POR TAREA",
           f"{len(risk_rows)} riesgos mapeados a tareas del cronograma | Contingencia total: Q 20,846 | PMBOK 8va")
add_table(slide, Emu(91440), Emu(1280160), Emu(12009160), Emu(5300000),
    ["Tarea\nUID", "Riesgo", "Descripción del Riesgo", "Prob.", "Impacto", "CE (Q)", "Respuesta"],
    risk_rows,
    [Emu(457200), Emu(457200), Emu(2926080), Emu(457200), Emu(731520), Emu(640080), Emu(3291840)],
    Pt(5))
print('✓ Risk register rebuilt — per-task analysis')

# ═══════════════════════════════════════════════════
# 4. REBUILD ESTIMACIONES DE COSTOS — per-task
# ═══════════════════════════════════════════════════
wb2 = openpyxl.load_workbook(f'{BASE}/LoRA_Cost_Estimates_Funding.xlsx', data_only=True)
ws2 = wb2['Estimación de Costos']

cost_rows = []
for row in ws2.iter_rows(min_row=5, max_row=ws2.max_row, values_only=True):
    if row[0] is not None and row[1] is not None:
        phase = str(row[0])
        uid = int(row[1]) if row[1] else None
        name = str(row[2])[:55] if row[2] else ''
        hours = row[3] if row[3] else 0
        cost = row[5] if row[5] else 0
        cost_rows.append((uid, phase, name, f"{hours}h", f"Q {cost:,.0f}"))

# Slide 45: rebuild with per-task costs
slide = prs.slides[44]
clear_slide(slide)
add_header(slide, 45, "ESTIMACIONES DE COSTOS — DESGLOSE POR TAREA",
           f"Costos directos por actividad | Tarifa: Q 80/hr | {len(cost_rows)} tareas | PMBOK 8va, Costos")
add_table(slide, Emu(457200), Emu(1280160), Emu(11277600), Emu(5300000),
    ["UID", "Fase", "Tarea", "Horas", "Costo (Q)"],
    cost_rows[:45],  # First 45 tasks
    [Emu(457200), Emu(914400), Emu(5029200), Emu(640080), Emu(914400)],
    Pt(7))

txBox = slide.shapes.add_textbox(MARGIN, Emu(6620000), CONTENT_W, Emu(180000))
tf = txBox.text_frame; p = tf.paragraphs[0]
p.text = f"45 de {len(cost_rows)} tareas | Costo total mano de obra: Q 102,960 | Continúa en diapositiva 46"
p.font.size = Pt(7); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = 'Calibri'

# Slide 46: remaining costs (overwrite "Base de las Estimaciones — Costos")
slide = prs.slides[45]
clear_slide(slide)
add_header(slide, 46, "ESTIMACIONES DE COSTOS (2/2) — BASE DE LAS ESTIMACIONES",
           f"Continuación + fundamento de estimaciones | PMBOK 8va, Costos")
add_table(slide, Emu(457200), Emu(1280160), Emu(11277600), Emu(3657600),
    ["UID", "Fase", "Tarea", "Horas", "Costo (Q)"],
    cost_rows[45:85],  # Next 40 tasks
    [Emu(457200), Emu(914400), Emu(5029200), Emu(640080), Emu(914400)],
    Pt(7))

add_card(slide, Emu(457200), Emu(5110000), Emu(10972800), Emu(1463040),
    'Base de las Estimaciones — Costos', [
        'Tarifa mano de obra Q 80/hr: basada en estipendio investigador junior Guatemala (ref. SENACYT 2025, ~Q 12,800/mes ÷ 160h).',
        'Costo de oportunidad: el tesista no recibe salario. | MacBook M1: Q 6,000 precio mercado GT (2025, depreciado 2 años).',
        'Materiales Q 800: cotización Copistería UNIS (Q 150/copia × 3) + posters Q 200 + papelería Q 150.',
        'Precisión: ±10% en paquetes nivel 4. Contingencia calculada como EMV de 25 riesgos (ver Registro de Riesgos).'
    ], Pt(10), Pt(8))

print('✓ Cost estimates rebuilt — per-task breakdown (2 slides)')

# ═══════════════════════════════════════════════════
# 5. FIX TABLE OVERFLOW — d9, d25
# ═══════════════════════════════════════════════════

# d9: Matriz de Comunicaciones — rebuild with tight columns
slide = prs.slides[8]
clear_slide(slide)
add_header(slide, 9, "MATRIZ DE COMUNICACIONES",
           "Detalle de flujos de información entre interesados | PMBOK 8va, Interesados")
comms = [
    ("Reunión semanal", "Presencial/Zoom", "Semanal", "Minuta de reunión", "PM, Asesor", "PM"),
    ("Avance mensual", "Email + PDF", "Mensual", "Informe de avance (EVM)", "PM, Asesor, Director", "PM"),
    ("Alerta de riesgo", "WhatsApp/Email", "Inmediato", "Notificación de riesgo", "PM, Asesor", "Cualquiera"),
    ("Solicitud de cambio", "Formato + Email", "Según necesidad", "RFC (Request for Change)", "PM, Asesor", "PM"),
    ("Convocatoria defensa", "Email formal", "Una vez (Sep 2026)", "Carta de convocatoria", "PM, Director, Jurado", "PM"),
    ("Publicación repo", "GitHub / HuggingFace", "Post-defensa", "README + paper + pesos", "Comunidad global", "PM"),
    ("Consulta normativa", "Email", "Según necesidad", "Solicitud formal", "PM, Coordinador Tesis", "PM"),
]
add_table(slide, MARGIN, Emu(1371600), CONTENT_W, Emu(4572000),
    ["Tipo", "Medio/Canal", "Frecuencia", "Formato/Entregable", "Audiencia", "Responsable"],
    comms,
    [Emu(1645920), Emu(1645920), Emu(1645920), Emu(2560320), Emu(2560320), Emu(1188720)], Pt(9))
print('✓ d9: Matriz de Comunicaciones rebuilt')

# d25: Línea Base del Cronograma — already rebuilt, check card overflow
# The "Notas sobre la Línea Base" card is at bottom. Let me fix it.
slide = prs.slides[24]
# Remove old card
for shape in list(slide.shapes):
    if shape.has_text_frame and 'Notas sobre la Línea Base' in shape.text_frame.text:
        slide.shapes._spTree.remove(shape._element)
# Add properly sized card
add_card(slide, Emu(457200), Emu(4297680), Emu(10972800), Emu(2286000),
    'Notas sobre la Línea Base', [
        '• La línea base del cronograma es la versión aprobada que se utiliza para comparar el desempeño real.',
        '• Fecha de congelamiento: posterior a la aprobación del anteproyecto por el comité evaluador (~Mar 25, 2026).',
        '• La línea base solo se modifica mediante solicitudes de cambio aprobadas formalmente.',
        '• M&C tiene actividades distribuidas a lo largo del proyecto; el bloque principal está en Feb-Mar 2026.',
        '• Duración total del proyecto: 225 días calendario (Feb 02 → Sep 14, 2026). Horas efectivas: 1,287h.',
        '• Se utiliza calendario Guatemala 2026 (lun-sáb 8h/día). Domingos y feriados no laborables.',
        '• Las estimaciones de duración provienen del archivo XML LoRA-HDL-QA_Proyecto_2026.xml.'
    ], Pt(10), Pt(8))
print('✓ d25: Línea Base card adjusted')

# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
out = f'{BASE}/LoRA_Plan_Direccion_Proyecto_Completo.pptx'
prs.save(out)
# Also save as v2
import shutil
shutil.copy(out, f'{BASE}/LoRA_Plan_Direccion_Proyecto_Completo_v2.pptx')
print(f'\n✅ Saved: {out} ({len(prs.slides)} slides)')
