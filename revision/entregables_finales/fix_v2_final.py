#!/usr/bin/env python3
"""Apply ALL fixes to v2 and split tasks into 5 slides."""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import xml.etree.ElementTree as ET

BASE = '/Users/pabloflores/Documents/estructura/revision/entregables_finales'
prs = Presentation(f'{BASE}/LoRA_Plan_Direccion_Proyecto_Completo_v2.pptx')

NAVY = RGBColor(0x1A, 0x27, 0x44)
BLUE = RGBColor(0x42, 0x85, 0xF4)
BLUE_LIGHT = RGBColor(0xC8, 0xDC, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
MED_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
ACCENT_BLUE = RGBColor(0x42, 0x85, 0xC8)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN = Emu(457200)
CONTENT_W = SLIDE_W - 2 * MARGIN  # 11277600

def clear_slide(slide):
    spTree = slide.shapes._spTree
    for sp in list(spTree):
        if sp.tag != qn('p:cSld'):
            spTree.remove(sp)

def add_header(slide, num, title, sub=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Emu(822960))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Emu(228600), Emu(182880), Emu(548640), Emu(502920))
    tf = txBox.text_frame; p = tf.paragraphs[0]
    p.text = f"{num:02d}"; p.font.size = Pt(36); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = 'Calibri'
    txBox2 = slide.shapes.add_textbox(Emu(868680), Emu(228600), Emu(10972800), Emu(457200))
    tf2 = txBox2.text_frame; p2 = tf2.paragraphs[0]
    p2.text = title; p2.font.size = Pt(20); p2.font.bold = True
    p2.font.color.rgb = WHITE; p2.font.name = 'Calibri'
    if sub:
        txBox3 = slide.shapes.add_textbox(MARGIN, Emu(868680), CONTENT_W, Emu(320040))
        tf3 = txBox3.text_frame; tf3.word_wrap = True
        p3 = tf3.paragraphs[0]; p3.text = sub
        p3.font.size = Pt(11); p3.font.color.rgb = GRAY_TEXT; p3.font.name = 'Calibri'

def add_card(slide, left, top, width, height, title, lines, ts=Pt(11), bs=Pt(9)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6); card.line.width = Pt(0.5)
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(68580); tf.margin_right = Emu(68580)
    tf.margin_top = Emu(45720); tf.margin_bottom = Emu(45720)
    p = tf.paragraphs[0]; p.text = title
    p.font.size = ts; p.font.bold = True; p.font.color.rgb = NAVY; p.font.name = 'Calibri'; p.space_after = Pt(3)
    for line in lines:
        p2 = tf.add_paragraph(); p2.text = line
        p2.font.size = bs; p2.font.color.rgb = MED_TEXT; p2.font.name = 'Calibri'; p2.space_after = Pt(1)
    return card

def add_table(slide, left, top, width, height, headers, rows, col_widths, fs=Pt(8)):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    ts = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = ts.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = fs; para.font.bold = True
            para.font.color.rgb = WHITE; para.font.name = 'Calibri'
        tcPr = cell._tc.get_or_add_tcPr()
        sf = cell._tc.makeelement(qn('a:solidFill'), {})
        sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': '4285C8'})
        sf.append(sc); tcPr.append(sf)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c); cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = fs; para.font.color.rgb = MED_TEXT; para.font.name = 'Calibri'
            if r % 2 == 0:
                tcPr = cell._tc.get_or_add_tcPr()
                sf = cell._tc.makeelement(qn('a:solidFill'), {})
                sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': 'F8F9FA'})
                sf.append(sc); tcPr.append(sf)
    return ts

def add_domain_slide(slide, domain_name, description):
    clear_slide(slide)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(2743200), SLIDE_W, Emu(91440))
    line.fill.solid(); line.fill.fore_color.rgb = BLUE; line.line.fill.background()
    txBox = slide.shapes.add_textbox(Emu(914400), Emu(1371600), Emu(10363200), Emu(914400))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = domain_name; p.font.size = Pt(40)
    p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Calibri'; p.alignment = 1
    txBox2 = slide.shapes.add_textbox(Emu(914400), Emu(3200400), Emu(10363200), Emu(914400))
    tf2 = txBox2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = description; p2.font.size = Pt(14)
    p2.font.color.rgb = BLUE_LIGHT; p2.font.name = 'Calibri'; p2.alignment = 1
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(6766560), SLIDE_W, Emu(91440))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

# ═══════════════════════════════════════════════════
# 1. FIX DOMAIN SLIDES (2,5,10,23,35,43,50)
# ═══════════════════════════════════════════════════
domains = {
    2: ("DOMINIO DE DESEMPEÑO:\nGOBERNANZA",
        "Un modelo de gestión alineado con los objetivos institucionales. Define estructuras de autoridad, toma de decisiones, y mecanismos de control del proyecto."),
    5: ("DOMINIO DE DESEMPEÑO:\nINTERESADOS",
        "Identificación, análisis, involucramiento y comunicación efectiva con todas las partes interesadas del proyecto."),
    10: ("DOMINIO DE DESEMPEÑO:\nALCANCE",
        "Definición, validación y control del alcance del proyecto y del producto. Incluye requisitos, EDT, y trazabilidad."),
    23: ("DOMINIO DE DESEMPEÑO:\nCRONOGRAMA",
        "Planificación, desarrollo, monitoreo y control del cronograma. Incluye definición de actividades, secuenciación, estimación de duraciones y línea base."),
    35: ("DOMINIO DE DESEMPEÑO:\nRECURSOS",
        "Identificación, adquisición y gestión de los recursos necesarios para ejecutar el proyecto."),
    43: ("DOMINIO DE DESEMPEÑO:\nCOSTOS / FINANZAS",
        "Planificación, estimación, presupuestación, financiamiento y control de los costos del proyecto. Incluye EVM y análisis de valor ganado."),
    50: ("DOMINIO DE DESEMPEÑO:\nRIESGOS",
        "Identificación, análisis cualitativo y cuantitativo, planificación de respuestas, implementación y monitoreo de riesgos del proyecto."),
}
for sn, (title, desc) in domains.items():
    add_domain_slide(prs.slides[sn - 1], title, desc)
print('✓ Domain slides fixed')

# ═══════════════════════════════════════════════════
# 2. FIX TABLES — remove ID columns
# ═══════════════════════════════════════════════════

# d4: Registro de Supuestos
slide = prs.slides[3]; clear_slide(slide)
add_header(slide, 4, "REGISTRO DE SUPUESTOS", "Supuestos documentados durante la iniciación | PMBOK 8va, Gobernanza")
assumptions = [
    ("Disponibilidad de asesor", "Se contará con un asesor asignado antes del 16 de febrero 2026", "Alta", "Director de Carrera"),
    ("Hardware adecuado", "MacBook Air M1 8GB es suficiente para fine-tuning LoRA con quantización 4-bit", "Media", "PM/Tesista"),
    ("Datasets públicos", "Existen datasets HDL/QA públicos de calidad suficiente para fine-tuning", "Media", "PM/Tesista"),
    ("Normativa estable", "El reglamento de graduación FING no cambiará durante el proyecto", "Alta", "Director de Carrera"),
    ("Disponibilidad del PM", "El tesista mantendrá disponibilidad 100% durante el período del proyecto", "Alta", "PM/Tesista"),
    ("Acceso a software", "Python, HuggingFace, MLX, y librerías asociadas estarán disponibles sin costo", "Alta", "PM/Tesista"),
    ("Plazos institucionales", "Los tiempos de respuesta institucional no excederán 10 días hábiles", "Media", "FING/UNIS"),
    ("Modelo base accesible", "Qwen3-4B será accesible en HuggingFace y cargable en 8GB RAM", "Media", "PM/Tesista"),
    ("VerilogEval funcional", "El benchmark VerilogEval v2 será ejecutable en entorno ARM/M1", "Media", "PM/Tesista"),
    ("Sin cambios mayores", "No habrá cambios en el alcance del proyecto una vez aprobado el anteproyecto", "Alta", "PM + Asesor"),
]
add_table(slide, MARGIN, Emu(1371600), CONTENT_W, Emu(4937760),
    ["Supuesto", "Descripción", "Certidumbre", "Responsable de Validación"],
    assumptions, [Emu(2560320), Emu(5303520), Emu(914400), Emu(2286000)], Pt(10))

# d6: Registro de Interesados
slide = prs.slides[5]; clear_slide(slide)
add_header(slide, 6, "REGISTRO DE INTERESADOS", "Identificar, clasificar y documentar a todos los interesados | PMBOK 8va, Interesados")
stakeholders = [
    ("Pablo Flores", "PM / Tesista", "Interno", "Alto", "Alto", "Gestionar de cerca"),
    ("Asesor", "Orientador técnico", "Interno", "Alto", "Alto", "Gestionar de cerca"),
    ("Director de Carrera FING", "Autoridad académica", "Interno", "Alto", "Medio", "Mantener satisfecho"),
    ("Comité Evaluador", "Aprobación anteproyecto", "Interno", "Alto", "Bajo", "Mantener satisfecho"),
    ("Jurado de Defensa", "Evaluación final", "Interno", "Alto", "Bajo", "Mantener satisfecho"),
    ("Coordinador de Tesis", "Procesos administrativos", "Interno", "Medio", "Medio", "Gestionar de cerca"),
    ("Biblioteca FING/UNIS", "Recepción y archivo", "Interno", "Bajo", "Bajo", "Monitorear"),
    ("Estudiantes FING", "Audiencia académica", "Externo", "Bajo", "Bajo", "Monitorear"),
    ("Comunidad ML/HW", "Investigadores en HDL", "Externo", "Bajo", "Bajo", "Monitorear"),
    ("HuggingFace / GitHub", "Plataformas de publicación", "Externo", "Medio", "Bajo", "Mantener informado"),
    ("UNIS — Rectoría", "Aval institucional", "Externo", "Bajo", "Bajo", "Monitorear"),
    ("Familia Flores", "Apoyo logístico y moral", "Externo", "Bajo", "Medio", "Mantener informado"),
]
add_table(slide, MARGIN, Emu(1371600), CONTENT_W, Emu(4937760),
    ["Nombre", "Rol en el Proyecto", "Clasificación", "Poder", "Interés", "Estrategia"],
    stakeholders, [Emu(2194560), Emu(2377440), Emu(1280160), Emu(914400), Emu(914400), Emu(2743200)], Pt(9))

# d13: Documentación de Requisitos
slide = prs.slides[12]; clear_slide(slide)
add_header(slide, 13, "DOCUMENTACIÓN DE REQUISITOS", "Registro y clasificación de requerimientos | PMBOK 8va, Alcance")
reqs = [
    ("2 configuraciones fine-tuning LoRA", "Comparar chunking vs filtrado", "A", "NE-02", "E-04, E-05", "2 archivos .safetensors"),
    ("Evaluar con VerilogEval v2 (156 prob.)", "Benchmark estandarizado", "A", "NE-02", "E-06", "Script + JSON 156 entradas"),
    ("Cuantificar corrupción en datos", "Reproducir hallazgo del paper", "A", "NE-02", "E-04, E-06", "verify_paper_numbers.py ejecutable"),
    ("Scripts verificación reproducción", "Reproducibilidad científica", "M", "NE-02", "E-04", "Scripts sin errores en entorno limpio"),
    ("Hardware real (M1 8GB)", "Condición experimental real", "A", "NE-02", "E-05", "Logs con ID hardware"),
    ("Informe formato FING + IMRaD", "Requisito de grado", "A", "NE-03", "E-07, E-10", "Aprobación asesor + tribunal"),
    ("Repositorio estructura clara", "Divulgación científica", "M", "NE-02", "E-04", "README instrucciones reproducción"),
    ("Cronograma máximo 360h", "Restricción institucional", "A", "NE-01", "E-08", "Registro horas ≤ 360"),
    ("Inscripción + asesor (Feb.)", "Procedimiento FING", "A", "NE-01", "E-01", "Constancia + carta asesor"),
    ("Anteproyecto aprobado (Mar.)", "Procedimiento FING", "A", "NE-01", "E-02", "Visto bueno Director + Asesor"),
    ("Avances ejecución (Abr-Ago)", "Procedimiento FING", "A", "NE-01", "E-03, E-07", "Documentos entregados y revisados"),
    ("Defensa + entrega final (Sep.)", "Procedimiento FING", "A", "NE-01", "E-09, E-10", "Acta defensa + constancia"),
]
add_table(slide, MARGIN, Emu(1371600), CONTENT_W, Emu(5029200),
    ["Descripción", "Propósito", "Prioridad", "Origen", "Entregable", "Criterio de Aceptación"],
    reqs, [Emu(2560320), Emu(1828800), Emu(640080), Emu(914400), Emu(1645920), Emu(3474720)], Pt(9))

# d52: Registro de Riesgos
slide = prs.slides[51]; clear_slide(slide)
add_header(slide, 52, "REGISTRO DE RIESGOS", "25 riesgos identificados | Contingencia total: Q 20,846.00 | PMBOK 8va, Riesgos")
risks = [
    ("Modelo base no carga en 8GB RAM (OOM)", "Técnico", "40%", "Q 4,000", "Q 1,600", "Quantización 4-bit GGUF; probar modelos más pequeños"),
    ("Arquitectura experimental no viable con hardware M1", "Técnico", "30%", "Q 5,000", "Q 1,500", "QLoRA 4-bit con MLX; gradient checkpointing"),
    ("Entrenamiento final interrumpido por falla hardware", "Técnico", "25%", "Q 6,000", "Q 1,500", "Checkpointing cada 100 steps; respaldo en nube"),
    ("Resultados del modelo final no superan baseline", "Técnico", "30%", "Q 5,000", "Q 1,500", "Ajustar rank LoRA (r=16→64); data augmentation"),
    ("Tiempo de entrenamiento excede estimado (throttling)", "Técnico", "45%", "Q 3,000", "Q 1,350", "Entrenar horarios nocturnos; ventilador externo"),
    ("Observaciones extensas del asesor requieren reescritura", "Organizacional", "25%", "Q 5,000", "Q 1,250", "Redacción iterativa temprana; anticipar correcciones"),
    ("Demora en resolución institucional", "Organizacional", "40%", "Q 3,120", "Q 1,248", "Escalar con Director; recordatorios cada 5 días"),
    ("Rechazo del anteproyecto por comité evaluador", "Organizacional", "20%", "Q 5,000", "Q 1,000", "Reducir alcance; revisión previa con Director"),
    ("No encontrar datasets HDL/QA de calidad suficiente", "Técnico", "30%", "Q 3,000", "Q 900", "Generar dataset sintético con GPT-4; web scraping"),
    ("Datos de entrenamiento con corrupción no detectada", "Calidad", "30%", "Q 3,000", "Q 900", "Pipeline de validación automatizada; revisión manual"),
    ("Retraso en asignación de asesor", "Recursos Humanos", "35%", "Q 2,520", "Q 882", "Contactar múltiples candidatos simultáneamente"),
    ("Entrenamiento piloto diverge o produce NaN", "Técnico", "35%", "Q 2,500", "Q 875", "Reducir learning rate; aumentar warmup steps"),
    ("Modelos SOTA publicados usan GPU/A100 no reproducible", "Técnico", "40%", "Q 2,000", "Q 800", "Documentar limitaciones; normalizar por recursos"),
    ("VerilogEval v2 no ejecuta en entorno ARM/M1", "Técnico", "30%", "Q 2,500", "Q 750", "Compilar dependencias nativas ARM; usar Rosetta 2"),
    ("Metodología requiere cambios mayores tras revisión asesor", "Organizacional", "25%", "Q 3,000", "Q 750", "Reuniones quincenales; enviar borradores por secciones"),
]
add_table(slide, Emu(274320), Emu(1371600), Emu(11643360), Emu(5029200),
    ["Descripción del Riesgo", "Cat.", "Prob.", "Impacto (Q)", "CE (Q)", "Respuesta Planificada"],
    risks, [Emu(3291840), Emu(822960), Emu(548640), Emu(914400), Emu(640080), Emu(3657600)], Pt(7))

# d54: Registro de Lecciones Aprendidas
slide = prs.slides[53]; clear_slide(slide)
add_header(slide, 54, "REGISTRO DE LECCIONES APRENDIDAS", "Conocimiento para mejorar desempeño actual y futuro | PMBOK 8va, Riesgos y General")
lessons = [
    ("Inicio", "Validar calendario antes de planificar", "Crash ProjectLibre por <Exceptions> sin <RecurrenceType>", "Probar XML en ProjectLibre antes de entregar; verificar compatibilidad"),
    ("Planificación", "Anticipar asignación de asesor", "Demora en asignación afectó revisiones tempranas", "Contactar candidatos desde la inscripción; tener opciones de respaldo"),
    ("Ejecución", "Validar hardware antes de diseñar arquitectura", "8GB RAM obligó quantización agresiva (4-bit)", "Hacer PoC de memoria antes de diseñar pipeline; probar modelo pequeño primero"),
    ("Ejecución", "Versionar datos y checkpoints desde inicio", "Corrupción de datos requirió reprocesar 31h", "Usar DVC/Git LFS; checkpointing automático cada 100 steps"),
    ("M&C", "Documentar minutas de todas las reuniones", "Observaciones verbales del asesor se olvidaron", "Enviar minuta en 24h post-reunión; confirmar entendimiento por email"),
    ("Cierre", "Reservar margen para correcciones del jurado", "Espera por nota y correcciones puede extenderse", "Solicitar defensa con 2 semanas de anticipación; tener semana buffer"),
]
add_table(slide, MARGIN, Emu(1371600), CONTENT_W, Emu(4937760),
    ["Fase", "Lección Aprendida", "Causa Raíz", "Recomendación para Futuro"],
    lessons, [Emu(1188720), Emu(2743200), Emu(3291840), Emu(4206240)], Pt(10))
print('✓ Tables fixed (IDs removed)')

# ═══════════════════════════════════════════════════
# 3. EDT: remove summary card, add WBS numbering
# ═══════════════════════════════════════════════════
slide = prs.slides[14]
for shape in list(slide.shapes):
    if shape.has_text_frame:
        t = shape.text_frame.text
        if 'Resumen EDT' in t or ('Nivel 1:' in t and 'proyecto' in t.lower()):
            slide.shapes._spTree.remove(shape._element)
# Add WBS numbering below phases
wbs_nums = [
    ("1.1", Emu(1280160)), ("1.2", Emu(3474720)), ("1.3", Emu(5669280)),
    ("1.4", Emu(7863840)), ("1.5", Emu(10058400)),
]
for code, x in wbs_nums:
    txBox = slide.shapes.add_textbox(x, Emu(2743200), Emu(2011680), Emu(228600))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = code; p.font.size = Pt(9)
    p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Calibri'; p.alignment = 1
print('✓ EDT: summary removed, WBS numbering added')

# ═══════════════════════════════════════════════════
# 4. d26: remove cronograma summary
# ═══════════════════════════════════════════════════
slide = prs.slides[25]
for shape in list(slide.shapes):
    if shape.has_text_frame and 'Resumen de Cronograma' in shape.text_frame.text:
        slide.shapes._spTree.remove(shape._element)
print('✓ Cronograma summary removed')

# ═══════════════════════════════════════════════════
# 5. d28: rebuild calendar slide
# ═══════════════════════════════════════════════════
slide = prs.slides[27]; clear_slide(slide)
add_header(slide, 28, "CALENDARIOS DEL PROYECTO", "Calendario guatemalteco completo 2026 con feriados oficiales | PMBOK 8va")
add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2286000),
    "Calendario Base — Guatemala 2026", [
        "Días laborables: Lunes a Sábado — Horario: 08:00-12:00 y 13:00-17:00 (8h/día efectivas)",
        "Domingo: No laborable (descanso obligatorio por ley)",
        "FERIADOS OFICIALES 2026 (no laborables):",
        "• 1 Ene — Año Nuevo | • 2-4 Abr — Semana Santa (Jueves, Viernes, Sábado Santo)",
        "• 1 May — Día del Trabajo | • 30 Jun — Día del Ejército",
        "• 15 Sep — Día de la Independencia | • 20 Oct — Día de la Revolución",
        "• 1 Nov — Día de Todos los Santos | • 24 Dic — Nochebuena (½ día)",
        "• 25 Dic — Navidad | • 31 Dic — Fin de Año (½ día)",
        "Feriados que afectan al proyecto (Feb—Sep 2026):",
        "Semana Santa (2-4 Abr), Día del Trabajo (1 May), Día del Ejército (30 Jun), Independencia (15 Sep)",
    ], Pt(10), Pt(8))
add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2286000),
    "Recurso — Pablo Flores", [
        "Disponibilidad: 100% (tiempo completo) | Máximo: 8h/día laborable | Flexibilidad: hasta 10h/día en semanas críticas",
        "Período activo: Feb 02 — Sep 14, 2026 (225 días) | Total horas planificadas: 1,287h efectivas | Promedio: ~5.7h/día",
        "",
        "Nota técnica (ProjectLibre 1.9.8): Bug con <Exceptions> sin <RecurrenceType> causa crash.",
        "Feriados documentados aquí pero NO en el XML por compatibilidad. Ajustar fechas manualmente en seguimiento.",
    ], Pt(10), Pt(8))
add_card(slide, Emu(457200), Emu(3880000), Emu(10972800), Emu(2560000),
    "Períodos de Espera Institucionales (modelados como tareas en el cronograma)", [
        "• Resolución institucional: 39h (Feb 6 → Feb 13) | • Aprobación comité: 31h (Mar 19 → Mar 25)",
        "• Retroalimentación 1ra revisión: 23h (Jul 29 → Ago 3) | • Retroalimentación 2da revisión: 39h (Ago 6 → Ago 12)",
        "• Resolución jurado: 23h (Sep 4 → Sep 9)",
        "Total esperas: ~155h (12% del total). Estas esperas son tiempo calendario, no trabajo activo del tesista.",
        "Impacto: Las esperas institucionales están en la ruta crítica. Cualquier extensión afecta la fecha de finalización.",
        "Plan de contingencia: Si la espera supera 2 semanas sobre lo planificado, escalar al Director de Carrera."
    ], Pt(10), Pt(8))
print('✓ Calendar rebuilt')

# ═══════════════════════════════════════════════════
# 6. TASKS: 5 slides (110 tasks total)
# ═══════════════════════════════════════════════════
tree = ET.parse(f'{BASE}/LoRA-HDL-QA_Proyecto_2026.xml')
root = tree.getroot()
nsxml = {'ns': 'http://schemas.microsoft.com/project'}
all_tasks = root.findall('.//ns:Task', nsxml)

l4_tasks = []
for task in all_tasks:
    uid = task.find('ns:UID', nsxml)
    uid_int = int(uid.text)
    if uid_int > 133: continue
    outline = task.find('ns:OutlineLevel', nsxml)
    if outline is None or int(outline.text) != 4: continue
    name = task.find('ns:Name', nsxml)
    dur = task.find('ns:Duration', nsxml)
    start = task.find('ns:Start', nsxml)
    finish = task.find('ns:Finish', nsxml)
    wbs = task.find('ns:WBS', nsxml)
    dur_text = dur.text if dur is not None else '0'
    d_hours = 0
    if dur_text and dur_text.startswith('PT'):
        dc = dur_text.replace('PT','').replace('H0M0S','').replace('H','')
        try: d_hours = int(dc)
        except: pass
    start_t = start.text[:10] if start is not None else ''
    finish_t = finish.text[:10] if finish is not None else ''
    wbs_t = wbs.text if wbs is not None else ''
    name_t = name.text if name is not None else ''
    phase = '?'
    if wbs_t.startswith('1.1'): phase = 'INICIO'
    elif wbs_t.startswith('1.2'): phase = 'PLANIF.'
    elif wbs_t.startswith('1.3'): phase = 'EJEC.'
    elif wbs_t.startswith('1.4'): phase = 'M&C'
    elif wbs_t.startswith('1.5'): phase = 'CIERRE'
    l4_tasks.append((uid_int, phase, name_t[:70], f"{d_hours}h", start_t, finish_t))

n = len(l4_tasks)
NSLIDES = 5
chunk = n // NSLIDES  # 22
rem = n % NSLIDES    # 0
chunks = []
idx = 0
for i in range(NSLIDES):
    extra = 1 if i < rem else 0
    chunks.append(l4_tasks[idx:idx + chunk + extra])
    idx += chunk + extra

cols_w = [Emu(457200), Emu(594360), Emu(5029200), Emu(502920), Emu(1828800), Emu(1828800)]

for i, ch in enumerate(chunks):
    slide = prs.slides[28 + i]  # indices 28-32 = slides 29-33
    clear_slide(slide)
    sn = 29 + i
    add_header(slide, sn, f"LISTA DE ACTIVIDADES ({i+1}/{NSLIDES})",
               f"{len(ch)} de {n} actividades nivel 4 | PMBOK 8va, Cronograma")
    add_table(slide, MARGIN, Emu(1280160), CONTENT_W, Emu(5300000),
        ["UID", "Fase", "Actividad", "Dur.", "Inicio", "Fin"], ch, cols_w, Pt(8))
    txBox = slide.shapes.add_textbox(MARGIN, Emu(6620000), CONTENT_W, Emu(180000))
    tf = txBox.text_frame; p = tf.paragraphs[0]
    p.text = f"{len(ch)} de {n} actividades nivel 4 | Fuente: LoRA-HDL-QA_Proyecto_2026.xml"
    p.font.size = Pt(7); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = 'Calibri'
print(f'✓ Tasks: {NSLIDES} slides, {[len(c) for c in chunks]}')

# ═══════════════════════════════════════════════════
# 7. Rebuild slides 34-37 (shifted after 5 task slides)
#   34: Atributos | 35: Base+Estimaciones | 36: Lista Hitos | 37: Dominio Recursos
# ═══════════════════════════════════════════════════

# Slide 34: Atributos
slide = prs.slides[33]; clear_slide(slide)
add_header(slide, 34, "ATRIBUTOS DE LA ACTIVIDAD",
           "Campos que describen cada actividad del cronograma | PMBOK 8va, Cronograma")
attr_data = [
    ("UID", "Identificador único en el XML de ProjectLibre", "Ej: 4, 5, 68"),
    ("Nombre", "Descripción de la actividad (verbo + objeto)", "Revisar normativo y reglamento de graduación"),
    ("Nivel EDT", "1=Proyecto, 2=Fase, 3=Entregable, 4=Paquete trabajo", "1, 2, 3, 4"),
    ("Código WBS", "Código de la EDT al que pertenece la actividad", "1.1.1.4 (Inscripción del Tema)"),
    ("Duración", "Horas de trabajo efectivas (formato ISO 8601)", "PT8H0M0S = 8 horas"),
    ("Fecha Inicio", "Fecha y hora de inicio planificada", "2026-02-02T08:00:00"),
    ("Fecha Fin", "Fecha y hora de finalización planificada", "2026-02-02T17:00:00"),
    ("Predecesoras", "Tareas que deben completarse antes (FS)", "4FS (termina UID 4, inicia esta)"),
    ("Recurso Asignado", "Recurso que ejecuta la tarea", "Pablo Flores (UID 1)"),
    ("Tipo de Tarea", "Fixed Work / Fixed Duration / Fixed Units", "Fixed Work (todas las tareas)"),
    ("Calendario", "Calendario aplicable a la tarea", "UID 11: Guatemala - Feriados 2026"),
    ("Hito", "TRUE si duración = 0, FALSE si no", "FALSE (110 tareas de trabajo)"),
]
add_table(slide, MARGIN, Emu(1371600), CONTENT_W, Emu(4937760),
    ["Atributo", "Descripción", "Valores de Ejemplo / Notas"],
    attr_data, [Emu(2560320), Emu(4937760), Emu(3657600)], Pt(10))

# Slide 35: Base + Estimaciones Duración
slide = prs.slides[34]; clear_slide(slide)
add_header(slide, 35, "BASE DE LAS ESTIMACIONES Y ESTIMACIONES DE DURACIÓN",
           "Fundamento de cada estimación + resumen por fase | PMBOK 8va, Cronograma")
add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2377440),
    'Base de las Estimaciones — Duraciones', [
        'Metodología: Bottom-up desde EDT nivel 4. Unidad: horas efectivas (netas, sin pausas).',
        'Fuentes: Juicio de expertos (asesor FING), análoga (tesis previas 150-360h),',
        'paramétrica (~15h por sección de 5pp), descomposición (paquetes de 3-31h).',
        '',
        'Ajustes aplicados:',
        '• Productividad: 8h/día netas | Calendario: lun-sáb, excluye domingos y feriados',
        '• Curva de aprendizaje: +20% buffer en primeras tareas de cada fase nueva',
        '• Hardware M1 8GB: batch=1 → 2-3x tiempo vs GPU, quantización 4-bit obligatoria',
        '• Precisión: ±10% en paquetes nivel 4; ±15% en fases nivel 3',
    ], Pt(10), Pt(8))
add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2377440),
    'Estimaciones de Duración por Fase', [
        'INICIO (148h): Inscripción Tema 78h + Asignación Asesor 80h + Plan Proyecto 60h',
        'PLANIFICACIÓN (235h): Anteproyecto 152h + Diseño Metodológico 83h',
        'EJECUCIÓN (744h): Marco Teórico 161h + Preprocesamiento 136h +',
        '  Entrenamiento LoRA 227h + Evaluación Benchmark 106h +',
        '  Análisis Resultados 90h + Redacción Informe 175h',
        'M&C (86h): Seguimiento + Control Cambios (actividades solapadas)',
        'CIERRE (160h): Correcciones Finales 75h + Defensa Pública 59h + Entrega 26h',
    ], Pt(10), Pt(8))
add_card(slide, Emu(457200), Emu(3931920), Emu(10972800), Emu(1645920),
    'Ejemplos de Cálculo de Duraciones', [
        'Investigación/Lectura: 2-3h/paper → 23h = ~10 papers | Redacción: 15h/sección 5pp | Programación: 15-23h/módulo',
        'Entrenamiento ML: 31h = ~3 épocas M1 8GB batch=1 modelo 4B | Revisión: 15h/ronda completa (leer+corregir+verificar)',
        'Períodos de espera: 31-39h (3-5 días hábiles FING) modelados como tareas de duración fija',
    ], Pt(10), Pt(8))

# Slide 36: Lista de Hitos
slide = prs.slides[35]; clear_slide(slide)
add_header(slide, 36, "LISTA DE HITOS",
           "Hitos principales con fechas planificadas y criterios de aceptación | PMBOK 8va, Cronograma")
hitos = [
    ("1", "Proyecto de Graduación Completado", "2026-09-14", "Aprobación formal por tribunal y constancia Biblioteca FING"),
    ("2", "INICIO DEL PROYECTO", "2026-02-02", "Carta de inscripción firmada por Director de Carrera"),
    ("3", "INICIO: INICIO", "2026-02-02", "Tema aprobado + asesor identificado formalmente"),
    ("4", "FIN: INICIO", "2026-02-26", "Plan de proyecto aprobado por el asesor"),
    ("5", "INICIO: PLANIFICACION", "2026-02-26", "Inicio formal del anteproyecto con asesor asignado"),
    ("6", "FIN: PLANIFICACION", "2026-04-08", "Anteproyecto y diseño metodológico aprobados por comité"),
    ("7", "INICIO: EJECUCION", "2026-04-08", "Inicio formal de experimentos y marco teórico"),
    ("8", "FIN: EJECUCION", "2026-08-17", "Informe final completo revisado y aprobado por asesor"),
    ("9", "INICIO: M&C", "2026-02-26", "Primera reunión de seguimiento con asesor"),
    ("10", "FIN: M&C", "2026-03-13", "Cierre formal de actividades de seguimiento intensivo"),
    ("11", "INICIO: CIERRE", "2026-08-17", "Inicio de correcciones finales post-revisión asesor"),
    ("12", "FIN: CIERRE", "2026-09-14", "Documento final depositado en Biblioteca FING"),
]
add_table(slide, MARGIN, Emu(1371600), CONTENT_W, Emu(4937760),
    ["#", "Hito", "Fecha", "Criterio de Aceptación"],
    hitos, [Emu(548640), Emu(3657600), Emu(1828800), Emu(5029200)], Pt(10))

# Slide 37: Recursos Domain Separator
add_domain_slide(prs.slides[36],
    "DOMINIO DE DESEMPEÑO:\nRECURSOS",
    "Identificación, adquisición y gestión de los recursos necesarios para ejecutar el proyecto. Incluye recursos humanos, equipamiento, materiales e infraestructura.")

# ═══════════════════════════════════════════════════
# 8. Expand cards in d38-d55 (shifted from 36-55)
#    Slide 38 (idx 37): Plan Gestión Recursos
#    Slide 39 (idx 38): Carta Equipo
#    etc.
# ═══════════════════════════════════════════════════

# Slide 38: Plan de Gestión de Recursos (was slide 36, now 38)
slide = prs.slides[37]; clear_slide(slide)
add_header(slide, 38, "PLAN DE GESTIÓN DE RECURSOS",
           "Define cómo se estimarán, adquirirán y gestionarán los recursos | PMBOK 8va")
add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2743200),
    "Identificación y Estimación de Recursos", [
        "Fuente primaria: EDT nivel 4 — cada una de las 110 actividades tiene recurso asignado",
        "Método: Bottom-up — suma de horas por recurso desde actividades nivel 4",
        "",
        "Recurso Humano:",
        "• Pablo Flores (PM/Tesista) — 1,287h planificadas @ Q 80/hr = Q 102,960",
        "• Asesor FING — ~40h estimadas para revisión técnica (sin costo directo para el proyecto)",
        "",
        "Recursos Físicos:",
        "• MacBook Air M1 2020 (8GB/256GB) — Q 6,000 (costo hundido, ya adquirido)",
        "• Internet 10Mbps (~Q 250/mes) y electricidad (~Q 300 total) — costos absorbidos",
        "",
        "Materiales:",
        "• Impresión 3 copias × Q 150 = Q 450 | Posters defensa: Q 200 | Papelería: Q 150",
        "• Total materiales: Q 800 — compra directa por bajo monto, sin licitación requerida",
    ], Pt(11), Pt(9))
add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2743200),
    "Estrategia de Gestión de Recursos", [
        "Asignación: PM dedicado 100% al proyecto (225 días, 1,287h efectivas)",
        "Calendario: lun-sáb 8h/día, flexibilidad ±2h en semanas críticas",
        "",
        "Adquisiciones:",
        "• MacBook ya adquirido — no requiere proceso de compra",
        "• Materiales: compra directa por bajo monto (< Q 1,000), sin licitación",
        "",
        "Monitoreo y Control:",
        "• Registro de horas reales vs planificadas como parte del EVM mensual",
        "• Reporte mensual de utilización de recursos al asesor",
        "",
        "Contingencias:",
        "• Si MacBook falla → usar laboratorio de cómputo FING (costo: Q 0)",
        "• Si PM enfrenta indisponibilidad > 48h → ajustar cronograma, notificar asesor",
    ], Pt(11), Pt(9))

# Slide 39: Carta del Equipo (was 37, now 39)
slide = prs.slides[38]; clear_slide(slide)
add_header(slide, 39, "CARTA DEL EQUIPO",
           "Valores, normas y acuerdos de trabajo | PMBOK 8va, Recursos")
add_card(slide, Emu(457200), Emu(1371600), Emu(10972800), Emu(5029200),
    "Carta del Equipo — Proyecto LoRA-HDL-QA", [
        "Contexto: Proyecto de graduación individual. El 'equipo' está compuesto por:",
        "• Pablo Flores — PM/Tesista (ejecutor principal, 100% del trabajo técnico y de gestión)",
        "• Asesor (por asignar) — Revisor técnico (~40h totales, sin costo directo para el proyecto)",
        "• Director de Carrera — Autoridad de aprobación y escalamiento de decisiones críticas",
        "",
        "Valores del Trabajo:",
        "• Excelencia académica: rigor científico en experimentación y documentación de resultados",
        "• Transparencia: reportar avances y desviaciones en tiempo real al asesor",
        "• Responsabilidad: cumplir hitos y plazos acordados; notificar riesgos con anticipación",
        "• Reproducibilidad: todo código, datos y resultados deben ser reproducibles por terceros",
        "",
        "Normas de Trabajo:",
        "• Reuniones: semanal con asesor (1h fija), quincenal de seguimiento formal con minuta",
        "• Comunicación: email para documentos formales, WhatsApp para coordinación ágil",
        "• Documentación: minutas de reunión en máximo 24h, reportes de avance mensuales con EVM",
        "• Código: versionado en GitHub, commits atómicos con mensajes descriptivos en español",
        "• Conflictos: escalar al Director de Carrera si PM y asesor no llegan a acuerdo en 5 días",
        "",
        "Nota: Por ser proyecto unipersonal, esta carta establece normas de trabajo entre el PM y el asesor.",
        "No se requiere team charter tradicional (no hay equipo multidisciplinario que gestionar)."
    ], Pt(11), Pt(9))

# Slide 44: Calendarios de Recursos (was 42, now 44)
slide = prs.slides[43]; clear_slide(slide)
add_header(slide, 44, "CALENDARIOS DE RECURSOS",
           "Disponibilidad de cada recurso a lo largo del proyecto | PMBOK 8va, Recursos")
add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2560320),
    "Calendario del Recurso Principal — Pablo Flores (PM/Tesista)", [
        "Período: Feb 02 — Sep 14, 2026 (225 días calendario)",
        "Días laborables: Lunes a Sábado (6 días/semana)",
        "Horario estándar: 08:00-12:00 y 13:00-17:00 (8h/día efectivas)",
        "Horas máximas por semana: 48h (6 días × 8h)",
        "",
        "Excepciones (no laborable para el recurso):",
        "• Todos los domingos del período",
        "• Jueves 2 — Sábado 4 de Abril (Semana Santa)",
        "• Viernes 1 de Mayo (Día del Trabajo)",
        "• Martes 30 de Junio (Día del Ejército, ½ día opcional)",
        "",
        "Disponibilidad real: ~190 días laborables × 8h = ~1,520h disponibles",
        "Horas planificadas: 1,287h → utilización del ~85% del tiempo disponible"
    ], Pt(10), Pt(8))
add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2011680),
    "Calendario del Equipamiento — MacBook Air M1", [
        "Disponible 24/7, compartido con uso personal del tesista",
        "Ventanas óptimas para entrenamiento ML:",
        "• Horario nocturno (18:00-06:00): menor temperatura ambiente",
        "• Fines de semana: bloques de 8-12h continuos disponibles",
        "Horas estimadas de uso para el proyecto: ~1,000h"
    ], Pt(10), Pt(8))
add_card(slide, Emu(457200), Emu(4114800), Emu(5486400), Emu(2011680),
    "Calendario del Asesor", [
        "Disponibilidad: Sujeta a carga académica del docente en FING",
        "Reuniones: 1h/semana en día y hora fijos (por acordar al inicio)",
        "Revisiones de hitos: bajo demanda, con 5 días hábiles de anticipación",
        "Período de mayor carga docente: inicio de semestre (Ene, Jul)",
        "Planificar revisiones de documentos extensos con 1-2 semanas de anticipación"
    ], Pt(10), Pt(8))

# Slide 51: Estrategia de Financiamiento (was 49, now 51)
slide = prs.slides[50]; clear_slide(slide)
add_header(slide, 51, "ESTRATEGIA DE FINANCIAMIENTO",
           "Cómo se obtendrán y administrarán los fondos | PMBOK 8va, Costos")
add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(5029200),
    "Estrategia de Financiamiento", [
        "Fuente primaria: Recursos propios del tesista (Pablo Flores)",
        "",
        "El proyecto no requiere financiamiento externo porque:",
        "• El hardware principal (MacBook M1) ya fue adquirido previamente (costo hundido)",
        "• El software es 100% open-source y gratuito (Python, HuggingFace, MLX, ProjectLibre)",
        "• La mano de obra es aportada por el tesista como parte de su formación académica",
        "• Los costos de materiales son menores a Q 1,000 (impresión, empastado, posters)",
        "",
        "Estructura de costos:",
        "• Costos hundidos (ya incurridos): MacBook M1 Q 6,000",
        "• Costos futuros de materiales: Q 800 (impresión 3 copias, posters, papelería)",
        "• Costos recurrentes absorbidos: internet Q 250/mes, electricidad ~Q 300 total",
        "• Costo de oportunidad (mano de obra): 1,287h × Q 80/hr = Q 102,960",
        "",
        "Total erogación en efectivo requerida durante el proyecto: Q 800 (materiales fase cierre)",
        "",
        "No se requiere plan de financiamiento con bancos, inversionistas, ni patrocinadores.",
        "El proyecto es autofinanciado por el tesista con recursos ya disponibles."
    ], Pt(10), Pt(8))
add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(5029200),
    "Administración de Fondos", [
        "Responsable: PM (Pablo Flores)",
        "",
        "Mecanismo de pago:",
        "• Cuenta bancaria personal existente (no se requiere cuenta dedicada)",
        "• Pago directo a proveedores (copistería) contra entrega del servicio",
        "",
        "Control financiero:",
        "• Registro de gastos en Excel vinculado al EVM mensual",
        "• Facturas y recibos escaneados y archivados en Google Drive del proyecto",
        "• Conciliación mensual: gasto real vs presupuesto planificado",
        "",
        "Flujo de fondos:",
        "• Feb 2026: Q 0 (MacBook ya fue pagado previamente)",
        "• Sep 2026: Q 800 (impresión de 3 copias + posters para defensa)",
        "• Reservas: no requieren depósito previo (solo se usan si se materializan riesgos)",
        "",
        "Cierre financiero (Sep 2026):",
        "• Conciliación final de gastos reales vs presupuesto",
        "• Documentación de lecciones aprendidas en gestión financiera",
        "• Informe final de costos para el expediente del proyecto de graduación"
    ], Pt(10), Pt(8))

# Slide 55: Reporte de Riesgos (was 53, now 55)
slide = prs.slides[54]; clear_slide(slide)
add_header(slide, 55, "REPORTE DE RIESGOS",
           "Resumen ejecutivo de la situación de riesgos | PMBOK 8va, Riesgos")
add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2560320),
    "Resumen General de Riesgos", [
        "Total riesgos identificados: 25",
        "Riesgos con plan de respuesta activo (CE > Q 1,500): 12",
        "Riesgos en monitoreo pasivo (CE < Q 500): 5",
        "Riesgos escalables al Director: 3",
        "",
        "Costo Esperado Total (EMV): Q 20,846.00",
        "Reserva de contingencia recomendada: Q 20,846.00",
        "Reserva de gestión: Q 700.00",
        "",
        "Perfil de riesgo del proyecto: MODERADO",
        "• Mayoría de riesgos son técnicos (15/25 = 60%)",
        "• Los riesgos organizacionales tienen alto impacto pero baja probabilidad",
        "• No se identificaron riesgos catastróficos (impacto > Q 10,000)",
        "• El proyecto es manejable con las reservas calculadas y planes de respuesta activos"
    ], Pt(10), Pt(8))
add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2560320),
    "Distribución por Categoría", [
        "Técnico: 15 riesgos | CE: Q 13,994 (67.1%)",
        "  → Críticos: OOM en 8GB, hardware M1 insuficiente, falla HW, no supera baseline",
        "Organizacional: 7 riesgos | CE: Q 4,950 (23.7%)",
        "  → Críticos: Reescritura por observaciones, demora institucional, rechazo comité",
        "Calidad: 1 riesgo | CE: Q 900 (4.3%) → Corrupción no detectada en datos",
        "Recursos Humanos: 1 riesgo | CE: Q 882 (4.2%) → Retraso asignación asesor",
        "Operativo: 1 riesgo | CE: Q 120 (0.6%) → Error en impresión/empastado"
    ], Pt(10), Pt(8))
add_card(slide, Emu(457200), Emu(4114800), Emu(10972800), Emu(2011680),
    "Top 5 Riesgos por Costo Esperado", [
        "1. Modelo OOM en 8GB RAM (CE: Q 1,600) → Quantización 4-bit GGUF, probar modelos más pequeños como Qwen3-1.5B",
        "2. Arquitectura no viable en M1 (CE: Q 1,500) → Evaluar QLoRA 4-bit con MLX; implementar gradient checkpointing",
        "3. Falla hardware durante entrenamiento (CE: Q 1,500) → Checkpointing cada 100 steps; backup en iCloud/Google Drive",
        "4. Modelo no supera baseline (CE: Q 1,500) → Ajuste de hiperparámetros (r=16→64, learning rates); data augmentation",
        "5. Thermal throttling en M1 (CE: Q 1,350) → Entrenamiento en horarios nocturnos; ventilador externo; dividir en sub-épocas"
    ], Pt(10), Pt(8))

print('✓ All cards expanded')

# ═══════════════════════════════════════════════════
# 9. Center all tables
# ═══════════════════════════════════════════════════
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_table:
            shape.width = CONTENT_W
            shape.left = MARGIN
print('✓ All tables centered')

# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
out = f'{BASE}/LoRA_Plan_Direccion_Proyecto_Completo_v2.pptx'
prs.save(out)
print(f'\n✅ Saved: {out}')
print(f'   Slides: {len(prs.slides)}')
print(f'   Tasks: {n} in {NSLIDES} slides ({[len(c) for c in chunks]})')
