#!/usr/bin/env python3
"""
Fix script for LoRA Plan de Dirección presentation.
Updates: Calendar, EDT/WBS tree, RBS, complete task list.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
import os, sys

# ── Constants ──
NAVY = RGBColor(0x1A, 0x27, 0x44)
BLUE = RGBColor(0x42, 0x85, 0xF4)
BLUE_LIGHT = RGBColor(0xC8, 0xDC, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
MED_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
ACCENT_BLUE = RGBColor(0x42, 0x85, 0xC8)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

def add_header_bar(slide, slide_num, title, subtitle=None):
    """Add standard navy header bar"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Emu(822960))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Emu(228600), Emu(182880), Emu(548640), Emu(502920))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"{slide_num:02d}"; p.font.size = Pt(36)
    p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Calibri'
    
    txBox2 = slide.shapes.add_textbox(Emu(868680), Emu(228600), Emu(10972800), Emu(457200))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]; p2.text = title; p2.font.size = Pt(22)
    p2.font.bold = True; p2.font.color.rgb = WHITE; p2.font.name = 'Calibri'
    
    if subtitle:
        txBox3 = slide.shapes.add_textbox(Emu(274320), Emu(868680), Emu(11612880), Emu(365760))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]; p3.text = subtitle; p3.font.size = Pt(14)
        p3.font.color.rgb = GRAY_TEXT; p3.font.name = 'Calibri'

def add_card(slide, left, top, width, height, title, body_lines, title_size=Pt(12), body_size=Pt(9)):
    """Add a rounded rectangle card"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = RGBColor(0xDE, 0xE2, 0xE6); card.line.width = Pt(0.5)
    
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(68580); tf.margin_right = Emu(68580)
    tf.margin_top = Emu(45720); tf.margin_bottom = Emu(45720)
    
    p = tf.paragraphs[0]; p.text = title
    p.font.size = title_size; p.font.bold = True; p.font.color.rgb = NAVY
    p.font.name = 'Calibri'; p.space_after = Pt(3)
    
    for line in body_lines:
        p2 = tf.add_paragraph(); p2.text = line
        p2.font.size = body_size; p2.font.color.rgb = MED_TEXT
        p2.font.name = 'Calibri'; p2.space_after = Pt(1)
    return card

def add_simple_table(slide, left, top, width, height, headers, rows, col_widths=None, font_size=Pt(8)):
    """Add a table with alternating row colors"""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = font_size; para.font.bold = True
            para.font.color.rgb = WHITE; para.font.name = 'Calibri'
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = cell._tc.makeelement(qn('a:solidFill'), {})
        srgbClr = cell._tc.makeelement(qn('a:srgbClr'), {'val': '4285C8'})
        solidFill.append(srgbClr); tcPr.append(solidFill)
    
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c); cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = font_size; para.font.color.rgb = MED_TEXT
                para.font.name = 'Calibri'
            if r % 2 == 0:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = cell._tc.makeelement(qn('a:solidFill'), {})
                srgbClr = cell._tc.makeelement(qn('a:srgbClr'), {'val': 'F8F9FA'})
                solidFill.append(srgbClr); tcPr.append(solidFill)
    
    return table_shape

def add_edt_node(slide, left, top, width, height, text, fill_color=ACCENT_BLUE, text_color=WHITE, font_size=Pt(8)):
    """Add an EDT tree node box"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = fill_color
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(22860); tf.margin_right = Emu(22860)
    tf.margin_top = Emu(11430); tf.margin_bottom = Emu(11430)
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line; p.font.size = font_size; p.font.bold = True
        p.font.color.rgb = text_color; p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
    return box

def add_connector(slide, x1, y1, x2, y2, color=RGBColor(0x99, 0x99, 0x99)):
    """Add a simple line connector"""
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color; conn.line.width = Pt(1)
    return conn

# ════════════════════════════════════════════════════════════════
# Load existing presentation
# ════════════════════════════════════════════════════════════════
src = '/Users/pabloflores/Documents/estructura/revision/entregables_finales/LoRA_Plan_Direccion_Proyecto_Completo.pptx'
prs = Presentation(src)
blank_layout = prs.slide_layouts[6]

# Get reference to slides we need to replace
# The slides to replace:
# - Slide 15 (index 14): EDT/WBS → tree diagram
# - Slide 28 (index 27): Calendarios del Proyecto → full calendar
# - Slide 29 (index 28): Lista de Actividades → ALL tasks (need to add extra slides)
# - Slide 41 (index 40): RBS → fixed: only MacBook + Q80/hr

# We'll delete and recreate these slides
# But python-pptx doesn't have a clean delete-and-insert. 
# Strategy: load the XML of the slide and replace its shapes.
# Actually, easier: use the slide's _element to clear shapes and rebuild.

def clear_slide(slide):
    """Remove all shapes from a slide"""
    spTree = slide.shapes._spTree
    for sp in list(spTree):
        if sp.tag != qn('p:cSld'):  
            spTree.remove(sp)

# ── Fix Slide 15: EDT/WBS Tree Diagram ──
slide = prs.slides[14]
clear_slide(slide)
add_header_bar(slide, 15, "EDT / WBS — ESTRUCTURA DE DESGLOSE DEL TRABAJO",
               "Descomposición jerárquica del alcance total en 16 paquetes de trabajo | PMBOK 8va, Alcance")

# "Diagrama EDT" label
txBox = slide.shapes.add_textbox(Emu(274320), Emu(1234440), Emu(1828800), Emu(274320))
tf = txBox.text_frame; p = tf.paragraphs[0]
p.text = "Diagrama EDT"; p.font.size = Pt(12); p.font.bold = True
p.font.color.rgb = NAVY; p.font.name = 'Calibri'

# Root node
edt_root_x = Emu(4572000)
add_edt_node(slide, edt_root_x, Emu(1554480), Emu(2743200), Emu(457200),
             "Proyecto de Graduación\nLoRA-HDL-QA", NAVY, WHITE, Pt(10))

# Phase level nodes (level 2)
phases = [
    ("INICIO", Emu(274320)),
    ("PLANIFICACIÓN", Emu(2468880)),
    ("EJECUCIÓN", Emu(4663440)),
    ("MONITOREO\nY CONTROL", Emu(6858000)),
    ("CIERRE", Emu(9052560)),
]
phase_y = Emu(2286000)
phase_w = Emu(2011680)
phase_h = Emu(457200)

for name, x in phases:
    add_edt_node(slide, x, phase_y, phase_w, phase_h, name, ACCENT_BLUE, WHITE, Pt(9))

# Connectors from root to phases
root_bottom = Emu(2011680)
root_center_x = Emu(5943600)
for name, x in phases:
    phase_center_x = x + phase_w // 2
    add_connector(slide, root_center_x, root_bottom, phase_center_x, phase_y)

# Work package nodes (level 3) - 16 packages
# INICIO children
wp_y = Emu(2926080)
wp_w = Emu(835000)
wp_h_small = Emu(502920)

inicio_children = [
    ("Inscripción\ndel Tema", Emu(354320)),
    ("Asignación\nde Asesor", Emu(1280160)),
    ("Plan del\nProyecto", Emu(2206000)),
]
for name, x in inicio_children:
    add_edt_node(slide, x, wp_y, wp_w, wp_h_small, name, BLUE_LIGHT, NAVY, Pt(7))

# PLANIFICACION children
plan_children = [
    ("Anteproyecto", Emu(2971800)),
    ("Diseño\nMetodológico", Emu(3840480)),
]
for name, x in plan_children:
    add_edt_node(slide, x, wp_y, wp_w, wp_h_small, name, BLUE_LIGHT, NAVY, Pt(7))

# EJECUCION children (6 packages - needs two rows)
ejec_children_row1 = [
    ("Marco\nTeórico", Emu(4937760)),
    ("Preproc.\nde Datos", Emu(5817600)),
    ("Entrenam.\nLoRA", Emu(6697440)),
]
ejec_children_row2 = [
    ("Evaluación\nBenchmark", Emu(4937760)),
    ("Análisis\nResultados", Emu(5817600)),
    ("Redacción\nInforme Final", Emu(6697440)),
]
for name, x in ejec_children_row1:
    add_edt_node(slide, x, wp_y, Emu(810000), wp_h_small, name, BLUE_LIGHT, NAVY, Pt(7))
for name, x in ejec_children_row2:
    add_edt_node(slide, x, Emu(3570480), Emu(810000), wp_h_small, name, BLUE_LIGHT, NAVY, Pt(7))

# M&C children
mc_children = [
    ("Seguimiento", Emu(6938000)),
    ("Control de\nCambios", Emu(7854000)),
]
for name, x in mc_children:
    add_edt_node(slide, x, wp_y, wp_w, wp_h_small, name, BLUE_LIGHT, NAVY, Pt(7))

# CIERRE children
cierre_children = [
    ("Correcciones\nFinales", Emu(9150000)),
    ("Defensa\nPública", Emu(10020000)),
    ("Entrega\ny Cierre", Emu(10890000)),
]
for name, x in cierre_children:
    add_edt_node(slide, x, wp_y, Emu(850000), wp_h_small, name, BLUE_LIGHT, NAVY, Pt(7))

# Phase → child connectors
for name, phase_x in phases:
    pcx = phase_x + phase_w // 2
    py_bottom = phase_y + phase_h
    children = []
    if name == "INICIO":
        children = [Emu(354320), Emu(1280160), Emu(2206000)]
    elif name == "PLANIFICACIÓN":
        children = [Emu(2971800), Emu(3840480)]
    elif name == "EJECUCIÓN":
        children = [Emu(4937760), Emu(5817600), Emu(6697440)]
    elif name == "MONITOREO\nY CONTROL":
        children = [Emu(6938000), Emu(7854000)]
    elif name == "CIERRE":
        children = [Emu(9150000), Emu(10020000), Emu(10890000)]
    
    for cx in children:
        child_center_x = cx + wp_w // 2
        add_connector(slide, pcx, py_bottom, child_center_x, wp_y)

# Additional connectors for EJECUCION row 2
ejec_phase_x = Emu(4663440)
ejec_phase_cx = ejec_phase_x + phase_w // 2
ejec_py_bottom = phase_y + phase_h
for cx in [Emu(4937760), Emu(5817600), Emu(6697440)]:
    child_center_x = cx + Emu(810000) // 2
    child_top_y = Emu(3570480)
    add_connector(slide, ejec_phase_cx, ejec_py_bottom, child_center_x, child_top_y)

# Summary stats below
add_card(slide, Emu(457200), Emu(4200000), Emu(10972800), Emu(822960),
         "Resumen EDT", [
             "Nivel 1: 1 proyecto (LoRA-HDL-QA) | Nivel 2: 5 fases | Nivel 3: 16 entregables | Nivel 4: 110 paquetes de trabajo",
             "Total horas: 1,287h | Rango por paquete nivel 4: 3h — 31h | Método: Bottom-up desde cronograma XML (ProjectLibre 1.9.8)",
             "Los 16 paquetes de trabajo del nivel 3 se desglosan en 110 actividades nivel 4. Diccionario EDT detalla cada una."
         ], Pt(10), Pt(8))

# ── Fix Slide 28 (index 27): Calendarios del Proyecto ──
slide = prs.slides[27]
clear_slide(slide)
add_header_bar(slide, 28, "CALENDARIOS DEL PROYECTO",
               "Calendario guatemalteco completo 2026 con feriados oficiales | PMBOK 8va, Cronograma")

add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2560320),
         "Calendario Base — Guatemala 2026", [
             "Días laborables: Lunes a Sábado",
             "Horario: 08:00-12:00 y 13:00-17:00 (8h/día efectivas, 48h/semana máx.)",
             "Domingo: No laborable (día de descanso obligatorio)",
             "",
             "FERIADOS OFICIALES 2026 (no laborables):",
             "• Miércoles 1 de Enero — Año Nuevo",
             "• Jueves 2 de Abril — Jueves Santo (Semana Santa)",
             "• Viernes 3 de Abril — Viernes Santo (Semana Santa)",
             "• Sábado 4 de Abril — Sábado de Gloria (Semana Santa)",
             "• Viernes 1 de Mayo — Día del Trabajo",
             "• Martes 30 de Junio — Día del Ejército",
             "• Martes 15 de Septiembre — Día de la Independencia",
             "• Martes 20 de Octubre — Día de la Revolución",
             "• Domingo 1 de Noviembre — Día de Todos los Santos",
             "• Jueves 24 de Diciembre — Nochebuena (medio día)",
             "• Viernes 25 de Diciembre — Navidad",
             "• Jueves 31 de Diciembre — Fin de Año (medio día)",
             "",
             "Feriados que afectan al proyecto (Feb—Sep 2026):",
             "Semana Santa (2-4 Abr), Día del Trabajo (1 May),",
             "Día del Ejército (30 Jun), Independencia (15 Sep)"
         ], Pt(10), Pt(8))

add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2560320),
         "Recurso — Pablo Flores", [
             "Disponibilidad: 100% (estudiante tiempo completo)",
             "Máximo: 8h/día laborable estándar",
             "Flexibilidad: hasta 10h/día en semanas críticas",
             "Período activo: Feb 02 — Sep 14, 2026 (225 días)",
             "Total horas planificadas: 1,287h efectivas",
             "Promedio: ~5.7h/día (incluye esperas y buffers)",
             "",
             "Notas sobre el Calendario:",
             "• Feriados se modelan como días no laborables",
             "• NOTA TÉCNICA: ProjectLibre 1.9.8 tiene un bug con",
             "  <Exceptions> sin <RecurrenceType> que causa crash.",
             "  Los feriados están documentados aquí pero NO en el",
             "  XML para garantizar compatibilidad.",
             "• Alternativa: ajustar manualmente fechas de tareas",
             "  que caen en feriados durante el seguimiento."
         ], Pt(10), Pt(8))

add_card(slide, Emu(457200), Emu(4114800), Emu(10972800), Emu(1371600),
         "Períodos de Espera Institucionales (modelados como tareas en el cronograma)", [
             "• Resolución institucional: 39h (Feb 6 — Feb 13) | • Aprobación comité: 31h (Mar 19 — Mar 25)",
             "• Retroalimentación 1ra revisión: 23h (Jul 29 — Ago 3) | • Retroalimentación 2da revisión: 39h (Ago 6 — Ago 12)",
             "• Resolución jurado: 23h (Sep 4 — Sep 9)",
             "Total esperas: ~155h (12% del total de horas). Estas esperas son tiempo calendario, no trabajo activo del tesista.",
             "Impacto en cronograma: Las esperas institucionales están en la ruta crítica. Cualquier extensión afecta la fecha de finalización."
         ], Pt(10), Pt(8))

# ── Fix Slide 29 (index 28): Lista de Actividades → ALL 110 tasks ──
# We need 2 slides for all 110 tasks. Replace slide 29 and add slide 30a.
slide = prs.slides[28]
clear_slide(slide)
add_header_bar(slide, 29, "LISTA DE ACTIVIDADES (1/2)",
               "Actividades nivel 4 del cronograma — 55 de 110 | PMBOK 8va, Cronograma")

# Import task data
import xml.etree.ElementTree as ET
tree = ET.parse('/Users/pabloflores/Documents/estructura/revision/entregables_finales/LoRA-HDL-QA_Proyecto_2026.xml')
root_xml = tree.getroot()
ns = {'ns': 'http://schemas.microsoft.com/project'}

all_tasks_xml = root_xml.findall('.//ns:Task', ns)
tasks_data = []
for task in all_tasks_xml:
    uid = task.find('ns:UID', ns)
    uid_int = int(uid.text)
    if uid_int > 133:
        continue
    name = task.find('ns:Name', ns)
    dur = task.find('ns:Duration', ns)
    start = task.find('ns:Start', ns)
    finish = task.find('ns:Finish', ns)
    outline = task.find('ns:OutlineLevel', ns)
    wbs = task.find('ns:WBS', ns)
    
    lvl = int(outline.text) if outline is not None else 0
    dur_text = dur.text if dur is not None else '0'
    d_hours = 0
    if dur_text and dur_text.startswith('PT'):
        d_clean = dur_text.replace('PT','').replace('H0M0S','').replace('H','')
        try: d_hours = int(d_clean)
        except: pass
    
    start_t = start.text[:10] if start is not None else ''
    finish_t = finish.text[:10] if finish is not None else ''
    wbs_t = wbs.text if wbs is not None else ''
    name_t = name.text if name is not None else ''
    tasks_data.append((uid_int, lvl, wbs_t, name_t, d_hours, start_t, finish_t))

# Filter level 4 tasks only
l4_tasks = [(uid, wbs, name, hrs, start, finish) for uid, lvl, wbs, name, hrs, start, finish in tasks_data if lvl == 4]

# Map WBS to phase name
def get_phase(wbs):
    if wbs.startswith('1.1.1') or wbs.startswith('1.1.2') or wbs.startswith('1.1.3'):
        return 'INICIO'
    elif wbs.startswith('1.2'):
        return 'PLANIF.'
    elif wbs.startswith('1.3'):
        return 'EJEC.'
    elif wbs.startswith('1.4'):
        return 'M&C'
    elif wbs.startswith('1.5'):
        return 'CIERRE'
    return '?'

# Split into two halves
half = len(l4_tasks) // 2
first_half = l4_tasks[:half]
second_half = l4_tasks[half:]

rows_slide1 = [(uid, get_phase(wbs), name[:60], f"{hrs}h", start, finish) for uid, wbs, name, hrs, start, finish in first_half]
add_simple_table(slide, Emu(91440), Emu(1371600), Emu(12009160), Emu(5200000),
    ["UID", "Fase", "Actividad", "Dur.", "Inicio", "Fin"],
    rows_slide1,
    [Emu(457200), Emu(640080), Emu(4937760), Emu(457200), Emu(1828800), Emu(1828800)],
    Pt(6))

txBox = slide.shapes.add_textbox(Emu(91440), Emu(6600000), Emu(12009160), Emu(228600))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "55 de 110 actividades nivel 4 | Lista completa en archivo XML LoRA-HDL-QA_Proyecto_2026.xml | Continúa en diapositiva 30"
p.font.size = Pt(8); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = 'Calibri'

# ── ADD Slide 29b (index 29): Lista de Actividades 2/2 ──
# Insert new slide after slide 29
# python-pptx doesn't easily insert slides at specific position, so we'll add at end
# and reorder later. Actually, let's just modify slide 30 (which was Atributos) and shift everything.
# Simpler: replace slide 30 with Lista 2/2 and recreate Atributos later.

slide = prs.slides[29]  # This was "Atributos de la Actividad" 
clear_slide(slide)
add_header_bar(slide, 30, "LISTA DE ACTIVIDADES (2/2)",
               "Actividades nivel 4 del cronograma — 55 de 110 | PMBOK 8va, Cronograma")

rows_slide2 = [(uid, get_phase(wbs), name[:60], f"{hrs}h", start, finish) for uid, wbs, name, hrs, start, finish in second_half]
add_simple_table(slide, Emu(91440), Emu(1371600), Emu(12009160), Emu(5200000),
    ["UID", "Fase", "Actividad", "Dur.", "Inicio", "Fin"],
    rows_slide2,
    [Emu(457200), Emu(640080), Emu(4937760), Emu(457200), Emu(1828800), Emu(1828800)],
    Pt(6))

txBox = slide.shapes.add_textbox(Emu(91440), Emu(6600000), Emu(12009160), Emu(228600))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "55 de 110 actividades nivel 4 | Total horas nivel 4: 1,287h | Fuente: LoRA-HDL-QA_Proyecto_2026.xml (ProjectLibre 1.9.8)"
p.font.size = Pt(8); p.font.italic = True; p.font.color.rgb = GRAY_TEXT; p.font.name = 'Calibri'

# ── Fix Slide 30 (index 29) is now taken by Lista 2/2, move Atributos to slide 31 (index 30) ──
slide = prs.slides[30]  # Was "Base de las Estimaciones — Duraciones"
clear_slide(slide)
add_header_bar(slide, 31, "ATRIBUTOS DE LA ACTIVIDAD",
               "Campos que describen cada actividad del cronograma | PMBOK 8va, Cronograma")

attr_data = [
    ("UID", "Identificador único en el XML de ProjectLibre", "Ej: 4, 5, 68"),
    ("Nombre", "Descripción de la actividad (verbo + objeto)", "Revisar normativo y reglamento"),
    ("Nivel EDT", "OutlineLevel: 1=Proyecto, 2=Fase, 3=Entregable, 4=Paquete", "1, 2, 3, o 4"),
    ("Código WBS", "Código de la EDT al que pertenece", "1.1.1.4"),
    ("Duración", "Horas de trabajo efectivas (formato PT...H)", "PT8H0M0S (8 horas)"),
    ("Fecha Inicio", "Fecha/hora de inicio planificada", "2026-02-02T08:00:00"),
    ("Fecha Fin", "Fecha/hora de finalización planificada", "2026-02-02T17:00:00"),
    ("Predecesoras", "Tareas que deben completarse antes (FS)", "4FS (termina 4, inicia esta)"),
    ("Recurso", "Recurso asignado (Pablo Flores UID 1)", "Pablo Flores (100%)"),
    ("Tipo", "Fixed Work (trabajo fijo)", "Fixed Work"),
    ("Calendario", "Calendario Guatemala 2026 (UID 11)", "Guatemala - Feriados 2026"),
    ("Hito", "TRUE si duración = 0, FALSE si no", "FALSE para tareas de trabajo"),
]
add_simple_table(slide, Emu(274320), Emu(1371600), Emu(11551920), Emu(4937760),
    ["Atributo", "Descripción", "Valores de Ejemplo"],
    attr_data,
    [Emu(2743200), Emu(5486400), Emu(3657600)],
    Pt(10))

# ── Fix Slide 41 (index 40): RBS — simplified ──
slide = prs.slides[40]
clear_slide(slide)
add_header_bar(slide, 41, "ESTRUCTURA DE DESGLOSE DE RECURSOS (RBS)",
               "Descomposición jerárquica de recursos del proyecto | PMBOK 8va, Recursos")

add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2743200),
         "1. Recurso Humano", [
             "1.1 Pablo Flores — PM / Tesista",
             "    • Rol: Ejecución total del proyecto (gestión + investigación + desarrollo)",
             "    • Tarifa: Q 80.00 / hora",
             "    • Base de tarifa: Investigador junior Guatemala (ref. SENACYT)",
             "    • Dedicación: 100% (1,287h planificadas)",
             "    • Período: Feb 02 — Sep 14, 2026 (225 días)",
             "    • Costo total mano de obra: Q 102,960.00",
             "",
             "1.2 Asesor (designado por FING)",
             "    • ~40h totales estimadas (sin costo directo)",
             "    • Honorarios cubiertos por la universidad",
         ], Pt(11), Pt(9))

add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2743200),
         "2. Equipamiento", [
             "2.1 MacBook Air M1 2020",
             "    • Especificaciones: Chip M1, 8GB RAM, 256GB SSD",
             "    • Precio de adquisición: Q 6,000.00",
             "    • Estado: Propiedad del tesista (costo hundido)",
             "    • Horas estimadas de uso: ~1,000h",
             "    • Uso: Entrenamiento LoRA, desarrollo, redacción",
             "",
             "2.2 Periféricos y accesorios (propios, sin costo)",
             "    • Monitor externo, teclado, mouse",
             "    • Disco externo 1TB para backups",
             "",
             "2.3 Software (100% open-source, sin costo)",
             "    • Python, HuggingFace Transformers, PEFT, MLX",
             "    • ProjectLibre 1.9.8, Git/GitHub, Google Drive",
         ], Pt(11), Pt(9))

add_card(slide, Emu(457200), Emu(4297680), Emu(10972800), Emu(1371600),
         "Resumen de Costos de Recursos", [
             "Mano de obra (Pablo Flores): 1,287h × Q 80/hr = Q 102,960.00",
             "Equipamiento (MacBook Air M1 2020): Q 6,000.00 (ya adquirido, costo hundido)",
             "Materiales (impresión, empastado, posters): Q 800.00",
             "TOTAL RECURSOS DIRECTOS: Q 109,760.00",
             "Nota: La mano de obra es costo de oportunidad. Solo equipamiento y materiales requieren desembolso en efectivo (Q 6,800.00)."
         ], Pt(10), Pt(8))

# ── Fix Slide 36 (index 35): Plan de Gestión de Recursos → simplified ──
slide = prs.slides[35]
clear_slide(slide)
add_header_bar(slide, 36, "PLAN DE GESTIÓN DE RECURSOS",
               "Define cómo se estimarán, adquirirán y gestionarán los recursos | PMBOK 8va")

add_card(slide, Emu(457200), Emu(1371600), Emu(5486400), Emu(2011680),
         "Identificación y Estimación de Recursos", [
             "Fuente primaria: EDT nivel 4 (110 actividades con recurso asignado)",
             "Método: Bottom-up — suma de horas por recurso desde actividades nivel 4",
             "",
             "Recurso Humano:",
             "• Pablo Flores (PM/Tesista) — 1,287h @ Q 80/hr = Q 102,960",
             "• Asesor FING — ~40h (sin costo directo)",
             "",
             "Recursos Físicos:",
             "• MacBook Air M1 2020 (8GB/256GB) — Q 6,000 (ya adquirido)",
             "• Internet 10Mbps, electricidad — costos absorbidos (~Q 550 total)",
             "",
             "Materiales:",
             "• Impresión (3 copias × Q 150), posters (Q 200), papelería (Q 150)",
             "• Total materiales: Q 800"
         ], Pt(11), Pt(9))

add_card(slide, Emu(6400800), Emu(1371600), Emu(5029200), Emu(2011680),
         "Estrategia de Gestión", [
             "Asignación: PM dedicado 100% al proyecto",
             "Calendario: lun-sáb 8h/día, flexibilidad ±2h",
             "",
             "Adquisiciones:",
             "• MacBook ya adquirido — no requiere compra",
             "• Materiales: compra directa por bajo monto (< Q 1,000)",
             "",
             "Monitoreo:",
             "• Registro de horas reales vs planificadas (EVM)",
             "• Reporte mensual de utilización de recursos",
             "",
             "Contingencias:",
             "• Si MacBook falla → usar laboratorio FING",
             "• Si PM indispuesto → ajustar cronograma, notificar asesor",
             "",
             "Optimización:",
             "• Bloques de trabajo 2-4h enfocados (Pomodoro)",
             "• Entrenamiento ML en horarios nocturnos (menor temperatura)"
         ], Pt(11), Pt(9))

# ── Fix Slide 39 (index 38): Requisitos de Recursos ──
slide = prs.slides[38]
clear_slide(slide)
add_header_bar(slide, 39, "REQUISITOS DE RECURSOS",
               "Humanos, equipamiento y materiales necesarios | PMBOK 8va, Recursos")

add_card(slide, Emu(457200), Emu(1371600), Emu(3657600), Emu(2743200),
         "Recurso Humano", [
             "Pablo Flores (PM / Tesista)",
             "",
             "Horas planificadas: 1,287h",
             "Tarifa: Q 80.00 / hora",
             "Costo total: Q 102,960.00",
             "",
             "Disponibilidad: 100% tiempo completo",
             "Período: Feb 02 — Sep 14, 2026",
             "Horario: Lun-Sáb 8h/día",
             "Flexibilidad: 10h/día en semanas críticas",
             "",
             "Rol: investigación, desarrollo,",
             "entrenamiento ML, redacción,",
             "gestión del proyecto, defensa oral."
         ], Pt(11), Pt(9))

add_card(slide, Emu(4389120), Emu(1371600), Emu(3657600), Emu(2743200),
         "Equipamiento", [
             "MacBook Air M1 2020",
             "",
             "Chip: Apple M1 (8-core CPU, 7-core GPU)",
             "RAM: 8GB LPDDR4X unificada",
             "Almacenamiento: 256GB SSD",
             "",
             "Precio: Q 6,000.00",
             "Estado: Propiedad del tesista",
             "Uso: ~1,000h para el proyecto",
             "",
             "Limitación principal:",
             "8GB RAM → quantización 4-bit,",
             "batch=1, modelos ≤ 4B parámetros"
         ], Pt(11), Pt(9))

add_card(slide, Emu(8229600), Emu(1371600), Emu(3657600), Emu(2743200),
         "Software y Materiales", [
             "Software (todo gratuito):",
             "• Python 3.11+",
             "• HuggingFace Transformers",
             "• PEFT / LoRA / MLX",
             "• ProjectLibre 1.9.8",
             "• Git / GitHub",
             "• Google Drive 100GB (EDU)",
             "",
             "Materiales:",
             "• Impresión 3 copias: Q 450",
             "• Posters defensa: Q 200",
             "• Papelería: Q 150",
             "• Total materiales: Q 800",
             "",
             "Infraestructura (absorbida):",
             "• Internet 10Mbps",
             "• Electricidad ~Q 300 total"
         ], Pt(10), Pt(8))

# ── Fix Slide 42 (index 41): Calendarios de Recursos ──
slide = prs.slides[41]
clear_slide(slide)
add_header_bar(slide, 42, "CALENDARIOS DE RECURSOS",
               "Disponibilidad de cada recurso a lo largo del proyecto | PMBOK 8va, Recursos")

add_card(slide, Emu(457200), Emu(1371600), Emu(10972800), Emu(1645920),
         "Calendario del Recurso Principal — Pablo Flores (PM/Tesista)", [
             "Período: Feb 02 — Sep 14, 2026 (225 días calendario)",
             "Días laborables: Lunes a Sábado (6 días/semana)",
             "Horario estándar: 08:00-12:00 y 13:00-17:00 (8h/día efectivas)",
             "Horas máximas por semana: 48h (6 días × 8h)",
             "",
             "Excepciones (no laborable para el recurso):",
             "• Todos los domingos del período",
             "• Semana Santa: Jueves 2 — Sábado 4 de Abril, 2026",
             "• Día del Trabajo: Viernes 1 de Mayo, 2026",
             "• Día del Ejército: Martes 30 de Junio, 2026 (½ día opcional)",
             "• Día de la Independencia: Martes 15 de Septiembre, 2026 (post-proyecto)",
             "",
             "Disponibilidad real: ~190 días laborables efectivos × 8h = ~1,520h disponibles",
             "Horas planificadas: 1,287h → utilización del ~85% del tiempo disponible"
         ], Pt(10), Pt(8))

add_card(slide, Emu(457200), Emu(3200400), Emu(5486400), Emu(1828800),
         "Calendario del Equipamiento — MacBook Air M1", [
             "Disponibilidad: 24/7, compartido con uso personal del tesista",
             "Ventanas óptimas para entrenamiento ML:",
             "• Horario nocturno (18:00-06:00): menor temperatura ambiente",
             "• Fines de semana: bloques de 8-12h continuos disponibles",
             "Horas totales estimadas de uso para el proyecto: ~1,000h",
             "Nota: no requiere calendario exclusivo; se comparte con",
             "actividades personales fuera del horario de proyecto."
         ], Pt(10), Pt(8))

add_card(slide, Emu(6400800), Emu(3200400), Emu(5029200), Emu(1828800),
         "Calendario del Asesor", [
             "Disponibilidad: Sujeta a carga académica FING",
             "Reuniones: 1h/semana (día y hora fijos por acordar)",
             "Revisiones de hitos: bajo demanda, con 5 días de anticipación",
             "Período de mayor carga: inicio de semestre (Ene, Jul)",
             "",
             "Nota: El asesor no tiene dedicación exclusiva.",
             "Se debe planificar con 1-2 semanas de anticipación",
             "para revisiones de documentos extensos."
         ], Pt(10), Pt(8))

# ── Save ──
output_path = '/Users/pabloflores/Documents/estructura/revision/entregables_finales/LoRA_Plan_Direccion_Proyecto_Completo_v2.pptx'
prs.save(output_path)
print(f"✅ Presentation updated: {output_path}")
print(f"   Slides: {len(prs.slides)}")
